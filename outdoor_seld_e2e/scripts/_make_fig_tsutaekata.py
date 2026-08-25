# -*- coding: utf-8 -*-
"""「伝え方」の図 ― 首元の振動子5個で方向、鳴り方で危険度。

背景: 2026-08-25の中間発表リハで「関連研究のところで、自分がどういうところを
付け加えるのか具体的なものを出してもらえると良い。例えば首のところで振動で伝えるなら、
ピーピーピーならそんなに怖がらなくていい、ピピピピピなら怖い、みたいな」と指摘された。
デバイス構想（首元に振動子5個・2026-08-24 本人）を図にしたもの。**未実装の構想**。

貼り先候補: p3「伝え方」の右、または提案②2層スライドの「首元振動デバイス」の箱
出力: md/seminar/図_伝え方_2026-08-25.pptx（コピペ用・1枚）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
FAINT = RGBColor(0xF4, 0xF5, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def shape(x, y, w, h, fill, line=None, lw=1.0, kind=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(kind, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


# ---- テンプレ骨格 ----
shape(44, 14, 1.4, H - 28, PUR)
shape(58, 46, 30, 6, PUR)
txt(98, 30, W - 160, 40, [("伝え方　―　方向は「場所」、危険度は「鳴り方」",
                           {"size": 22, "bold": True, "spc": 220})])
shape(44, 90, W - 90, 1.2, INK)

txt(70, 102, W - 140, 24,
    [("首元に振動子を5個並べる構想。", {"size": 13.5, "color": SUB}),
     ("鳴っている位置で方向を、鳴り方の強さと間隔で危険度を表す。",
      {"size": 13.5, "bold": True})])

# ================= 左: 方向 =================
LX, LY, LW, LH = 70, 140, 400, 250
shape(LX, LY, LW, LH, WHITE, line=LINE, lw=1.0)
shape(LX, LY, 4, LH, PUR)
txt(LX + 20, LY + 14, LW - 40, 22, [("方向 ＝ どの振動子が鳴るか", {"size": 14, "bold": True})])

# 首を上から見た弧
NECK_CX, NECK_CY, RR = LX + LW / 2, LY + 168, 96.0
shape(NECK_CX - 34, NECK_CY - 34, 68, 68, FAINT, line=LINE, lw=1.0,
      kind=MSO_SHAPE.OVAL)
txt(NECK_CX - 40, NECK_CY - 12, 80, 22,
    [("首", {"size": 13, "bold": True, "color": SUB})], align=PP_ALIGN.CENTER)

import math
pts = [("左", 200), ("左前", 235), ("前", 270), ("右前", 305), ("右", 340)]
for i, (lab, deg) in enumerate(pts):
    a = math.radians(deg)
    cx = NECK_CX + RR * math.cos(a)
    cy = NECK_CY + RR * math.sin(a)
    on = (lab == "右前")
    shape(cx - 13, cy - 13, 26, 26, GOLD if on else CHIP,
          line=(RED if on else LINE), lw=(1.6 if on else 1.0), kind=MSO_SHAPE.OVAL)
    txt(cx - 30, cy - 34, 60, 18,
        [(lab, {"size": 10.5, "bold": on, "color": (RED if on else MUTED)})],
        align=PP_ALIGN.CENTER)

txt(LX + 20, LY + LH - 42, LW - 40, 34,
    [("右前から車 → ", {"size": 12, "color": SUB}),
     ("右前の1個だけが鳴る", {"size": 12, "bold": True, "color": RED})],
    align=PP_ALIGN.CENTER)

# ================= 右: 危険度 =================
RX, RY, RW, RH2 = 490, 140, 400, 250
shape(RX, RY, RW, RH2, WHITE, line=LINE, lw=1.0)
shape(RX, RY, 4, RH2, GOLD)
txt(RX + 20, RY + 14, RW - 40, 22, [("危険度 ＝ 強さと間隔", {"size": 14, "bold": True})])

# 3段のパターン
rows = [
    ("至近警告", RED, [(0, 26), (30, 26), (60, 26), (90, 26), (120, 26), (150, 26)],
     "強く・短い間隔"),
    ("注意", GOLD, [(0, 26), (60, 26), (120, 26)], "弱く・広い間隔"),
    ("抑制", MUTED, [], "鳴らさない"),
]
ry = RY + 52
for lab, col, bars, note in rows:
    txt(RX + 20, ry + 6, 76, 20, [(lab, {"size": 12.5, "bold": True, "color": col})])
    bx0 = RX + 106
    if bars:
        for ox, bw in bars:
            h = 26 if col is RED else 17
            shape(bx0 + ox, ry + 6 + (26 - h) / 2, bw, h, col)
    else:
        shape(bx0, ry + 14, 176, 2, CHIP)
    txt(RX + 300, ry + 8, 84, 20,
        [(note, {"size": 10.5, "color": MUTED})], align=PP_ALIGN.RIGHT)
    ry += 56

txt(RX + 20, RY + RH2 - 42, RW - 40, 34,
    [("同じ「近づいてくる」でも、", {"size": 12, "color": SUB}),
     ("慌てるべきかが指先で分かる", {"size": 12, "bold": True})], align=PP_ALIGN.CENTER)

# ---- 締め ----
shape(70, 410, 6, 44, PUR)
shape(76, 410, W - 146, 44, WHITE, line=LINE, lw=1.0)
txt(94, 418, W - 180, 30,
    [("音を「聞かせ直す」のではなく、", {"size": 13.5, "color": SUB}),
     ("触覚に置き換えて伝える", {"size": 13.5, "bold": True}),
     ("。デバイスは未実装で、本発表では構想。",
      {"size": 13.5, "color": SUB})], anchor=MSO_ANCHOR.MIDDLE)

# ---- フッター ----
txt(58, 506, 90, 20, [("2026/08/30", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "背景・目的")
    c = shape(x, 502, wd, 22, NAVY if active else CHIP)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(10.5)
    r.font.bold = active
    r.font.color.rgb = WHITE if active else MUTED
    r.font.name = "Meiryo"
    meiryo(r)
    x += wd + 8

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"図_伝え方_2026-08-25.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
