# -*- coding: utf-8 -*-
"""【修正スライド1枚だけ】夏ゼミの振り返り — いただいた4件に回答を付けた版。

本人の依頼（2026-09-03）: 7枚目はコメントを並べただけなので、それぞれに回答を付けたい。
デザインは元のまま（案Aは優先度を下げる判断・同日）。

回答の出典:
 1件目 → out/notify_v42_sweep2/q2_anzen/q2_table.md（型別の数値）
 2件目 → md/research/触覚方向提示_先行調査_2026-08-30.md（間隔と弁別率）
 3件目 → md/design/実録ハンドブック_2026-08-13.md §8（前方欄の追加・2026-08-30決定）
 4件目 → out/joycon_demo/（2026-09-02 実機で振動成功）

出力: md/seminar/修正_夏ゼミの振り返り_2026-09-15.pptx（1枚のみ）
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF5, 0xF3, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])
SECTIONS = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


# ---- 骨格 ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 170, 40,
    [("夏 ゼ ミ の 振 り 返 り", {"size": 22, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

txt(70, 100, 820, 20,
    [("いただいた4件すべてに手を動かしました。", {"size": 12.5, "bold": True})])

rows = [
    ("狭いところをすり抜けていく車と対向してくる車が、同じ至近警告になって"
     "区別できないのではないか",
     [("方位の変化率で区別する仕組みを実装。", {"bold": True}),
      ("危ない対向車への強い警告は 81→90%、鳴らすべきでない対向車で"
       "黙れた割合は 33→71%。", {})],
     "8・9枚目", RED),
    ("振動による方向指示のデバイスは、全盲や弱視の方向けの製品やノウハウも"
     "いろいろあるので、参考になるかもしれません",
     [("先行研究を調査。", {"bold": True}),
      ("腰の振動子は間隔107mmで8方向を98%当てられるが、72mmでは74%に落ちる。"
       "首は腰より鈍いので、", {}),
      ("首元は4〜6方向が現実的", {"bold": True}),
      ("と分かった。", {})],
     "調査済", GREEN),
    ("前方を検出対象外としているが、研究上はとりあえず全方位録ってみるで"
     "いいのでは。データを見たら思いがけないことが分かるかも",
     [("記録紙に前方イベントの欄を追加した。", {"bold": True}),
      ("マイクは最初から全方位録っているので追加費用はゼロ。採点では既定で除外し"
       "別集計にする（録らなければ二度と取れないため）。", {})],
     "反映済", GREEN),
    ("振動デバイスを一から作るとハードルが高そう。Joy-conとUnityとかで"
     "簡単に作ってみてもいいかも",
     [("そのとおりに作った。9/2に実機で振動成功。", {"bold": True}),
      ("ハンダごて不要・部品の購入ゼロ。通知層の出力をそのまま振動に変えている。", {})],
     "10枚目", RED),
]

y = 132.0
for quote, answer, tag, col in rows:
    # いただいたコメント
    rect(70, y + 6, 6, 6, MUTED)
    txt(88, y, 660, 34, [(quote, {"size": 11, "color": SUB})], line=1.35)
    b = rect(790, y - 1, 100, 20, None, line=col, lw=1.1)
    label(b, tag, size=10, color=col)
    # やったこと
    txt(88, y + 38, 30, 20, [("→", {"size": 13, "bold": True, "color": PUR})])
    runs = [(t, dict(o, size=o.get("size", 11.5),
                     color=o.get("color", INK if o.get("bold") else SUB)))
            for t, o in answer]
    txt(116, y + 38, 774, 40, [runs], line=1.4)
    y += 88
    if y < 440:
        rect(70, y - 12, 820, 0.8, HAIR)

# ---- フッター ----
txt(58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
ws = [16 + len(s) * 12.0 for s in SECTIONS]
x = (W - (sum(ws) + 8 * (len(SECTIONS) - 1))) / 2
for s, wd in zip(SECTIONS, ws):
    on = (s == "振り返り")
    c = rect(x, 502, wd, 22, NAVY if on else CHIP)
    label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("7", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"修正_夏ゼミの振り返り_2026-09-15.pptx")
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
