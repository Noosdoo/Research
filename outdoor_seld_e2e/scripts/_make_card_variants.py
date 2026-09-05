# -*- coding: utf-8 -*-
"""カードの「種類の示し方」を3案つくって見比べる（9枚目を素材にする）。

経緯（2026-09-05）:
 上端の色帯 → 左端の縦線 に変えたが、本人「いやそれも色帯じゃん」。
 **箱に色の棒をつけて種類分けする発想そのもの**が Claude の癖、という指摘。
 根拠にした4枚目の縦線も、元は同じ生成物だったので循環していた。
 → 棒を一切使わない方法を3つ、実物で比較する。

6案（いずれも色の棒を使わない）:
  A 枠線そのものを色にする        箱の輪郭が色を持つ。何も貼り付けない
  B 見出し行だけ薄く色を敷く      表の見出し行と同じ発想。本文は白のまま
  C 箱全体を薄いグレーで塗る      枠線なし。metropolis の block=fill と同じ
  D 右と下の枠線を消す            上と左だけの「かぎ形」。パワポ研「枠線デザイン9選」⑦
  E 薄い背景＋同系色の少し濃い枠線  Webクリエイターボックスの定石。上品で色が残る
  F 薄い色の背景だけ（枠線なし）    Cの色つき版。②と③の区別が残る

素材は本人の編集中ファイルの9枚目（②と③のカードが並んでいる枚）をそのまま複製して
加工する。文言・配置には触らない。

使い方:
  python scripts/_make_card_variants.py <入力.pptx> [出力.pptx] [ページ番号]
"""
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

BAR_H_EMU = int(Pt(3.5))
MIN_W_PT = 128.0
GREY = RGBColor(0xF2, 0xF3, 0xF5)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)


def light(rgb, k=0.14):
    """白に寄せた薄い色をつくる。"""
    return RGBColor(*[int(255 - (255 - c) * k) for c in tuple(rgb)])


def fill_rgb(sh):
    try:
        if sh.fill.type == 1:
            return sh.fill.fore_color.rgb
    except Exception:
        pass
    return None


def clone_slide(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(
            copy.deepcopy(shp._element), "p:extLst")
    return dst


def cards_of(slide):
    """(帯, カード本体, 色) の組を返す。"""
    out = []
    for bar in list(slide.shapes):
        if bar.height != BAR_H_EMU or bar.width / 12700.0 < MIN_W_PT:
            continue
        rgb = fill_rgb(bar)
        if rgb is None:
            continue
        card = None
        for sh in slide.shapes:
            if sh is bar or sh.height < int(Pt(20)):
                continue
            if abs(sh.left - bar.left) <= int(Pt(6)) and \
               abs(sh.top - bar.top) <= int(Pt(6)) and \
               abs(sh.width - bar.width) <= int(Pt(6)):
                if card is None or sh.height > card.height:
                    card = sh
        out.append((bar, card, rgb))
    return out


def drop(sh):
    sh._element.getparent().remove(sh._element)


def note(slide, text):
    tb = slide.shapes.add_textbox(Pt(70), Pt(58), Pt(820), Pt(24))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    r.font.name = "Meiryo"


def apply_A(slide):
    """枠線そのものを色にする。"""
    for bar, card, rgb in cards_of(slide):
        if card is not None:
            card.line.color.rgb = rgb
            card.line.width = Pt(1.5)
        drop(bar)


def apply_B(slide):
    """見出し行だけ薄く色を敷く（本文は白のまま）。"""
    for bar, card, rgb in cards_of(slide):
        if card is not None:
            band = slide.shapes.add_shape(
                1, card.left, card.top, card.width, Pt(42))
            band.fill.solid()
            band.fill.fore_color.rgb = light(rgb)
            band.line.fill.background()
            band.shadow.inherit = False
            # 見出しの文字より後ろへ回す
            sp = band._element
            sp.getparent().remove(sp)
            card._element.addnext(sp)
        drop(bar)


def apply_C(slide):
    """箱全体を薄いグレーで塗り、枠線を消す。"""
    for bar, card, rgb in cards_of(slide):
        if card is not None:
            card.fill.solid()
            card.fill.fore_color.rgb = GREY
            card.line.fill.background()
        drop(bar)


def apply_D(slide):
    """右と下の枠線を消す＝上と左だけの「かぎ形」にする（パワポ研⑦）。"""
    for bar, card, rgb in cards_of(slide):
        if card is not None:
            card.line.fill.background()
            for x, y, w, h in [(card.left, card.top, card.width, Pt(1.6)),
                               (card.left, card.top, Pt(1.6), card.height)]:
                ln = slide.shapes.add_shape(1, x, y, w, h)
                ln.fill.solid()
                ln.fill.fore_color.rgb = rgb
                ln.line.fill.background()
                ln.shadow.inherit = False
        drop(bar)


def apply_E(slide):
    """薄い背景＋同系色の少し濃い枠線（Webクリエイターボックスの定石）。"""
    for bar, card, rgb in cards_of(slide):
        if card is not None:
            card.fill.solid()
            card.fill.fore_color.rgb = light(rgb, 0.07)
            card.line.color.rgb = light(rgb, 0.30)
            card.line.width = Pt(1.0)
        drop(bar)


def apply_F(slide):
    """薄い色の背景だけ・枠線なし（Cの色つき版）。"""
    for bar, card, rgb in cards_of(slide):
        if card is not None:
            card.fill.solid()
            card.fill.fore_color.rgb = light(rgb, 0.10)
            card.line.fill.background()
        drop(bar)


def main() -> int:
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2]) if len(sys.argv) > 2 else \
        src.with_name("カード3案.pptx")
    page = int(sys.argv[3]) if len(sys.argv) > 3 else 9

    s = Presentation(str(src))
    base = list(s.slides)[page - 1]

    out = Presentation()
    out.slide_width = s.slide_width
    out.slide_height = s.slide_height
    for label, fn in [("案A　枠線そのものを色にする", apply_A),
                      ("案B　見出し行だけ薄く色を敷く", apply_B),
                      ("案C　箱全体を薄いグレー・枠線なし", apply_C),
                      ("案D　右と下の枠線を消す（上と左だけのかぎ形）", apply_D),
                      ("案E　薄い背景＋同系色の少し濃い枠線", apply_E),
                      ("案F　薄い色の背景だけ・枠線なし", apply_F)]:
        sl = clone_slide(base, out)
        fn(sl)
        note(sl, label)
        print("作成:", label)
    out.save(str(dst))
    print("保存:", dst)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
