# -*- coding: utf-8 -*-
"""質疑用の付録スライド図5枚を生成（本体デッキと同配色・メイリオ・ネイティブ図形）。

1. 実録の測定方法 — 1テイクの流れ＋撮れ高120本（2026-08-14改訂）
2. 実録の正解ラベル — 合成との役割分担
3. 距離推定の統合 — 出力ヘッドの拡張
4. 音だけで距離が分かる4つの手がかり
5. 同一クラス複数台の解消（前回指摘への回答）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x16, 0x23, 0x3A)
INK2 = RGBColor(0x3A, 0x46, 0x58)
MUTED = RGBColor(0x66, 0x70, 0x7F)
PAPER2 = RGBColor(0xEF, 0xEF, 0xEA)
HAIR = RGBColor(0xD9, 0xDA, 0xD2)
AMBER = RGBColor(0xE8, 0xA2, 0x00)
AMBERD = RGBColor(0xB3, 0x7E, 0x00)
RED = RGBColor(0xC4, 0x43, 0x2B)
GREEN = RGBColor(0x2F, 0x7D, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMB18 = RGBColor(0xF4, 0xE8, 0xC7)
RED22 = RGBColor(0xEC, 0xCF, 0xC7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def meiryo(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", "メイリオ")


def txt(sl, x, y, w, h, lines, size=12.5, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, mono=False, gap=4):
    """lines: str または (text, {size,bold,color}) のリスト。"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        if isinstance(ln, str):
            ln = (ln, {})
        t, o = ln
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.color.rgb = o.get("color", color)
        f.name = "Consolas" if mono else "Meiryo"
        if not mono:
            meiryo(r)
    return tb


def box(sl, x, y, w, h, fill=WHITE, line=HAIR, line_w=1.0, radius=0.08):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def arrow_r(sl, x, y, w=0.32, h=0.3, color=AMBER):
    sh = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    sh.rotation = 90
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def title(sl, text):
    txt(sl, 0.55, 0.32, W - 1.1, 0.6, [(text, {"size": 21, "bold": True})])
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.95),
                              Inches(1.6), Inches(0.045))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()
    bar.shadow.inherit = False


def new_slide():
    return prs.slides.add_slide(BLANK)


# ============ 1. 実録の測定方法 — 1テイクの流れ ============
sl = new_slide()
title(sl, "付録：実録の測定方法 — 1テイクの流れ（全テイク共通の骨格）")

steps = [
    ("①", "儀式 90秒", "録音開始→手拍子1回（時刻同期）→騒音計と並べ無言60秒・読み値は計測後にスレート（絶対較正）→ベル4方位×1打（方位校正）", AMBER),
    ("②", "本番（無言）", "立って待つだけ。録りっぱなしで1通過＝1テイク。統制試行は台本どおり協力者が走行", INK2),
    ("③", "直後 30秒", "小声スレート「試行3終了・車1台・右から・徐行・横2m」→そのまま注釈CSVの1行に", INK2),
    ("④", "現場QC 5分", "イヤホンで確認：4ch入っているか・音割れないか・校正音が録れたか →OKなら次の地点へ", GREEN),
]
cx = 0.55
for i, (no, name, body, col) in enumerate(steps):
    bw = 2.92
    box(sl, cx, 1.25, bw, 2.9)
    hd = box(sl, cx, 1.25, bw, 0.62, fill=col, line=None, radius=0.16)
    txt(sl, cx + 0.12, 1.33, bw - 0.24, 0.5,
        [(f"{no} {name}", {"size": 14, "bold": True, "color": WHITE})])
    txt(sl, cx + 0.16, 2.05, bw - 0.32, 2.0, body, size=11.5, color=INK2)
    if i < 3:
        arrow_r(sl, cx + bw + 0.02, 2.55)
    cx += bw + 0.32

box(sl, 0.55, 4.55, W - 1.1, 2.35, fill=PAPER2, line=None)
txt(sl, 0.85, 4.75, 8.5, 0.4, [("撮れ高 計100本＋静止/歩行対比20本＋負例は露出時間100分以上"
                                "（実働3日＋注釈1日・計画値）",
                                {"size": 13.5, "bold": True})])
chips = [
    ("20", "走行車の統制録音", "caution帯10・safe帯10（後方6以上）"),
    ("20", "弱点検証", "EV5・発進5・自転車5・見通し不良の角5"),
    ("20", "負例", "静穏・繁華街・雨上がり・風（誤警告率）"),
    ("20", "固定・機会", "踏切8＋バック音統制4＋機会枠8"),
    ("20", "キックボード", "横1.0〜1.5mの至近帯を統制"),
]
cx = 0.85
for n, name, note in chips:
    bw = 2.28
    box(sl, cx, 5.25, bw, 1.4, fill=WHITE, line=HAIR)
    txt(sl, cx + 0.12, 5.36, bw - 0.24, 0.55,
        [(n + "本", {"size": 17, "bold": True, "color": AMBERD})])
    txt(sl, cx + 0.12, 5.86, bw - 0.24, 0.75,
        [(name, {"size": 11.5, "bold": True}), (note, {"size": 9.5, "color": MUTED})], gap=1)
    cx += bw + 0.11

# ============ 2. 実録の正解ラベル — 役割分担 ============
sl = new_slide()
title(sl, "付録：実録の正解ラベル — 合成との役割分担")

box(sl, 0.55, 1.25, 6.0, 4.35)
box(sl, 0.55, 1.25, 6.0, 0.66, fill=INK, line=None, radius=0.12)
txt(sl, 0.75, 1.35, 5.6, 0.5, [("合成データ ＝ 精密な定規", {"size": 14.5, "bold": True, "color": WHITE})])
txt(sl, 0.85, 2.15, 5.4, 3.3, [
    ("正解はシミュレーションが自動生成（誤差ゼロ）", {"bold": True}),
    ("・0.1秒ごとの（クラス・方位・仰角・距離）", {}),
    ("・方向誤差 中央2.0°（車）、至近距離誤差 0.21m の精密検証はこちらの役割", {}),
    ("・学習もすべて合成（1万200クリップ）", {}),
], size=12.5, color=INK2, gap=7)

box(sl, 6.85, 1.25, 5.95, 4.35)
box(sl, 6.85, 1.25, 5.95, 0.66, fill=AMBERD, line=None, radius=0.12)
txt(sl, 7.05, 1.35, 5.5, 0.5, [("実録 ＝ 本物の物差し（テイク単位の粗い注釈）", {"size": 14.5, "bold": True, "color": WHITE})])
box(sl, 7.1, 2.15, 5.45, 1.35, fill=PAPER2, line=HAIR)
txt(sl, 7.3, 2.3, 5.1, 1.1, [
    ("trial=3, event_id=1, class=車, 象限=右,", {}),
    ("t_cpa=00:41(±1s), 横距離=2m(コーン実測),", {}),
    ("速度=徐行(目視), LAeq=52.3, 装着=チェスト", {}),
], size=11, color=INK2, mono=True, gap=2)
txt(sl, 7.1, 3.7, 5.5, 1.8, [
    ("答え合わせする問い（粗さに合わせる）", {"bold": True}),
    ("・そのクラスを検出できたか／方向の象限は合っているか", {}),
    ("・最接近の何秒前に警告が出たか／負例で誤警告は何回/時か", {}),
], size=12.5, color=INK2, gap=6)

box(sl, 0.55, 5.9, W - 1.1, 0.95, fill=AMB18, line=None)
txt(sl, 0.85, 6.05, W - 1.7, 0.7,
    [("フレーム単位の x,y,z・距離を手入力することはない（人間には不可能・設計上も不要）。"
      "注釈は1イベント1行（複数イベントは行を追加）＋時刻合わせで半日の作業量", {"size": 13, "bold": True})])

# ============ 3. 距離推定の統合 ============
sl = new_slide()
title(sl, "付録：距離推定の統合 — 出力ヘッドに「距離」を1軸追加")

box(sl, 0.55, 1.5, 2.1, 1.5, fill=PAPER2, line=HAIR)
txt(sl, 0.55, 1.85, 2.1, 0.9, [("4ch FOA音声", {"bold": True, "size": 13}),
                               ("（0.1秒フレーム）", {"size": 10.5, "color": MUTED})],
    align=PP_ALIGN.CENTER, gap=2)
arrow_r(sl, 2.75, 2.1)
box(sl, 3.2, 1.4, 2.7, 1.7, fill=INK, line=None)
txt(sl, 3.2, 1.8, 2.7, 1.0, [("PSELDNets", {"bold": True, "size": 14, "color": WHITE}),
                             ("（事前学習済み・FT）", {"size": 10.5, "color": HAIR})],
    align=PP_ALIGN.CENTER, gap=2)
arrow_r(sl, 6.0, 2.1)

txt(sl, 6.5, 1.15, 6.3, 0.35, [("出力（クラスごと・同時最大3トラック）", {"size": 11.5, "color": MUTED})])
for t in range(3):
    y = 1.55 + t * 0.62
    txt(sl, 6.5, y + 0.05, 1.05, 0.4, [(f"トラック{t+1}", {"size": 10.5, "color": MUTED})])
    for j, lab in enumerate(["x", "y", "z"]):
        b = box(sl, 7.6 + j * 0.62, y, 0.56, 0.5, fill=PAPER2, line=HAIR)
        txt(sl, 7.6 + j * 0.62, y + 0.06, 0.56, 0.38, [(lab, {"size": 12, "bold": True})],
            align=PP_ALIGN.CENTER)
    box(sl, 9.55, y, 0.85, 0.5, fill=AMBER, line=None)
    txt(sl, 9.55, y + 0.06, 0.85, 0.38, [("距離", {"size": 12, "bold": True, "color": WHITE})],
        align=PP_ALIGN.CENTER)
txt(sl, 10.6, 1.5, 2.4, 1.9, [
    ("x,y,z＝矢印", {"bold": True, "size": 11.5}),
    ("向き＝方向／長さ＝確度", {"size": 10.5, "color": MUTED}),
    ("距離＝メートル", {"bold": True, "size": 11.5}),
    ("log符号化（近くほど細かく）", {"size": 10.5, "color": MUTED}),
], gap=3)

rules = [
    ("① 鳴っている時だけ学習", "距離の誤差は音が活動中のフレームだけ採点（活動マスク）。無音区間で距離を当てさせない"),
    ("② 方向の学習を壊さない", "出力と正解の対応付けは従来どおり方向(x,y,z)だけで決定し、距離誤差は後から重み付きで加算。距離学習を切った対照実験で等価性を確認済み"),
    ("③ 10m頭打ちの教訓", "旧版は出力のtanhで距離が±10mに飽和。「距離だけ線形出力＋log符号化」に修正し飽和を解消（fold11遠方の開発評価ではGT303mのサイレンを検出。遠距離の距離精度の主張ではない）。至近5m以内は中央0.21m"),
]
cx = 0.55
for name, body in rules:
    bw = 4.0
    box(sl, cx, 3.85, bw, 2.55)
    txt(sl, cx + 0.18, 4.0, bw - 0.36, 0.45, [(name, {"size": 13.5, "bold": True, "color": AMBERD})])
    txt(sl, cx + 0.18, 4.55, bw - 0.36, 1.75, body, size=11.5, color=INK2)
    cx += bw + 0.36

txt(sl, 0.55, 6.65, W - 1.1, 0.5,
    [("通知層はこの（クラス・方向・距離）を読み、1.5m×2フレーム連続→至近警告／≤3.0m→注意／>3.2m→抑制 に振り分ける（同一物体は方位連結±60°の近似・完全な追跡ではない）",
      {"size": 12, "color": MUTED})])

# ============ 4. 距離の手がかり4種 ============
sl = new_slide()
title(sl, "付録：音だけでなぜ距離が分かるのか — 4つの手がかり")

cues = [
    ("①", "クラス条件つき音量", "音源の音量を法規・実測レンジに固定してあるので「サイレンなのに小さい＝遠い」が成立。種類が分かることが距離の前提＝同時推定型にした理由", AMBER),
    ("②", "大気吸収の音色", "高い音ほど空気に吸われる → 遠い音は「こもる」。音量が同じ60dBでも音色が違う", INK2),
    ("③", "動きの幾何", "近い音源ほど方位と音量の変化が速い（目の前を通る車は一瞬、遠くの車はゆっくり）", INK2),
    ("④", "地面反射", "直接音と地面反射音の干渉パターンが距離で変わる", INK2),
]
cx = 0.55
for no, name, body, col in cues:
    bw = 2.92
    box(sl, cx, 1.3, bw, 3.6)
    txt(sl, cx + 0.16, 1.45, bw - 0.32, 0.5, [(no, {"size": 20, "bold": True, "color": AMBERD})])
    txt(sl, cx + 0.16, 2.05, bw - 0.32, 0.6, [(name, {"size": 13.5, "bold": True})])
    txt(sl, cx + 0.16, 2.7, bw - 0.32, 2.1, body, size=11, color=INK2)
    cx += bw + 0.32

box(sl, 0.55, 5.2, W - 1.1, 1.65, fill=PAPER2, line=None)
txt(sl, 0.85, 5.38, W - 1.7, 1.3, [
    ("残る原理的な曖昧さ：同じクラス内の音量幅（車で約7dB ≒ 距離2倍分）→ 至近警告 約7割の一因", {"size": 13, "bold": True}),
    ("9月のablationのうち「大気吸収なし」「地面反射なし」の2本は、手がかり②④の寄与をそれぞれ直接測る実験になっている", {"size": 12.5, "color": INK2}),
], gap=8)

# ============ 5. 複数台の解消 ============
sl = new_slide()
title(sl, "付録：同一クラス複数台の問題（前回指摘）→ 解消")

box(sl, 0.55, 1.3, 5.6, 4.5)
txt(sl, 0.85, 1.45, 5.0, 0.4, [("車の検出率 vs 同時台数（評価1,800クリップ実測）",
                                {"size": 12.5, "bold": True})])
base_y, plot_h = 5.1, 2.6
vals = [("1台", 99.2), ("2台", 99.0), ("3台以上", 99.3)]
for i, (lab, v) in enumerate(vals):
    bx = 1.15 + i * 1.65
    bh = (v - 90.0) / 10.0 * plot_h
    b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(base_y - bh),
                            Inches(1.0), Inches(bh))
    b.fill.solid()
    b.fill.fore_color.rgb = AMBER
    b.line.fill.background()
    b.shadow.inherit = False
    txt(sl, bx - 0.2, base_y - bh - 0.42, 1.4, 0.35, [(f"{v}%", {"size": 13, "bold": True})],
        align=PP_ALIGN.CENTER)
    txt(sl, bx - 0.2, base_y + 0.08, 1.4, 0.35, [(lab, {"size": 11.5, "color": INK2})],
        align=PP_ALIGN.CENTER)
ax = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(base_y),
                         Inches(4.7), Inches(0.02))
ax.fill.solid()
ax.fill.fore_color.rgb = INK2
ax.line.fill.background()
ax.shadow.inherit = False
txt(sl, 0.85, 5.45, 5.0, 0.35, [("縦軸は90〜100%に拡大表示。台数によらずフラット",
                                 {"size": 10.5, "color": MUTED})])

box(sl, 6.5, 1.3, 6.3, 2.15)
txt(sl, 6.7, 1.45, 5.9, 1.9, [
    ("原因と対策", {"size": 13.5, "bold": True, "color": AMBERD}),
    ("原因は学習データに同一クラスの多重シーンが薄かったこと。複数車の同時シーンを追加し、追加の有無だけを変えた対照実験で効果を確定", {"size": 12, "color": INK2}),
], gap=5)
box(sl, 6.5, 3.65, 6.3, 2.15)
txt(sl, 6.7, 3.8, 5.9, 1.9, [
    ("仕組み", {"size": 13.5, "bold": True, "color": AMBERD}),
    ("出力がクラスごとに最大3トラック（矢印3本）あり、学習時に出力と正解の対応付けを自動で解くため、同じクラスの車3台まで別々の方向・距離で同時に指せる", {"size": 12, "color": INK2}),
], gap=5)

box(sl, 0.55, 6.15, W - 1.1, 0.95, fill=AMB18, line=None)
txt(sl, 0.85, 6.3, W - 1.7, 0.7,
    [("残る粒度の課題（正直に）：2台同時のフレームで「1台以上検出」99.92%（存在見逃し0.08%）、「2台として分離」は77.7% ＝ 見逃しはほぼ無いが完全分離は今後",
      {"size": 12.5, "bold": True})])

OUT = r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/図_付録QA用_2026-08-13.pptx"
prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
