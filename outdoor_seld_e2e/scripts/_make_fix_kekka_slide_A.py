# -*- coding: utf-8 -*-
"""【デザイン案A】「その結果」— 罫線と余白だけで組む版（カード・塗り囲みを使わない）。

本人の指摘（2026-09-03）:
 「白カード＋色帯＋薄色の囲み＋丸バッジ」の作りが毎回同じで "Claudeすぎる"。
 案A＝罫線だけの学会スライド風で1枚作って見せること。ただし
 「シンプルすぎてつまらないスライドにはしたくない」。

つまらなくしないための仕掛け:
 数字を表に並べる代わりに、0〜100%の同じ物差しの上に
 前(v4.1)＝白丸 / 今(v4.2)＝塗り丸 を置き、線でつなぐ（ダンベル図）。
 4行の線の長さの差がそのまま「どこが効いてどこが効かなかったか」になる。
 とくに最終行（62→64）は線がほぼ点になり、弱点が文章より先に目に入る。

色: 本文=黒 / 補助=灰 / アクセント=テンプレートの紫1色のみ。
 伸びなかった行だけ灰にして、色相ではなく明度で差をつける。

数値の正本 = out/notify_v42_sweep2/q2_anzen/q2_table.md
出力: md/seminar/修正_その結果_案A_2026-09-15.pptx（1枚のみ）
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x1E, 0x22, 0x2C)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x9A, 0x9E, 0xA8)
FAINT = RGBColor(0xC9, 0xCB, 0xD2)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
PUR = RGBColor(0x7E, 0x6F, 0x98)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


# ---- 骨格（本人のテンプレートのまま） ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 170, 40, [("そ の 結 果", {"size": 22, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

txt(70, 102, 820, 22,
    [[("危ない車には強く鳴るようになり、鳴らすべきでない対向車では黙るようになった。",
       {"size": 13, "bold": True})]])

# ---- 物差し（0〜100%）を一度だけ引く ----
X0, X1 = 268.0, 772.0          # 0% と 100% の位置
SY = 140.0


def px(v):
    return X0 + (X1 - X0) * v / 100.0


for v in (0, 25, 50, 75, 100):
    rect(px(v), SY + 8, 0.8, 5, FAINT)
for v, t in ((0, "0"), (50, "50"), (100, "100%")):
    txt(px(v) - 24, SY - 6, 48, 14, [(t, {"size": 8.5, "color": MUTED})],
        align=PP_ALIGN.CENTER)
rect(X0, SY + 13, X1 - X0, 0.8, HAIR)

# ---- 2つの区分 × 2行 ----
sections = [
    ("危ない車（最接近 1.5m 以内）", "強い警告が届いた割合",
     [("対向型", 460, 81, 90, True),
      ("すり抜け型", 43, 77, 93, True)]),
    ("鳴らすべきでない車（最接近 3.2m 超）", "黙っていられた割合",
     [("対向型", 52, 33, 71, True),
      ("すり抜け型", 501, 62, 64, False)]),
]

y = SY + 34
for title, metric, rows in sections:
    rect(70, y, 820, 1.0, HAIR)
    txt(70, y + 10, 400, 20,
        [[(title, {"size": 12.5, "bold": True}),
          ("　" + metric, {"size": 11, "color": SUB})]])
    y += 40
    for name, n, old, new, moved in rows:
        col = INK if moved else MUTED
        dot = PUR if moved else MUTED
        txt(88, y - 1, 120, 20, [(name, {"size": 12, "color": col})])
        txt(196, y + 2, 60, 16,
            [("n=%d" % n, {"size": 9, "color": MUTED})], align=PP_ALIGN.RIGHT)
        # 前→今 をつなぐ線
        rect(px(old), y + 8, px(new) - px(old), 2.2, dot)
        # 前＝白丸（輪郭のみ）
        rect(px(old) - 5, y + 4, 10, 10, WHITE, line=MUTED, lw=1.1,
             shape=MSO_SHAPE.OVAL)
        # 今＝塗り丸
        rect(px(new) - 5.5, y + 3.5, 11, 11, dot, shape=MSO_SHAPE.OVAL)
        # 前の値（丸の左下に小さく）
        txt(px(old) - 46, y + 14, 40, 14,
            [("%d" % old, {"size": 9, "color": MUTED})], align=PP_ALIGN.RIGHT)
        # 増分（線の上に小さく）
        txt((px(old) + px(new)) / 2 - 30, y - 12, 60, 14,
            [(("+%d pt" % (new - old)), {"size": 9.5, "color": dot})],
            align=PP_ALIGN.CENTER)
        # 今の値（大きく）
        txt(786, y - 5, 90, 28,
            [("%d%%" % new, {"size": 21 if moved else 17, "bold": True,
                             "color": col})])
        y += 42
    y += 6

rect(70, y - 2, 820, 1.0, HAIR)

# ---- 読み取りと弱点（囲まず、罫線と余白だけで分ける） ----
ty = y + 14
txt(88, ty, 380, 40,
    [[("読み取り", {"size": 10.5, "color": PUR, "bold": True})],
     [("方位の変化率を使う経路を足したことが、両方の型で効いた。",
       {"size": 11, "color": SUB})]], line=1.45)
rect(496, ty + 2, 0.8, 36, HAIR)
txt(520, ty, 370, 40,
    [[("残っている弱点", {"size": 10.5, "color": MUTED, "bold": True})],
     [("すり抜けていく安全な車での無駄鳴りは、ほぼ変わっていない。",
       {"size": 11, "color": SUB})]], line=1.45)

txt(88, ty + 52, 800, 16,
    [("合成データ・調整用の検証セット（1,800本）での値。確定評価とは別枠。",
      {"size": 9, "color": MUTED})])

# ---- フッター ----
txt(58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
steps = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]
x = 300.0
ws = [16 + len(s) * 12.0 for s in steps]
x = (W - (sum(ws) + 8 * (len(steps) - 1))) / 2
for s, wd in zip(steps, ws):
    active = (s == "振り返り")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    label(c, s, size=10.5, color=WHITE if active else MUTED, bold=active)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("9", {"size": 11, "color": MUTED})],
    align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"修正_その結果_案A_2026-09-15.pptx")
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
