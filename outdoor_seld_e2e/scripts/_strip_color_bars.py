# -*- coding: utf-8 -*-
"""できあがった .pptx から「カード上端の色帯」だけを外し、色を見出しの文字へ移す。

本人の指示（2026-09-05）:
 「そんな装飾なくしてシンプルにしたくない気持ちもある。**色帯だけを変えたい、
  それ以外はそのまま**」

やること（この3つだけ。ほかは一切触らない）:
 1. 高さ 3.5pt・幅128pt以上の塗りつぶし長方形 ＝ 色帯 を削除
 2. その帯のすぐ下にある見出しのテキストを、帯と同じ色にする
    （色帯が持っていた「どのカードか」の識別を、文字の色が引き継ぐ）
 3. 全幅（600pt超）の帯は、下に見出しが1つに定まらないので削除のみ

巻き込み防止: 高さは 3.5pt ちょうど（44450 EMU）で判定する。
 道路の帯（h=6.0）・軌跡の線（h=2.4）・白線（h=3.0）は対象外になる。

使い方:
  python scripts/_strip_color_bars.py <入力.pptx> [出力.pptx]
出力を省略すると <入力>_色帯なし.pptx を同じ場所に作る。原本は変更しない。
"""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

BAR_H_EMU = int(Pt(3.5))        # 44450
MIN_W_PT = 128.0                # これ未満は帯とみなさない
FULLWIDTH_PT = 600.0            # これ以上は「箱全体の帯」＝色を移す先が定まらない
TITLE_DX_PT = 40.0              # 見出しは帯の左端から右へこの範囲
TITLE_DY_PT = 40.0              # 見出しは帯の上端から下へこの範囲


def fill_rgb(sh):
    try:
        if sh.fill.type == 1:                     # MSO_FILL.SOLID
            return sh.fill.fore_color.rgb
    except Exception:
        pass
    return None


def find_bars(slide):
    bars = []
    for sh in slide.shapes:
        if sh.height != BAR_H_EMU:
            continue
        if sh.width / 12700.0 < MIN_W_PT:
            continue
        rgb = fill_rgb(sh)
        if rgb is None:
            continue
        bars.append((sh, rgb))
    return bars


def find_title(slide, bar):
    """帯のすぐ下・左端がほぼ揃っているテキスト枠のうち、一番上のものを見出しとみなす。"""
    bx, by = bar.left / 12700.0, bar.top / 12700.0
    best, best_y = None, None
    for sh in slide.shapes:
        if not sh.has_text_frame or sh is bar:
            continue
        if not sh.text_frame.text.strip():
            continue
        x, y = sh.left / 12700.0, sh.top / 12700.0
        if bx - 4 <= x <= bx + TITLE_DX_PT and by < y <= by + TITLE_DY_PT:
            if best_y is None or y < best_y:
                best, best_y = sh, y
    return best


def recolor(tb, rgb):
    """先頭段落の太字ラン（＝見出し）を帯の色にする。太字が無ければ全ランを塗る。"""
    p = tb.text_frame.paragraphs[0]
    runs = [r for r in p.runs if r.font.bold] or list(p.runs)
    for r in runs:
        r.font.color.rgb = rgb
    return len(runs)


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 2
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2]) if len(sys.argv) > 2 else \
        src.with_name(src.stem + "_色帯なし" + src.suffix)
    shutil.copyfile(src, dst)

    prs = Presentation(str(dst))
    n_bar = n_col = 0
    for i, slide in enumerate(prs.slides, 1):
        for bar, rgb in find_bars(slide):
            w = bar.width / 12700.0
            note = ""
            if w < FULLWIDTH_PT:
                tb = find_title(slide, bar)
                if tb is not None:
                    k = recolor(tb, rgb)
                    n_col += 1
                    note = "→ 見出し「%s」を #%02X%02X%02X に（%dラン）" % (
                        tb.text_frame.text.split("\n")[0][:22], *tuple(rgb), k)
                else:
                    note = "→ 見出しが見つからず削除のみ"
            else:
                note = "→ 全幅の帯。削除のみ"
            bar._element.getparent().remove(bar._element)
            n_bar += 1
            print("p%-3d 帯 w=%6.1f  %s" % (i, w, note))
    prs.save(str(dst))
    print("\n帯を %d 本削除、見出し %d 個を着色。" % (n_bar, n_col))
    print("保存:", dst)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
