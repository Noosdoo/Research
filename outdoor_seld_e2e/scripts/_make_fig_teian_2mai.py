# -*- coding: utf-8 -*-
"""提案手法を2枚に割ったスライド（①2層の構成 / ②通知層）を生成。

背景: 2026-08-25。教員の指摘で「合成データの自作は提案手法ではなく実験の設計」と
なり、そのスライドを検証側へ移した結果、提案手法が1枚だけになって薄く見える問題が出た。
本人「さすがに提案手法のスライドが1枚しかないのは薄い」。
そこで**中身を水増しせず**、詰まりすぎていた現p8（2層）を2枚に割る。
話す量は同じなので時間はほぼ増えず、CPA図を大きく使えるぶん説明しやすくなる。

規則の正 = scripts/step12_notify_v4_ttc.py（v4.1・最接近予測）
  強 = (d_cpa≤1.0m かつ t_cpa≤2.5s) が4フレーム ∪ (推定距離≤1.5m) が2フレーム（保険）
  中 = (d_cpa≤2.0m かつ t_cpa≤4.0s) が4フレーム ∪ (推定距離≤3.2m) が2フレーム（保険）
  対象は DIST_CLASSES={car,kick,bike} の3クラスのみ。同一物体の判定は前フレーム±60°の近似
出力: md/seminar/図_提案手法2枚_2026-08-25.pptx（コピペ用・2枚）
"""
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
FAINT = RGBColor(0xF4, 0xF5, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
SL = None


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = SL.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def shape(x, y, w, h, fill, line=None, lw=1.0, kind=MSO_SHAPE.RECTANGLE, rot=None):
    sh = SL.shapes.add_shape(kind, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if rot is not None:
        sh.rotation = rot
    return sh


def conn(x1, y1, x2, y2, color, lw=1.0, dash=None):
    c = SL.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Pt(x1), Pt(y1), Pt(x2), Pt(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    if dash:
        c.line.dash_style = dash
    return c


def skeleton(title, active, page):
    shape(44, 14, 1.4, H - 28, PUR)
    shape(58, 46, 30, 6, PUR)
    txt(98, 30, W - 160, 40, [(title, {"size": 23, "bold": True, "spc": 250})])
    shape(44, 90, W - 90, 1.2, INK)
    txt(58, 506, 90, 20, [("2026/08/30", {"size": 11, "color": MUTED})])
    x = 300.0
    for s in ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]:
        wd = 16 + len(s) * 12.0
        on = (s == active)
        c = shape(x, 502, wd, 22, NAVY if on else CHIP)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        tf.margin_top = tf.margin_bottom = Pt(1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = s
        r.font.size = Pt(10.5)
        r.font.bold = on
        r.font.color.rgb = WHITE if on else MUTED
        r.font.name = "Meiryo"
        meiryo(r)
        x += wd + 8
    txt(W - 90, 506, 40, 20, [(page, {"size": 11, "color": MUTED})],
        align=PP_ALIGN.RIGHT)


# ============================================================
# 1枚目: 提案手法① ― 2層の構成
# ============================================================
SL = prs.slides.add_slide(prs.slide_layouts[6])
skeleton("提案手法①　―　2層の構成", "提案手法", "8")

txt(70, 102, W - 140, 24,
    [("音を聞いて", {"size": 13.5, "color": SUB}),
     ("「何が・どこに」を出す層", {"size": 13.5, "bold": True}),
     ("と、それを", {"size": 13.5, "color": SUB}),
     ("「いつ・どの強さで伝えるか」に変える層", {"size": 13.5, "bold": True}),
     ("に分けた。", {"size": 13.5, "color": SUB})])

# 入力
shape(70, 140, 150, 40, FAINT, line=LINE, lw=1.0)
txt(70, 150, 150, 22, [("入力：FOA 4ch", {"size": 12.5, "bold": True})],
    align=PP_ALIGN.CENTER)
conn(145, 180, 145, 198, LINE, lw=1.2)

LAYERS = [
    (PUR, "知覚層（SELD＋SDE）", "本研究で距離の出力を1軸追加",
     ["8クラスの検出・方向・距離を0.1秒ごとに推定"],
     "種類 ・ 方位 ・ 距離"),
    (GOLD, "通知層（3段階）", "本研究の中核",
     ["このまま進むとどこまで近づくかを予測し、",
      "至近警告 / 注意 / 抑制に振り分ける"],
     "至近警告 ・ 注意 ・ 抑制"),
    (MUTED, "首元振動デバイス", "構想・未実装",
     ["鳴る位置で方向を、鳴り方で危険度を伝える"],
     "首元の振動"),
]
y = 198
for i, (col, name, tag, body, out) in enumerate(LAYERS):
    bh = 90 if len(body) == 2 else 76
    shape(70, y, W - 140, bh, WHITE, line=LINE, lw=1.0)
    shape(70, y, 5, bh, col)
    txt(92, y + 12, 330, 22, [(name, {"size": 15, "bold": True})])
    txt(92 + 5, y + 36, 400, 20 * len(body),
        [[(t, {"size": 12, "color": SUB})] for t in body], line=1.35)
    # 右: 出す情報
    shape(600, y + 14, 250, bh - 28, FAINT)
    txt(600, y + 20, 250, 18,
        [("出す情報", {"size": 10, "color": MUTED})], align=PP_ALIGN.CENTER)
    txt(600, y + 38, 250, 22,
        [(out, {"size": 12.5, "bold": True, "color": col if col is not MUTED else SUB})],
        align=PP_ALIGN.CENTER)
    txt(430, y + 13, 160, 18, [(tag, {"size": 10, "color": MUTED})])
    if i < 2:
        conn(145, y + bh, 145, y + bh + 18, LINE, lw=1.2)
        shape(139, y + bh + 14, 12, 9, LINE, kind=MSO_SHAPE.ISOSCELES_TRIANGLE,
              rot=180)
    y += bh + 22

shape(70, 452, 6, 44, GOLD)
shape(76, 452, W - 146, 44, WHITE, line=LINE, lw=1.0)
txt(94, 460, W - 180, 30,
    [("危険のたびに全部鳴らすと使い物にならない。", {"size": 13.5, "color": SUB}),
     ("通知の頻度そのものを設計対象にしている", {"size": 13.5, "bold": True})],
    anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 2枚目: 提案手法② ― 通知層
# ============================================================
SL = prs.slides.add_slide(prs.slide_layouts[6])
skeleton("提案手法②　―　通知層（最接近の予測で出し分け）", "提案手法", "9")

txt(70, 102, W - 140, 24,
    [("距離のしきい値ではなく、", {"size": 13.5, "color": SUB}),
     ("このまま進むとどこまで近づくか（最接近）を予測", {"size": 13.5, "bold": True}),
     ("して段を決める。", {"size": 13.5, "color": SUB})])

# ---- 左: 3段階 ----
LX, LW = 70, 400
shape(LX, 138, LW, 172, WHITE, line=GOLD, lw=1.4)
rows = [
    ("至近警告", RED, "予測最接近 1.0m以内 ・ 2.5秒以内", "が4フレーム"),
    ("注意", GOLD, "予測最接近 2.0m以内 ・ 4秒以内", "が4フレーム"),
    ("抑制", GREEN, "横を通り過ぎるだけの対象", "は鳴らさない"),
]
ry = 156
for lab, col, cond, tail in rows:
    txt(LX + 18, ry, 84, 20, [(lab, {"size": 13, "bold": True, "color": col})])
    txt(LX + 106, ry + 1, LW - 124, 34,
        [[(cond, {"size": 11.5, "bold": True})],
         [(tail, {"size": 11, "color": SUB})]], line=1.28)
    ry += 52

txt(LX + 4, 322, LW - 8, 96,
    [[("※ 予測が出せないときの保険", {"size": 11, "bold": True, "color": SUB})],
     [("　 推定距離 1.5m（至近警告）/ 3.2m（注意）を切ったら2フレームで通知",
       {"size": 10.5, "color": SUB})],
     [("※ 距離を使うのは 車の走行音・キックボード・バイク の3クラス",
       {"size": 10.5, "color": SUB})],
     [("　 サイレンなど残り5クラスは距離を使わず、検出したら通知",
       {"size": 10.5, "color": SUB})],
     [("※ 同じ物体かは前フレームとの方位±60°で近似判定（物体追跡ではない）",
       {"size": 10.5, "color": MUTED})]], line=1.42)

# ---- 右: CPA図 ----
FX, FY, FW, FH = 496, 138, W - 140 - (496 - 70), 300
shape(FX, FY, FW, FH, WHITE, line=LINE, lw=1.0)
txt(FX, FY + 12, FW, 22,
    [("同じ「近づいてくる」でも、鳴らすべき相手は片方だけ",
      {"size": 12.5, "bold": True})], align=PP_ALIGN.CENTER)

PX, PY = FX + 96, FY + 236          # 歩行者
for r in (30.0, 58.0):
    c = shape(PX - r, PY - r, 2 * r, 2 * r, None, line=LINE, lw=1.0,
              kind=MSO_SHAPE.OVAL)
    c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
shape(PX - 5, PY - 5, 10, 10, NAVY, kind=MSO_SHAPE.OVAL)
txt(PX - 50, PY + 12, 100, 18,
    [("歩行者", {"size": 10.5, "color": SUB})], align=PP_ALIGN.CENTER)

# 緑: 横を通り過ぎる（方位が横に流れる）
GY = FY + 66
conn(FX + 34, GY, FX + FW - 24, GY, GREEN, lw=2.0)
shape(FX + 28, GY - 6, 12, 12, GREEN, kind=MSO_SHAPE.ISOSCELES_TRIANGLE, rot=270)
for gx in (FX + 150, FX + 300):
    shape(gx - 5, GY - 5, 10, 10, GREEN, kind=MSO_SHAPE.OVAL)
    conn(PX, PY, gx, GY, LINE, lw=0.8, dash=MSO_LINE_DASH_STYLE.ROUND_DOT)
txt(FX + 120, GY - 32, FW - 130, 18,
    [("横を通り過ぎる → 鳴らさない", {"size": 11.5, "bold": True, "color": GREEN})],
    align=PP_ALIGN.RIGHT)
txt(FX + 20, GY + 8, 160, 18,
    [("方位が横に流れる", {"size": 10, "color": GREEN})])

# 赤: 正面に来る（方位が変わらない）＝歩行者と同じ高さ
RY2 = PY
conn(PX + 76, RY2, FX + FW - 24, RY2, RED, lw=2.0)
shape(PX + 70, RY2 - 6, 12, 12, RED, kind=MSO_SHAPE.ISOSCELES_TRIANGLE, rot=270)
for rx in (PX + 150, PX + 260):
    shape(rx - 5, RY2 - 5, 10, 10, RED, kind=MSO_SHAPE.OVAL)
txt(FX + 120, RY2 - 34, FW - 130, 18,
    [("正面に来る → 至近警告", {"size": 11.5, "bold": True, "color": RED})],
    align=PP_ALIGN.RIGHT)
txt(FX + 120, RY2 + 10, FW - 130, 18,
    [("方位が変わらないまま近づく", {"size": 10, "color": RED})],
    align=PP_ALIGN.RIGHT)

txt(FX, FY + FH - 26, FW, 18,
    [("点線の円 ＝ 保険の距離しきい値", {"size": 10, "color": MUTED})],
    align=PP_ALIGN.CENTER)

shape(496, 452, 6, 44, PUR)
shape(502, 452, W - 70 - 502, 44, WHITE, line=LINE, lw=1.0)
txt(518, 460, W - 100 - 502, 30,
    [("方位が変わらないまま近づく相手＝", {"size": 12.5, "color": SUB}),
     ("衝突コース", {"size": 12.5, "bold": True, "color": RED})],
    anchor=MSO_ANCHOR.MIDDLE)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"図_提案手法2枚_2026-08-25.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
