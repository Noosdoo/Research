# -*- coding: utf-8 -*-
"""p15「どこで、どんな条件で録るのか」の「記録」の行を、1行に収まる現行手順の文言に直す（2026-09-07）。

現状（本人が一括修正を適用後に調整した版）: 「記録 :イベントごとに車種・速度・横距離を表に手打ち」
→ 「記録 : 1件ごとに向き・車種・速度・横距離をスマホに手打ち」
  - 「向き」= 象限（記入用CSVの必須列）を足す
  - 「スマホに」= 紙ではなくスマホの表、を1語で
  - 全角換算 28 字。元の行（29 字）が1行に収まっていたので同幅に収まる

本人のデッキから p15 を複製し、該当段落の本文だけ差し替える。書式は元のまま。
出力: md/seminar/修正_記録の行_2026-09-15.pptx（1枚）
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = ROOT / "md/seminar/修正_記録の行_2026-09-15.pptx"
TITLE = "どこで、どんな条件で録るのか"
BODY = "1件ごとに向き・車種・速度・横距離をスマホに手打ち"


def find(prs, title):
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and title in sh.text_frame.text and sh.top / 12700 < 70:
                return sl
    raise SystemExit("見出し「%s」のスライドが見つからない" % title)


def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(copy.deepcopy(shp._element), "p:extLst")
    return dst


def fix_footer(sl):
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == "2026/09/15":
            sh.left, sh.top, sh.width, sh.height = int(Pt(66)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))
        elif t.isdigit() and sh.top / 12700 > 480:
            sh.left, sh.top, sh.width, sh.height = int(Pt(678)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))


def main() -> int:
    src = Presentation(str(SRC))
    out = Presentation()
    out.slide_width, out.slide_height = src.slide_width, src.slide_height
    sl = clone(find(src, TITLE), out)

    panel = [s for s in sl.shapes if s.has_text_frame and "中止の基準" in s.text_frame.text]
    assert len(panel) == 1, "ルールの枠が %d 個" % len(panel)
    paras = [p for p in panel[0].text_frame.paragraphs
             if "".join(r.text for r in p.runs).strip().startswith("記録")]
    assert len(paras) == 1, "「記録」の段落が %d 個" % len(paras)
    p = paras[0]
    assert len(p.runs) >= 2, "run の並びが想定と違う: %r" % [r.text for r in p.runs]
    p.runs[0].text = "記録 "
    p.runs[1].text = ": " + BODY
    for r in list(p.runs[2:]):
        r._r.getparent().remove(r._r)
    fix_footer(sl)

    dst = OUT
    try:
        out.save(str(dst))
    except PermissionError:
        dst = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
        out.save(str(dst))
    print("保存:", dst)
    print("記録 :", BODY)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
