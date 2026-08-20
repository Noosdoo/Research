# -*- coding: utf-8 -*-
"""検証スライド「通知層の結果」を1枚だけ作る（2026-08-20）。

## なぜ要るか

現行デッキは p3「危険な接近だけを3段階で通知する」・p8「通知層の設計」・p13まとめ
と通知層を中核に置いているのに、検証(p9/p10)は検出・方向・距離だけで、
**通知層の数値が1つも出てこない**。これでは「作ったが効くかは不明」に見える。

## 数値の扱い（ここを間違えない）

- 表の値はすべて**確定評価セット・同一採点器**（step12_notify_v33.py・分母508台）。
  v4.1は規則だけ差し替えて採点し直したもの（設計文書 §6-quinquies）
- **キックボード88.7%は距離しきい値のみ(v3.4)の値**。確定評価の通知採点器は車専用なので
  v4.1でのキックボードの値は存在しない。v4.1の列に混ぜてはいけない
- 表の値は**判定時刻より後の音も参照している**。実機を想定した因果推論では75.0%。
  ここを書かずに85.0%だけ出すと、「リアルタイムで動くのか」に答えられなくなる

出力: md/seminar/図_検証_通知層_2026-08-20.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

W, H = 960, 540
INK = RGBColor(0x16, 0x23, 0x3A)
INK2 = RGBColor(0x3A, 0x46, 0x58)
MUTED = RGBColor(0x66, 0x70, 0x7F)
HAIR = RGBColor(0xD9, 0xDA, 0xD2)
AMBERD = RGBColor(0xB3, 0x7E, 0x00)
RED = RGBColor(0xC4, 0x43, 0x2B)
GREEN = RGBColor(0x2F, 0x7D, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHIP = RGBColor(0xF0, 0xF0, 0xEC)
JP = "游ゴシック"

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(int(W * 12700)), Emu(int(H * 12700))
sl = prs.slides.add_slide(prs.slide_layouts[6])


def _font(run, size, bold, color):
    f = run.font
    f.name, f.size, f.bold = JP, Pt(size), bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", JP)


def box(x, y, w, h):
    tb = sl.shapes.add_textbox(Emu(int(x * 12700)), Emu(int(y * 12700)),
                               Emu(int(w * 12700)), Emu(int(h * 12700)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb


def para(tf, runs, size=13, bold=False, color=INK2, align=PP_ALIGN.LEFT,
         first=False, line=1.3, space=2):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment, p.line_spacing, p.space_after = align, line, Pt(space)
    for text, kw in ([(runs, {})] if isinstance(runs, str) else runs):
        r = p.add_run()
        r.text = text
        _font(r, kw.get("size", size), kw.get("bold", bold), kw.get("color", color))
    return p


def rect(x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = sl.shapes.add_shape(shape, Emu(int(x * 12700)), Emu(int(y * 12700)),
                            Emu(int(w * 12700)), Emu(int(h * 12700)))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(10)
    tf.margin_top = tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return s


def label(x, y, w, text, size=10.5, color=MUTED, align=PP_ALIGN.LEFT, bold=False):
    tb = box(x, y, w, 18)
    tb.text_frame.word_wrap = False
    para(tb.text_frame, text, size=size, bold=bold, color=color, align=align,
         first=True, space=0)


# ============ 見出し ============
conn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(54 * 12700),
                               Emu(96 * 12700), Emu((W - 54) * 12700),
                               Emu(96 * 12700))
conn.shadow.inherit = False
conn.line.color.rgb = INK
conn.line.width = Pt(1.5)
tb = box(54, 40, 860, 48)
para(tb.text_frame, "検 証 ― 通 知 層 （ 危 険 な 接 近 を 伝 え ら れ た か ）",
     size=26, bold=True, color=INK, first=True, space=0)

label(54, 104, 860,
      "確定評価セット（学習・検証に未使用の1,800クリップ）／危険な車508台／採点器・分母は同一。"
      "規則だけを差し替えて比較", size=11.5, color=INK2)

# ============ 表 ============
ROWS = [
    ("", "距離しきい値のみ", "最接近の予測", ""),
    ("1.5m以内まで来る車への至近警告", "69.9%", "85.0%", "up"),
    ("警告から最接近までの余裕（中央値）", "0.00秒", "0.80秒", "up"),
    ("余裕が0.5秒以上あった割合", "0.6%", "65.3%", "up"),
    ("安全な車を鳴らさなかった割合", "90.4%", "85.2%", "down"),
    ("安全な車への誤った至近警告", "1.3%", "3.7%", "down"),
]
TX, TY, TW = 54, 132, 596
CW = (250, 168, 178)
rh = 40
for i, (a, b, c, mark) in enumerate(ROWS):
    y = TY + (rh + 4 if i else 0) + (i - 1) * (rh - 4) if i else TY
    y = TY + i * 43
    head = i == 0
    rect(TX, y, TW, 43, fill=(CHIP if head else WHITE), line=HAIR)
    if head:
        for j, (t, w) in enumerate(zip((a, b, c), CW)):
            label(TX + sum(CW[:j]) + 12, y + 13, w - 20, t, size=12.5,
                  color=INK, bold=True,
                  align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        continue
    label(TX + 12, y + 13, CW[0] - 20, a, size=12, color=INK2)
    label(TX + CW[0], y + 12, CW[1], b, size=14, color=MUTED,
          align=PP_ALIGN.CENTER)
    label(TX + CW[0] + CW[1], y + 10, CW[2], c, size=17, bold=True,
          color=(GREEN if mark == "up" else RED), align=PP_ALIGN.CENTER)

label(TX, TY + 6 * 43 + 8, TW,
      "静音なキックボードへの至近警告 88.7%（距離しきい値のみで測定。"
      "確定評価の通知採点は車が対象のため、最接近予測での値は未測定）",
      size=10.5, color=MUTED)

# ============ 右：意味・代償・実時間 ============
RX, RW = 674, 232
c = rect(RX, 132, RW, 96, fill=WHITE, line=GREEN, lw=1.5)
para(c.text_frame, "何が変わったか", size=13, bold=True, color=GREEN, first=True,
     space=4)
para(c.text_frame, [("旧方式は中央値で", {}),
                    ("最接近のちょうどその瞬間", {"bold": True, "color": INK}),
                    ("に鳴っていた＝間に合っていない", {})], size=11.5, line=1.25)

c = rect(RX, 240, RW, 108, fill=WHITE, line=RED, lw=1.5)
para(c.text_frame, "代償", size=13, bold=True, color=RED, first=True, space=4)
para(c.text_frame, "早く鳴らすことと、余計に鳴らさないことは両立しない。"
     "抑制が5ポイント下がり、誤った至近警告が増える", size=11.5, line=1.25)

c = rect(RX, 360, RW, 120, fill=WHITE, line=AMBERD, lw=1.5)
para(c.text_frame, "実時間で動かすと", size=13, bold=True, color=AMBERD,
     first=True, space=4)
para(c.text_frame, [("上の値は判定時刻より", {}),
                    ("後の音も参照", {"bold": True, "color": INK}),
                    ("している。各時刻までの音だけで判定すると", {}),
                    ("75.0%", {"bold": True, "color": INK, "size": 13}),
                    ("（因果学習後）", {})], size=11.5, line=1.25)

# ============ 下：いつ鳴るかの帯（この1枚の要点） ============
LY = 418
rect(54, LY, 596, 62, fill=WHITE, line=HAIR)
ax0, ax1 = 178, 612          # 左端=1.5秒前 / 右端=最接近
AXY = LY + 36


def _t(sec):                 # 0=最接近, 1.5=1.5秒前
    return ax1 - (ax1 - ax0) * sec / 1.5


c2 = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(ax0 * 12700)),
                             Emu(int(AXY * 12700)), Emu(int(ax1 * 12700)),
                             Emu(int(AXY * 12700)))
c2.shadow.inherit = False
c2.line.color.rgb = INK
c2.line.width = Pt(1.5)
for sec in (1.5, 1.0, 0.5, 0.0):
    x = _t(sec)
    cc = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x * 12700)),
                                 Emu(int((AXY - 4) * 12700)),
                                 Emu(int(x * 12700)), Emu(int((AXY + 4) * 12700)))
    cc.shadow.inherit = False
    cc.line.color.rgb = MUTED
    cc.line.width = Pt(1.0)
label(_t(0.0) - 46, AXY + 8, 92, "最接近", size=10, color=INK,
      align=PP_ALIGN.CENTER, bold=True)
label(_t(1.5) - 40, AXY + 8, 92, "1.5秒前", size=10, color=MUTED,
      align=PP_ALIGN.CENTER)
label(60, LY + 8, 120, "いつ鳴るか", size=11.5, color=INK, bold=True)

# 旧方式=最接近と同時
rect(_t(0.0) - 6, AXY - 20, 12, 12, fill=RED, shape=MSO_SHAPE.OVAL)
label(_t(0.0) - 200, AXY - 26, 190, "距離しきい値のみ ＝ 同時", size=11,
      color=RED, align=PP_ALIGN.RIGHT, bold=True)
# 新方式=0.8秒前
rect(_t(0.8) - 6, AXY + 8, 12, 12, fill=GREEN, shape=MSO_SHAPE.OVAL)
label(_t(0.8) - 190, AXY + 26, 186, "最接近の予測 ＝ 0.80秒前", size=11,
      color=GREEN, align=PP_ALIGN.RIGHT, bold=True)

label(54, 494, 860,
      "「全部鳴らす」でも「近づくまで黙る」でもなく、"
      "鳴らす相手と鳴らす時刻を選べることを数値で確認した",
      size=12.5, color=INK, bold=True)

out = ROOT / "md/seminar/図_検証_通知層_2026-08-20.pptx"
try:
    prs.save(out)
except PermissionError:
    out = out.with_name(out.stem + "_新" + out.suffix)
    prs.save(out)
    print("※ 元ファイルが開かれていたため別名で保存した")
print(f"saved: {out} ({out.stat().st_size // 1024} KB)")
