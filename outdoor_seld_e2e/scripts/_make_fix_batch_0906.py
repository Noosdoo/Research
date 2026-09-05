# -*- coding: utf-8 -*-
"""9/06 のレビューで直すと決めた 8 枚を、本人のデッキから複製して一括で直す。

対象と直す内容（PDF 2026-09-06 04:21 版のレビュー結果）:
  p5  提案手法②   至近警告 1.0m→1.3m・注意 2.0m→1.6m（採用構成 cs1.3/cm1.6。1.0/2.0 は v4.1 の値）
  p8  ①その結果   脚注に「正解の軌跡で」を戻し、1行に収まるよう 14pt に
  p11 ④触覚デモ   「本物のモデル出力5本と自作28本」→「62場面。通知は全て本物のモデル出力」
  p12 ④同じ場面   脚注を追加: 通知は本物の検出層の出力（オラクル版ではない）・通知規則 v4.3
  p14 どんな音を  「210クリップ」→「210本」、D「固定・機会」→「警告音」、
                  「サイレン等の遭遇待ち8」→「サイレンの遭遇待ち4 / クラクションの統制4」
  p15 どこで      装着「全テイク静止」→「静止で統一、歩行対比だけ歩行」、記録「紙」→「スマホの表」
  p16 収録の流れ  9/05 決定に更新: 儀式=騒音計60秒（初日だけ正面で手拍子1打）、
                  本番=録音とストップウォッチ同時開始→真横でラップ、直後=スマホに手打ち、
                  帰宅後=較正・ラップから切り出し・表をCSVへ
  p17 機材相談    「唯一の機種」→「一体型で満たす現実的な機種」、
                  96kHz の理由を「40kHz帯まで原本に残せる（探索用。検出は主張しない）」に
                  （統制超音波評価は撤回済み・監査 R08）

p9（p8 の古い複製）と p13（Unity の静止画）はここでは扱わない（削除／差し替えは本人）。
図形は位置ではなく本文で探し、想定と違えば止まる。書式は元のまま。
出力: md/seminar/修正_9-06一括_2026-09-15.pptx（8枚・デッキ順）
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = ROOT / "md/seminar/修正_9-06一括_2026-09-15.pptx"
MUTED = RGBColor(0x8A, 0x8F, 0x9A)


# ---------------------------------------------------------------- helpers
def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(copy.deepcopy(shp._element), "p:extLst")
    return dst


def slide_by_title(prs, needle, nth=0):
    hits = []
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and needle in sh.text_frame.text and sh.top / 12700 < 70:
                hits.append(sl)
                break
    assert len(hits) > nth, "見出し「%s」が %d 枚（%d枚目が欲しい）" % (needle, len(hits), nth + 1)
    return hits[nth]


def shape_with(sl, needle):
    hits = [s for s in sl.shapes if s.has_text_frame and needle in s.text_frame.text]
    assert len(hits) == 1, "「%s」を含む図形が %d 個" % (needle, len(hits))
    return hits[0]


def run_with(sp, needle):
    hits = [r for p in sp.text_frame.paragraphs for r in p.runs if needle in r.text]
    assert len(hits) == 1, "run「%s」が %d 個" % (needle, len(hits))
    return hits[0]


def set_para(sp, idx, text):
    """段落 idx の本文を差し替える（1つ目の run の書式を保つ）。"""
    p = sp.text_frame.paragraphs[idx]
    p.runs[0].text = text
    for r in list(p.runs[1:]):
        r._r.getparent().remove(r._r)


def add_para_like(sp, text, like=0):
    """段落 like の1つ目の run と同じ書式で段落を足す。"""
    src = sp.text_frame.paragraphs[like].runs[0]
    p = sp.text_frame.add_paragraph()
    p.alignment = sp.text_frame.paragraphs[like].alignment
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.name = src.font.size, src.font.bold, src.font.name
    try:
        r.font.color.rgb = src.font.color.rgb
    except Exception:
        pass
    for ea in src._r.findall(qn("a:rPr") + "/" + qn("a:ea")):
        r._r.get_or_add_rPr().append(copy.deepcopy(ea))
    return p


def note_box(sl, x, y, w, h, text, size=12):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size, r.font.name = Pt(size), "Meiryo"
    r.font.color.rgb = MUTED
    ea = r._r.get_or_add_rPr().makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    r._r.get_or_add_rPr().append(ea)
    return tb


def fix_footer(sl):
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == "2026/09/15":
            sh.left, sh.top, sh.width, sh.height = int(Pt(66)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))
        elif t.isdigit() and sh.top / 12700 > 480:
            sh.left, sh.top, sh.width, sh.height = int(Pt(678)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))


# ---------------------------------------------------------------- fixes
def fix_p5(sl):
    sp = shape_with(sl, "予測最接近")
    assert run_with(sp, "1.0m") and run_with(sp, "2.0m")
    run_with(sp, "1.0m").text = "1.3m"
    run_with(sp, "2.0m").text = "1.6m"


def fix_p8(sl):
    sp = shape_with(sl, "型の分け方")
    assert "正解の軌跡" not in sp.text_frame.text
    r = run_with(sp, "最接近の")
    r.text = "正解の軌跡で、最接近の"
    for p in sp.text_frame.paragraphs:                 # 1行に収める
        for rr in p.runs:
            rr.font.size = Pt(14)


def fix_p11(sl):
    sp = shape_with(sl, "鳴らす中身")
    paras = sp.text_frame.paragraphs
    p2 = [p for p in paras if "鳴らす中身" in "".join(r.text for r in p.runs)][0]
    p3 = [p for p in paras if "自作した歩道のシナリオ" in "".join(r.text for r in p.runs)][0]
    p2.runs[1].text = "：62場面（歩道・実地図・評価用）"
    for r in list(p2.runs[2:]):
        r._r.getparent().remove(r._r)
    keep = [r for r in p3.runs if r.text.strip() == ""]          # 字下げの空白は残す
    body = [r for r in p3.runs if r.text.strip() != ""]
    body[0].text = "通知は全て本物のモデル出力"
    for r in body[1:]:
        r._r.getparent().remove(r._r)
    assert keep, "字下げの空白 run が無い"


def fix_p12(sl):
    assert not any(s.has_text_frame and "オラクル" in s.text_frame.text for s in sl.shapes)
    note_box(sl, 74, 474, 820, 18,
             "※ 通知は本物の検出層の出力（正解の位置から作ったオラクル版ではない）。通知規則は v4.3。")


def fix_p14(sl):
    sp = shape_with(sl, "屋外で計210")
    run_with(sp, "クリップ").text = "本"
    run_with(sp, "分").text = ""
    run_with(shape_with(sl, "固定・機会"), "固定・機会").text = "警告音"
    run_with(shape_with(sl, "サイレン等の遭遇待ち8"), "サイレン等の遭遇待ち8").text = \
        "サイレンの遭遇待ち4 / クラクションの統制4"


def fix_p15(sl):
    sp = shape_with(sl, "全テイク静止")
    run_with(sp, "全テイク静止").text = "ヘルメット。静止で統一、歩行対比だけ歩行"
    r1 = run_with(sp, ": 1")
    r1.text = ": "
    run_with(sp, "テイクごとに車種").text = "イベントごとに車種・速度・横距離をスマホの表に手打ち"


def fix_p16(sl):
    run_with(shape_with(sl, "全210回で共通"), "全210回で共通").text = \
        "全地点で共通。④だけを通るたびに繰り返す。"
    sp = shape_with(sl, "定位置に立つ")
    add_para_like(sp, "水準器を確認")
    sp.height = int(Pt(38.7))
    set_para(shape_with(sl, "90秒"), 0, "60秒")
    sp = shape_with(sl, "手拍子・騒音計・ベル4方位")
    set_para(sp, 0, "騒音計と並べて")
    add_para_like(sp, "無言で60秒")
    sp = shape_with(sl, "1通過＝1テイク")
    set_para(sp, 0, "無言で待つ")
    sp = shape_with(sl, "声で記録")
    set_para(sp, 0, "5秒待って停止")
    set_para(sp, 1, "スマホに手打ち")
    run_with(shape_with(sl, "儀式の中身"), "儀式の中身").text = "儀式 と ④ 本番のやり方"
    sp = shape_with(sl, "手拍子を1回")
    set_para(sp, 0, "騒音計と並べて無言で60秒")
    set_para(sp, 1, "何デシベルだったかの基準を取る")
    # 上で 40 を書き換えた直後なので、元の 42（同じ本文）を区別して取る
    cands = [s for s in sl.shapes if s.has_text_frame and s.text_frame.text.startswith("騒音計と並べて無言で60秒")]
    assert len(cands) == 2
    sp42 = max(cands, key=lambda s: s.top)
    set_para(sp42, 0, "初日だけ、正面で手拍子1打")
    set_para(sp42, 1, "マイクの正面のズレを度で読む。±10°超は補正")
    sp = shape_with(sl, "ベルを前後左右から1打ずつ")
    set_para(sp, 0, "録音とストップウォッチを同時に開始")
    set_para(sp, 1, "真横に来たらラップ。声は出さない")
    sp = shape_with(sl, "紙の記録を表に打ち込む")
    set_para(sp, 0, "録音を変換し、騒音計の値で音量を較正する")
    set_para(sp, 1, "波形とラップから最接近を決め、10秒に切り出す")
    set_para(sp, 2, "スマホの表を注釈CSVにまとめる")


def fix_p17(sl):
    sp = shape_with(sl, "唯一の機種")
    run_with(sp, "唯一の機種").text = "本体。右の条件を一体型で満たす現実的な機種"
    sp = shape_with(sl, "40kHz")
    p = [p for p in sp.text_frame.paragraphs if "40kHz" in "".join(r.text for r in p.runs)][0]
    p.runs[0].text = "40kHz帯まで原本に残せる（探索用。超音波の検出は主張しない）。"
    for r in list(p.runs[1:]):
        r._r.getparent().remove(r._r)


PLAN = [  # (見出し, 何枚目のその見出しか, 修正関数)
    ("提案手法②", 0, fix_p5),
    ("①その結果", 0, fix_p8),
    ("④触覚デモを動かした", 0, fix_p11),
    ("④同じ場面で、鳴る車と鳴らない車", 0, fix_p12),
    ("どんな音を、どれだけ録るのか", 0, fix_p14),
    ("どこで、どんな条件で録るのか", 0, fix_p15),
    ("収録の流れ", 0, fix_p16),
    ("機材購入の相談", 0, fix_p17),
]


def main() -> int:
    src = Presentation(str(SRC))
    out = Presentation()
    out.slide_width, out.slide_height = src.slide_width, src.slide_height
    for title, nth, fn in PLAN:
        sl = clone(slide_by_title(src, title, nth), out)
        fn(sl)
        fix_footer(sl)
        print("済:", title)
    dst = OUT
    try:
        out.save(str(dst))
    except PermissionError:
        dst = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
        out.save(str(dst))
    print("保存:", dst)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
