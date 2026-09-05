# -*- coding: utf-8 -*-
"""「①すり抜けと対向を、どう区別するか」の左の図を描き直す（2026-09-06 本人指摘）。

旧図は 対向型=縦↓・すり抜け型=横→ と直交させていたため「違いは進む向き」に見えた。
実際の違いは**進路が装着者からどれだけ横にずれているか**だけで、どちらも同じ道路を
同じ向きに走る。よって2本の線を平行（どちらも↓）に描き、横のずれだけ変える。

副産物: 緑の車への方位線を等しい時間間隔で引くと、遠いうちはほぼ平行（方位が動かない
ように見える）、近づくと扇が一気に開く。これが下段「なぜ15mで切るのか」（dθ/dt ∝ 1/d²）
の絵による説明になる。

あわせて、本人がPowerPoint上で済ませた右パネルの修正（未保存の可能性があるため）も
同じ内容で当てる: 「救済ルートを1本足した」→「追加ルート(4フレーム)」、
②の小見出し「頑健な傾きがマイナス」→「直近0.5秒の距離の傾き(最小二乗)がマイナス」
（採用構成 brg5+mn4/4+rc(0.10,15)+link+cs1.3/cm1.6 は robust_slope=False。
 step12_notify_v43.py L60 は v4.closing_speed=最小二乗を直接呼ぶ）。
キャプションは前回の修正どおり「判定は毎フレーム、直近0.5秒の方位の傾きで行う」。

本人のデッキから p7 を複製し、左の図の中身だけ消して描き直す。書式は元のまま。
出力: md/seminar/修正_すり抜けと対向_v2_2026-09-15.pptx（1枚）
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = ROOT / "md/seminar/修正_すり抜けと対向_v2_2026-09-15.pptx"
TITLE = "すり抜けと対向を、どう区別するか"

RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)

CAPTION_NEW = "判定は毎フレーム、直近0.5秒の方位の傾きで行う（最接近がいつかは知らない）"
PANEL_TITLE_NEW = "追加ルート(4フレーム)"
COND2_SUB_NEW = "直近0.5秒の距離の傾き(最小二乗)がマイナス"

# 図の寸法（装着者の中心 (OX, OY) からの相対値・pt）
DX_GREEN = 58.0                                   # すり抜け型の横ずれ
DY_MARKS = (160, 132, 104, 76, 48, 20)            # 等しい時間間隔の車の位置
Y_LINE_TOP = 164.0                                # 線の上端（カード見出しの直下）


# ---------------------------------------------------------------- helpers
def find(prs, title):
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and title in sh.text_frame.text and sh.top / 12700 < 70:
                return sl
    raise SystemExit("見出し「%s」のスライドが見つからない" % title)


def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(copy.deepcopy(shp._element), "p:extLst")
    return dst


def meiryo(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)


def txt(sl, x, y, w, h, paras, align=PP_ALIGN.LEFT, line=None):
    """paras = [[(text, {size,bold,color}), ...], ...]"""
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in runs:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(o.get("size", 10))
            r.font.bold = o.get("bold", False)
            r.font.color.rgb = o.get("color", INK)
            r.font.name = "Meiryo"
            meiryo(r)
    return tb


def rect(sl, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, rot=0.0):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.rotation = rot
    return sh


def dashed(sl, x1, y1, x2, y2, color):
    ln = sl.shapes.add_connector(1, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(0.8)
    ln.line.dash_style = 4
    return ln


def set_para(sp, idx, text):
    p = sp.text_frame.paragraphs[idx]
    p.runs[0].text = text
    for r in list(p.runs[1:]):
        r._r.getparent().remove(r._r)


def pick(sl, needle):
    hits = [s for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip().startswith(needle)]
    assert len(hits) == 1, "「%s」が %d 個ある" % (needle, len(hits))
    return hits[0]


def fix_footer(sl):
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == "2026/09/15":
            sh.left, sh.top, sh.width, sh.height = int(Pt(66)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))
        elif t.isdigit() and sh.top / 12700 > 480:
            sh.left, sh.top, sh.width, sh.height = int(Pt(678)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))


# ---------------------------------------------------------------- main
def redraw(sl):
    # --- 目印になる図形を特定 ---
    card_title = pick(sl, "装着者から見た")
    wearer_lbl = [s for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip() == "装着者"]
    assert len(wearer_lbl) == 1, "「装着者」ラベルが %d 個" % len(wearer_lbl)
    wearer_lbl = wearer_lbl[0]
    caption = [s for s in sl.shapes if s.has_text_frame and
               (s.text_frame.text.startswith("判定に使うのは") or s.text_frame.text.startswith("判定は毎フレーム"))]
    assert len(caption) == 1, "キャプションが %d 個" % len(caption)
    caption = caption[0]
    cards = [s for s in sl.shapes if not s.has_text_frame or not s.text_frame.text.strip()]
    cards = [s for s in cards if s.left / 12700 < 120 and s.width / 12700 > 380 and s.height / 12700 > 200]
    assert len(cards) == 1, "左のカードが %d 個" % len(cards)
    card = cards[0]
    ovals = []
    for s in sl.shapes:
        try:
            if "OVAL" in str(s.auto_shape_type) and abs(s.width / 12700 - 16) < 1:
                ovals.append(s)
        except Exception:
            pass
    assert len(ovals) == 1, "装着者の丸が %d 個" % len(ovals)
    dot = ovals[0]
    OX, OY = dot.left / 12700 + 8, dot.top / 12700 + 8
    CL, CT = card.left / 12700, card.top / 12700
    CR, CB = CL + card.width / 12700, CT + card.height / 12700

    # --- カードの中身を消す（枠・見出し・丸・装着者・キャプション以外） ---
    keep = {id(card._element), id(card_title._element), id(dot._element),
            id(wearer_lbl._element), id(caption._element)}
    removed = []
    for s in list(sl.shapes):
        L, T = s.left / 12700, s.top / 12700
        R, B = L + s.width / 12700, T + s.height / 12700
        inside = CL - 2 <= L and R <= CR + 2 and CT - 2 <= T and B <= CB + 2
        if inside and id(s._element) not in keep:
            removed.append("%s L%.0f T%.0f %s" % (str(s.shape_type).split(" ")[0], L, T,
                                                  s.text_frame.text[:14] if s.has_text_frame else ""))
            s._element.getparent().remove(s._element)

    # --- 描き直し ---
    GX = OX + DX_GREEN
    # 対向型: 装着者に向かう線＋矢頭（下向き）
    rect(sl, OX - 1.2, Y_LINE_TOP, 2.4, (OY - 8) - Y_LINE_TOP, RED)
    rect(sl, OX - 7, OY - 8 - 14, 14, 14, RED, MSO_SHAPE.ISOSCELES_TRIANGLE, rot=180)
    # すり抜け型: 平行に横へずらした線＋矢頭（装着者の横を通り過ぎて下へ）
    rect(sl, GX - 1.2, Y_LINE_TOP, 2.4, (OY + 14) - Y_LINE_TOP, GREEN)
    rect(sl, GX - 7, OY + 14, 14, 14, GREEN, MSO_SHAPE.ISOSCELES_TRIANGLE, rot=180)
    # 等しい時間間隔の車の位置と、装着者から見た方位線（緑だけ扇になる）
    for dy in DY_MARKS:
        y = OY - dy
        rect(sl, OX - 2.5, y - 2.5, 5, 5, RED)
        rect(sl, GX - 2.5, y - 2.5, 5, 5, GREEN)
        dashed(sl, OX, OY, GX, y, GREEN)
    # ラベル
    txt(sl, OX - 8 - 140, Y_LINE_TOP, 140, 34,
        [[("対向型", {"size": 11.5, "bold": True, "color": RED})],
         [("方位が変わらない", {"size": 9.5, "color": SUB})]], align=PP_ALIGN.RIGHT, line=1.25)
    txt(sl, GX + 10, Y_LINE_TOP, 140, 34,
        [[("すり抜け型", {"size": 11.5, "bold": True, "color": GREEN})],
         [("方位が速く変わる", {"size": 9.5, "color": SUB})]], line=1.25)
    txt(sl, GX + 10, Y_LINE_TOP + 40, 140, 30,
        [[("遠いうちは緑も動かない", {"size": 9, "color": MUTED})],
         [("→ 15mで切る理由", {"size": 9, "color": MUTED})]], line=1.2)
    txt(sl, CL + 16, OY - 30, 150, 16,
        [[("■ 等しい時間間隔の車の位置", {"size": 8.5, "color": MUTED})]])
    # 装着者ラベルは緑の矢頭とぶつかるので、丸の左に置く
    wearer_lbl.left, wearer_lbl.top = int(Pt(OX - 100)), int(Pt(OY - 10.7))
    wearer_lbl.width, wearer_lbl.height = int(Pt(88)), int(Pt(21.4))
    wearer_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    # キャプション
    set_para(caption, 0, CAPTION_NEW)
    return removed


def fix_right_panel(sl):
    done = []
    for s in sl.shapes:
        if not s.has_text_frame:
            continue
        t = s.text_frame.text
        if t.strip().startswith("救済ルートを1本足した"):
            set_para(s, 0, PANEL_TITLE_NEW); done.append("panel title")
        elif t.strip().startswith("距離が縮んでいる") and "頑健な傾き" in t:
            set_para(s, 1, COND2_SUB_NEW); done.append("cond2 sub")
    return done


def main() -> int:
    src = Presentation(str(SRC))
    out = Presentation()
    out.slide_width, out.slide_height = src.slide_width, src.slide_height
    sl = clone(find(src, TITLE), out)
    removed = redraw(sl)
    panel = fix_right_panel(sl)
    fix_footer(sl)

    dst = OUT
    try:
        out.save(str(dst))
    except PermissionError:
        dst = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
        out.save(str(dst))
    print("保存:", dst)
    print("消した図形 %d 個:" % len(removed))
    for r in removed:
        print("  -", r)
    print("右パネル:", panel or "（既に修正済みのため変更なし）")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
