# -*- coding: utf-8 -*-
"""【デザイン案A】9/15ゼミ 8〜14枚目を、罫線と余白だけで組み直した版。

本人の指示（2026-09-03）:
 「白カード＋色帯＋薄色の囲み＋丸バッジ」の作りが毎回同じで "Claudeすぎる"。
 案A（罫線だけの学会スライド風）で1枚見せた結果 → 全体をこの作りに変える。
 ただし「シンプルすぎてつまらないスライドにはしたくない」。

案Aの決めごと:
 1. 塗りつぶしのカード・色帯・薄色の囲み・丸バッジを使わない
 2. 区切りは細い横罫（0.8pt）と余白だけ。表は上下の太罫＋見出し下の細罫（booktabs式）
 3. 色は 黒 / 灰2段階 / 紫（アクセント1色）のみ。強調は大きさと太さで作る
 4. つまらなくしないため、数字の羅列は「同じ物差しの上に置く図」に置き換える
    - 9枚目: 前→今 のダンベル図
    - 10枚目: 距離の数直線
    - 13枚目: 時間の帯

出力: md/seminar/修正_案A_8-14枚_2026-09-15.pptx（7枚）
元の内容は md/seminar/ゼミ相談スライド_2026-09-15.pptx と同じ。作りだけ変えた。
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x1E, 0x22, 0x2C)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x9A, 0x9E, 0xA8)
FAINT = RGBColor(0xC9, 0xCB, 0xD2)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
PUR = RGBColor(0x7E, 0x6F, 0x98)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
SECTIONS = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(sl, x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(sl, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


def rule(sl, x, y, w, weight=0.8, color=HAIR):
    """横罫。案Aの唯一の区切り手段。"""
    return rect(sl, x, y, w, weight, color)


def vrule(sl, x, y, h, weight=0.8, color=HAIR):
    return rect(sl, x, y, weight, h, color)


def sect(sl, x, y, w, title, note=""):
    """区分の見出し（上に細罫を1本引き、その下に見出し）。"""
    rule(sl, x, y, w)
    runs = [(title, {"size": 12.5, "bold": True})]
    if note:
        runs.append(("　" + note, {"size": 11, "color": SUB}))
    txt(sl, x, y + 10, w, 20, [runs])
    return y + 36


def item(sl, x, y, w, head, note, hsize=11.5, nsize=9, gap=6):
    """項目名＋その下に小さな補足。ダッシュ1本だけを頭につける。"""
    rect(sl, x, y + 7, 7, 1.4, PUR)
    txt(sl, x + 16, y, w - 16, 18, [(head, {"size": hsize})])
    if note:
        txt(sl, x + 16, y + 17, w - 16, 16, [(note, {"size": nsize, "color": MUTED})])
    return y + (17 + nsize + gap if note else 20 + gap)


def new_slide(title, page, active):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 44, 14, 1.4, H - 28, PUR)
    rect(sl, 58, 46, 30, 6, PUR)
    txt(sl, 98, 30, W - 170, 40,
        [(title, {"size": 22, "bold": True, "spc": 250})])
    rect(sl, 44, 90, W - 90, 1.2, INK)
    txt(sl, 58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
    ws = [16 + len(s) * 12.0 for s in SECTIONS]
    x = (W - (sum(ws) + 8 * (len(SECTIONS) - 1))) / 2
    for s, wd in zip(SECTIONS, ws):
        on = (s == active)
        c = rect(sl, x, 502, wd, 22, NAVY if on else CHIP)
        label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
        x += wd + 8
    txt(sl, W - 90, 506, 40, 20,
        [(str(page), {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)
    return sl


def lead(sl, text, size=13):
    txt(sl, 70, 102, 820, 22, [(text, {"size": size, "bold": True})])


# ==================================================== p8 すり抜けと対向
sl = new_slide("す り 抜 け と 対 向 を 、 ど う 区 別 す る か", 8, "振り返り")
lead(sl, "夏ゼミで「三角関数で区別する」とお答えした件。方位の変化率で分ける。")

# 左：軌跡の図（枠を描かず、罫線と線だけで）
FX, FY = 70, 140
txt(sl, FX, FY, 380, 18,
    [("装着者から見た「方向」の変わり方", {"size": 11.5, "bold": True, "color": PUR})])
OX, OY = FX + 200, FY + 218
rect(sl, OX - 6, OY - 6, 12, 12, INK, shape=MSO_SHAPE.OVAL)
txt(sl, OX - 70, OY + 10, 140, 16,
    [("装着者", {"size": 10, "color": SUB})], align=PP_ALIGN.CENTER)
rect(sl, OX - 1.1, FY + 62, 2.2, 148, INK)
rect(sl, OX - 6, FY + 54, 12, 12, INK, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
txt(sl, OX - 116, FY + 44, 104, 32,
    [[("対向型", {"size": 11.5, "bold": True})],
     [("方位が変わらない", {"size": 9, "color": MUTED})]],
    align=PP_ALIGN.RIGHT, line=1.25)
rect(sl, FX + 26, FY + 128, 340, 1.6, PUR)
rect(sl, FX + 332, FY + 122, 12, 12, PUR, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
txt(sl, FX + 214, FY + 90, 170, 32,
    [[("すり抜け型", {"size": 11.5, "bold": True, "color": PUR})],
     [("方位が速く変わる", {"size": 9, "color": MUTED})]], line=1.25)
for dx in (-140, -56, 56, 140):
    ln = sl.shapes.add_connector(1, Pt(OX), Pt(OY), Pt(OX + dx), Pt(FY + 129))
    ln.line.color.rgb = FAINT
    ln.line.width = Pt(0.8)
txt(sl, FX, FY + 244, 380, 16,
    [("判定に使うのは、最接近の 2.5〜1.5 秒前の区間", {"size": 9.5, "color": MUTED})])

vrule(sl, 484, 140, 264)

# 右：3条件
y = sect(sl, 520, 140, 370, "救済ルートを1本足した", "3つが4フレーム続いたら")
for i, (a, b) in enumerate([
    ("方位がほとんど変わらない", "|dθ/dt| ≤ 0.10 rad/s"),
    ("距離が縮んでいる", "頑健な傾きがマイナス"),
    ("15m 以内にいる", "距離ゲート"),
]):
    txt(sl, 520, y, 24, 20, [("%d" % (i + 1), {"size": 13, "bold": True, "color": PUR})])
    txt(sl, 546, y + 1, 344, 18, [(a, {"size": 12})])
    txt(sl, 546, y + 20, 344, 16, [(b, {"size": 9, "color": MUTED})])
    y += 44
rule(sl, 520, y + 2, 370)
txt(sl, 520, y + 12, 370, 34,
    [("船の衝突回避と同じ原則。方位が変わらないまま近づいてくる相手は、衝突コースにいる。",
      {"size": 10.5, "color": SUB})], line=1.4)

rule(sl, 70, 424, 820)
txt(sl, 70, 434, 130, 18,
    [("なぜ15mで切るか", {"size": 10.5, "bold": True, "color": PUR})])
txt(sl, 206, 434, 684, 34,
    [("方位の変化率は距離の2乗に反比例する。遠くをまっすぐ通り過ぎる安全な車も"
      "「方位が変わらない」ように見えるため、距離で切らないと遠い車まで強い警告になる。",
      {"size": 10.5, "color": SUB})], line=1.4)

# ==================================================== p9 その結果
sl = new_slide("そ の 結 果", 9, "振り返り")
lead(sl, "危ない車には強く鳴るようになり、鳴らすべきでない対向車では黙るようになった。")

X0, X1, SY = 268.0, 772.0, 140.0
px = lambda v: X0 + (X1 - X0) * v / 100.0
for v in (0, 25, 50, 75, 100):
    rect(sl, px(v), SY + 8, 0.8, 5, FAINT)
for v, t in ((0, "0"), (50, "50"), (100, "100%")):
    txt(sl, px(v) - 24, SY - 6, 48, 14, [(t, {"size": 8.5, "color": MUTED})],
        align=PP_ALIGN.CENTER)
rule(sl, X0, SY + 13, X1 - X0)

y = SY + 34
for title, metric, rows in [
    ("危ない車（最接近 1.5m 以内）", "強い警告が届いた割合",
     [("対向型", 460, 81, 90, True), ("すり抜け型", 43, 77, 93, True)]),
    ("鳴らすべきでない車（最接近 3.2m 超）", "黙っていられた割合",
     [("対向型", 52, 33, 71, True), ("すり抜け型", 501, 62, 64, False)]),
]:
    y = sect(sl, 70, y, 820, title, metric) + 4
    for name, n, old, new, moved in rows:
        col = INK if moved else MUTED
        dot = PUR if moved else MUTED
        txt(sl, 88, y - 1, 120, 20, [(name, {"size": 12, "color": col})])
        txt(sl, 196, y + 2, 60, 16,
            [("n=%d" % n, {"size": 9, "color": MUTED})], align=PP_ALIGN.RIGHT)
        rect(sl, px(old), y + 8, px(new) - px(old), 2.2, dot)
        rect(sl, px(old) - 5, y + 4, 10, 10, WHITE, line=MUTED, lw=1.1,
             shape=MSO_SHAPE.OVAL)
        rect(sl, px(new) - 5.5, y + 3.5, 11, 11, dot, shape=MSO_SHAPE.OVAL)
        txt(sl, px(old) - 46, y + 14, 40, 14,
            [("%d" % old, {"size": 9, "color": MUTED})], align=PP_ALIGN.RIGHT)
        txt(sl, (px(old) + px(new)) / 2 - 30, y - 12, 60, 14,
            [("+%d pt" % (new - old), {"size": 9.5, "color": dot})],
            align=PP_ALIGN.CENTER)
        txt(sl, 786, y - 5, 90, 28,
            [("%d%%" % new, {"size": 21 if moved else 17, "bold": True, "color": col})])
        y += 42
    y += 6

rule(sl, 70, y - 2, 820)
txt(sl, 88, y + 12, 380, 40,
    [[("読み取り", {"size": 10.5, "color": PUR, "bold": True})],
     [("方位の変化率を使う経路を足したことが、両方の型で効いた。",
       {"size": 11, "color": SUB})]], line=1.45)
vrule(sl, 496, y + 14, 36)
txt(sl, 520, y + 12, 370, 40,
    [[("残っている弱点", {"size": 10.5, "color": MUTED, "bold": True})],
     [("すり抜けていく安全な車での無駄鳴りは、ほぼ変わっていない。",
       {"size": 11, "color": SUB})]], line=1.45)
txt(sl, 88, y + 64, 800, 16,
    [("合成データ・調整用の検証セット（1,800本）での値。確定評価とは別枠。",
      {"size": 9, "color": MUTED})])

# ==================================================== p10 触覚デモ
sl = new_slide("触 覚 デ モ を 動 か し た", 10, "振り返り")
lead(sl, "「Joy-conとUnityで作ってみては」という助言を、そのまま実行した。")

y = sect(sl, 70, 140, 390, "実機で振動した", "9/2")
for a, b in [("SwitchのJoy-con と Unity 6.6", "ハンダごて不要・部品の購入ゼロ"),
             ("本物のモデル出力5本＋自作シナリオ3本", "対向とすり抜けの場面も含む"),
             ("強＝4連打／中＝2発／警告音＝単発", "パルスの数で危険度を表す")]:
    y = item(sl, 70, y, 390, a, b, hsize=11.5, gap=10)

vrule(sl, 492, 140, 190)

txt(sl, 520, 140, 370, 18,
    [("距離だけを変えた3本のシナリオ", {"size": 12.5, "bold": True})])
DX0, DX1, DY = 548.0, 858.0, 208.0
rule(sl, DX0 - 20, DY, (DX1 - DX0) + 50, 1.0, FAINT)
for d, t, on in [(2.0, "強く鳴る（4連打）", True), (2.8, "中くらい（2発）", True),
                 (3.5, "鳴らさない", False)]:
    x = DX0 + (DX1 - DX0) * (d - 2.0) / 1.5
    rect(sl, x - 1, DY - 9, 2, 18, PUR if on else MUTED)
    txt(sl, x - 30, DY - 30, 60, 18,
        [("%.1fm" % d, {"size": 12, "bold": True, "color": PUR if on else MUTED})],
        align=PP_ALIGN.CENTER)
    txt(sl, x - 62, DY + 14, 124, 16,
        [(t, {"size": 9.5, "color": SUB if on else MUTED})], align=PP_ALIGN.CENTER)
txt(sl, 520, DY + 56, 370, 32,
    [("同じ場面で最接近距離だけを変え、鳴り方が切り替わるところを体で確かめた。",
      {"size": 10.5, "color": SUB})], line=1.4)

y = sect(sl, 70, 356, 820, "この先に作る首元デバイス", "構想・未実装")
CX, CY = 232, y + 46
rect(sl, CX - 74, CY - 30, 148, 60, None, line=INK, lw=6.5, shape=MSO_SHAPE.OVAL)
rect(sl, CX - 34, CY - 15, 68, 30, WHITE, line=FAINT, lw=0.8, shape=MSO_SHAPE.OVAL)
txt(sl, CX - 34, CY - 8, 68, 16, [("首", {"size": 10, "color": SUB})],
    align=PP_ALIGN.CENTER)
for dx, dy, on in [(-70, 8, 0), (-46, 26, 0), (0, 32, 0), (46, 26, 1), (70, 8, 0)]:
    rect(sl, CX + dx - 6, CY + dy - 6, 12, 12,
         PUR if on else WHITE, line=PUR if on else MUTED, lw=1.2,
         shape=MSO_SHAPE.OVAL)
txt(sl, CX - 90, CY - 54, 180, 16,
    [("歩く向き（前）", {"size": 9.5, "color": MUTED})], align=PP_ALIGN.CENTER)
txt(sl, 400, y + 22, 490, 52,
    [[("振動子5個と4chマイクを、同じバンドに載せる。", {"size": 12, "bold": True})],
     [("「マイクから見た右」と「体から見た右」が一致するので、"
       "頭の向きを測るセンサーが要らない。", {"size": 10.5, "color": SUB})]], line=1.45)

# ==================================================== p11 どんな音を
sl = new_slide("ど ん な 音 を 、 ど れ だ け 録 る の か", 11, "実録の計画")
lead(sl, "検出対象の8種類を、屋外で計210回ぶん録る。内訳は録る前に固定してある。")

txt(sl, 70, 132, 820, 18,
    [("サイレン　クラクション　バック音　自転車ベル　車　踏切・列車　"
      "キックボード　バイク", {"size": 11, "color": SUB, "spc": 60})])

TY = 162
rule(sl, 70, TY, 820, 1.4, INK)
txt(sl, 84, TY + 8, 160, 18, [("区分", {"size": 11, "bold": True})])
txt(sl, 246, TY + 8, 60, 18,
    [("本数", {"size": 11, "bold": True})], align=PP_ALIGN.RIGHT)
txt(sl, 336, TY + 8, 540, 18, [("中身", {"size": 11, "bold": True})])
rule(sl, 70, TY + 30, 820)
ry = TY + 40
for a, n, b in [
    ("A 走行車", "20", "横2.0〜3.2mを10本（うち後方から6本以上）／横5〜15mを10本"),
    ("B 弱点", "20", "EV・HVの接近5／停車→発進5／自転車の至近通過5／見通し不良の角5"),
    ("C 負例", "20", "静穏6／繁華街・工事6／雨上がり4／風あり4　＝鳴ってはいけない場面"),
    ("D 固定・機会", "20", "実在の踏切8／バック音の統制4／サイレン等の遭遇待ち8"),
    ("E キックボード", "20", "横1.0/1.2/1.5m × 速度2段 × 左右＝12／後方4／前方4"),
    ("F バイク", "10", "交通量の多い道での待ち受け"),
    ("歩行対比", "100", "5条件×10反復を、止まって録る／歩いて録るの対で比較"),
]:
    txt(sl, 84, ry, 160, 18, [(a, {"size": 11.5})])
    txt(sl, 246, ry - 2, 60, 20,
        [(n, {"size": 13, "bold": True, "color": PUR})], align=PP_ALIGN.RIGHT)
    txt(sl, 336, ry, 540, 18, [(b, {"size": 10.5, "color": SUB})])
    ry += 30
rule(sl, 70, ry + 2, 820, 1.4, INK)
txt(sl, 84, ry + 14, 820, 18,
    [[("これに加えて、鳴ってはいけない場面の連続録音を100分以上。", {"size": 11}),
      ("　「1時間あたり何回の誤警告か」を出すために要る。",
       {"size": 11, "color": SUB})]])

# ==================================================== p12 どこで・どんな条件で
sl = new_slide("ど こ で 、 ど ん な 条 件 で 録 る の か", 12, "実録の計画")
lead(sl, "場所を変えないと録れない音があり、天候の条件も決めてある。")

y = sect(sl, 70, 140, 380, "6か所を回る")
for p in ["静かな道（交通量の少ない住宅路）", "交通量の多い道（車・バイク）",
          "繁華街（雑踏＝鳴ってはいけない場面）", "工事現場の近く（高い騒音）",
          "見通しの悪い角（塀のある交差点）", "実在の踏切（時刻表どおり＝確実）"]:
    rect(sl, 70, y + 7, 7, 1.4, PUR)
    txt(sl, 86, y, 364, 18, [(p, {"size": 11, "color": SUB})])
    y += 24

vrule(sl, 484, 140, 200)

y = sect(sl, 520, 140, 370, "録るときの決めごと")
for a, b in [("風速3m/s超と雨の日は録らない", "中止の基準"),
             ("ヘルメットに装着。全テイク静止で統一", "装着"),
             ("家で全部固定し、現地で押すのは録音ボタンだけ", "録音設定"),
             ("1テイクごとに車種・速度・横距離・騒音計の値を紙に書く", "記録")]:
    txt(sl, 520, y, 70, 16, [(b, {"size": 9, "color": MUTED})])
    txt(sl, 520, y + 15, 370, 18, [(a, {"size": 11})])
    y += 40

GY = 372
y2 = sect(sl, 70, GY, 820, "車を記録するときの位置関係", "速度：徐行〜約30km/h")
rule(sl, 210, y2 + 16, 560, 1.0, FAINT)
rule(sl, 210, y2 + 46, 560, 1.0, FAINT)
c = rect(sl, 470, y2 + 20, 48, 22, None, line=INK, lw=1.2)
label(c, "車", size=11, color=INK)
txt(sl, 524, y2 + 20, 30, 20, [("→", {"size": 14, "color": INK})])
rect(sl, 306, y2 + 46, 1.4, 30, PUR)
rect(sl, 300, y2 + 76, 13, 13, INK, shape=MSO_SHAPE.OVAL)
txt(sl, 322, y2 + 50, 420, 18,
    [("横距離 2.0〜3.2m（鳴るべき車）／ 5〜15m（鳴ってはいけない車）",
      {"size": 11, "bold": True, "color": PUR})])
txt(sl, 210, y2 + 92, 220, 16,
    [("装着者（頭上に4方向マイク・静止）", {"size": 9.5, "color": MUTED})])

# ==================================================== p13 収録の流れ
sl = new_slide("収 録 の 流 れ", 13, "実録の計画")
lead(sl, "1テイクあたり6つの手順。この骨格を全210回で共通にする。")

STEPS = [("① 前日", "15分", "設定を家で全部固定する"),
         ("② 到着", "2分", "定位置に立つ"),
         ("③ 儀式", "90秒", "手拍子・騒音計・ベル4方位"),
         ("④ 本番", "20〜30秒", "無言。1通過＝1テイク"),
         ("⑤ 直後", "30秒", "5秒数えてから声で記録"),
         ("⑥ 確認", "5分", "その場で再生して品質確認")]
SY2, SW = 146, 136.0
rule(sl, 70, SY2, 820, 1.4, INK)
for i, (a, t, b) in enumerate(STEPS):
    x = 70 + i * SW
    if i:
        vrule(sl, x - 6, SY2 + 12, 76)
    txt(sl, x + 6, SY2 + 12, SW - 16, 20,
        [(a, {"size": 12.5, "bold": True, "color": PUR})])
    txt(sl, x + 6, SY2 + 32, SW - 16, 16, [(t, {"size": 9.5, "color": MUTED})])
    txt(sl, x + 6, SY2 + 52, SW - 16, 40,
        [(b, {"size": 10, "color": SUB})], line=1.3)
rule(sl, 70, SY2 + 98, 820)

y = sect(sl, 70, 290, 390, "③ 儀式の中身", "90秒")
for i, (a, b) in enumerate([
    ("手拍子を1回", "録音と動画の時刻を合わせる"),
    ("騒音計と並べて無言で60秒", "何デシベルだったかの基準を取る"),
    ("ベルを前後左右から1打ずつ", "マイクの正面と体の正面のズレを測る"),
]):
    txt(sl, 70, y, 20, 18, [("%d" % (i + 1), {"size": 11, "bold": True, "color": PUR})])
    txt(sl, 92, y, 368, 18, [(a, {"size": 11})])
    txt(sl, 92, y + 18, 368, 16, [(b, {"size": 9, "color": MUTED})])
    y += 40

vrule(sl, 492, 290, 130)
y = sect(sl, 520, 290, 370, "帰宅後、その日のうちに")
for t in ["録音を学習データと同じ形式に変換する", "1件ずつの10秒クリップに切り出す",
          "紙の記録を表に打ち込む", "検査スクリプトで形式ミスを機械的に見つける"]:
    rect(sl, 520, y + 7, 7, 1.4, PUR)
    txt(sl, 536, y, 354, 18, [(t, {"size": 11, "color": SUB})])
    y += 26

rule(sl, 70, 432, 820)
txt(sl, 70, 442, 820, 18,
    [[("この流れは9/1に手元で通し確認まで済ませてある。", {"size": 11, "bold": True}),
      ("　機材が届いたその日から本番に入れる。", {"size": 11, "color": SUB})]])

# ==================================================== p14 機材購入の相談
sl = new_slide("機 材 購 入 の 相 談", 14, "相談")
lead(sl, "実環境で210回録るために、満たさなければならない条件が5つある。")

y = sect(sl, 70, 148, 360, "購入をお願いしたい品目")
for a, b in [("Zoom H3-VR（4方向マイク）", "本体。右の条件を満たす唯一の機種"),
             ("騒音計（平均音量が出るもの）", "学内でお借りできれば購入は不要"),
             ("小物（ハーネス・変換金具ほか）", "自費でもまかなえる範囲")]:
    y = item(sl, 70, y, 360, a, b, hsize=12, nsize=9.5, gap=14)
txt(sl, 70, 380, 360, 16,
    [("※ 見積りは別紙でご用意しています", {"size": 9.5, "color": MUTED})])

vrule(sl, 466, 148, 268)

y = sect(sl, 506, 148, 384, "この機種を選んだ理由",
         "ZOOMの現行10機種を確認")
for a, b in [
    ("上下を含む4方向で録れる", "ステレオでは高さの軸が無く、学習データと形式が合わない"),
    ("96kHz で録れる", "バックする車が出す40kHzの超音波を、同じ録音で拾うため"),
    ("音量の自動調整を切れる", "録音中に変わると、何デシベルだったかを計算できない"),
    ("身に着けて持ち歩ける重さ", "6か所を歩いて回る。公道では三脚を立てられない"),
    ("マイクが較正済み", "自作だと、誤差がモデルのせいかマイクのせいか分からない"),
]:
    y = item(sl, 506, y, 384, a, b, hsize=11.5, nsize=9, gap=6)

rule(sl, 506, y + 4, 384)
txt(sl, 506, y + 14, 384, 18,
    [[("4方向で8機種が外れ、残る1機種は576gで重い。", {"size": 10.5, "color": SUB}),
      ("該当はH3-VRだけ。", {"size": 10.5, "bold": True})]])

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"修正_案A_8-14枚_2026-09-15.pptx")
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT, len(prs.slides._sldIdLst), "枚")
