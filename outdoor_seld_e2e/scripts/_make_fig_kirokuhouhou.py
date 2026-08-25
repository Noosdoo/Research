# -*- coding: utf-8 -*-
"""実録スライドの「記録方法」の図を、道路らしい見た目に作り直す。

背景: 2026-08-25、本人「ここの図の道路？もう少し道路っぽい図にしてほしい」。
旧版は灰色の帯に車の箱を置いただけで、道路に見えなかった。
車道（アスファルト）・センターライン・路側線・縁石・歩道を描き分ける。

文字は本人指示により**すべて18pt以上**（2026-08-25「基本18pt以上」）。
2026-08-25 追記: 本人が現行p13の「記録方法」枠に合わせたいとのことで、
その枠（左161pt・上325pt・幅673pt・高162pt）にぴったり収まる寸法で描き直した。
出力: md/seminar/図_記録方法_道路_2026-08-25.pptx（コピペ用・1枚）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
ASPHALT = RGBColor(0xC2, 0xC7, 0xCD)      # 車道
CURB = RGBColor(0x97, 0x9D, 0xA7)         # 縁石
WALK = RGBColor(0xEC, 0xEE, 0xF1)         # 歩道

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=18, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def shape(x, y, w, h, fill, line=None, lw=1.0, kind=MSO_SHAPE.RECTANGLE, rot=None):
    sh = sl.shapes.add_shape(kind, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if rot is not None:
        sh.rotation = rot
    return sh


def conn(x1, y1, x2, y2, color, lw=1.0, dash=None):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Pt(x1), Pt(y1), Pt(x2), Pt(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    if dash:
        c.line.dash_style = dash
    return c


# ============ 枠（現行p13の「記録方法」ボックスと同寸） ============
FX, FY, FW, FH = 161.0, 325.0, 673.0, 162.0
shape(FX, FY, FW, FH, WHITE, line=LINE, lw=1.0)

txt(FX + 16, FY + 8, 200, 26, [("記録方法", {"size": 18, "bold": True})])
txt(FX + FW - 300, FY + 8, 284, 26,
    [("速度：徐行〜約30km/h", {"size": 18, "color": SUB})], align=PP_ALIGN.RIGHT)

# ---- 道路 ----
RX0, RX1 = FX + 14, FX + FW - 14
ROAD_T, ROAD_B = FY + 38, FY + 92
CURB_B = ROAD_B + 6
WALK_B = CURB_B + 42

shape(RX0, ROAD_T, RX1 - RX0, ROAD_B - ROAD_T, ASPHALT)
shape(RX0, ROAD_B, RX1 - RX0, CURB_B - ROAD_B, CURB)
shape(RX0, CURB_B, RX1 - RX0, WALK_B - CURB_B, WALK)

cy = (ROAD_T + ROAD_B) / 2                 # センターライン（白破線）
x = RX0 + 12
while x < RX1 - 12:
    shape(x, cy - 2, 26, 4, WHITE)
    x += 46
shape(RX0, ROAD_T + 6, RX1 - RX0, 3, WHITE)    # 路側線
shape(RX0, ROAD_B - 9, RX1 - RX0, 3, WHITE)

txt(RX0 + 8, CURB_B + 8, 90, 26, [("歩道", {"size": 18, "color": SUB})])

# ---- 車 ----
CX, CY = FX + 150, ROAD_T + 10
shape(CX - 38, CY, 76, 34, GOLD, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(CX - 38, CY + 3, 76, 28, [("車", {"size": 18, "bold": True, "color": WHITE})],
    align=PP_ALIGN.CENTER)
shape(CX + 46, CY + 8, 32, 18, NAVY, kind=MSO_SHAPE.RIGHT_ARROW)

# ---- 装着者 ----
PY = CURB_B + 26
shape(CX - 8, PY - 8, 16, 16, NAVY, kind=MSO_SHAPE.OVAL)
shape(CX - 15, PY - 15, 30, 30, None, line=GOLD, lw=2.0, kind=MSO_SHAPE.OVAL)

# ---- 横距離 ----
conn(CX, CY + 38, CX, PY - 18, RED, lw=1.8, dash=MSO_LINE_DASH_STYLE.DASH)
txt(CX + 26, CY + 40, FX + FW - (CX + 34), 56,
    [[("横距離 2.0〜3.2m", {"size": 18, "bold": True, "color": RED}),
      (" ＝ 鳴るべき車", {"size": 18, "color": SUB})],
     [("　　　　 5〜15m", {"size": 18, "bold": True, "color": SUB}),
      (" ＝ 鳴ってはいけない車", {"size": 18, "color": SUB})]], line=1.25)

txt(CX + 26, PY - 4, 300, 26,
    [("装着者（胸部に4chマイク）", {"size": 18})])

# ---- フッター ----
txt(58, 506, 110, 24, [("2026/08/30", {"size": 18, "color": MUTED})])
txt(W - 90, 506, 40, 24, [("13", {"size": 18, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"図_記録方法_道路_2026-08-25.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
