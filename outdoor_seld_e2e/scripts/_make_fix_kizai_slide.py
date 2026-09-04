# -*- coding: utf-8 -*-
"""【修正スライド1枚だけ】機材購入の相談 — 品目表の右に「選定の理由」を控えめに添える版。

本人の依頼（2026-09-03）:
 初版（要件テーブル＋絞り込みの図）は主張が強すぎたため差し戻し。
 「もっと控えめで、これ（既存の品目表）と併せて」＝ 品目表は左に残し、
 空いている右半分に理由を静かに置く。

意匠の方針: 濃い見出し帯・色付きの図・大きな数字を使わない。
 罫線1本と灰色の本文で構成し、強調は結論の1行だけ。

要件の正本 = md/design/実録機材リサーチ_2026-08-12.md §1（R1〜R3）
 ＋ 2026-09-03に公式仕様で確認した項目（マッチング済マイク×4・166g・H8は576g）。

⚠️ 金額は載せない（2026-09-03本人「金額は載せたくない」）。品目だけを挙げ、
 空いた分は「なぜその品目が要るか」の一言で埋める。見積りは別紙で持参する想定。

出力: md/seminar/修正_機材購入の相談_2026-09-15.pptx（1枚のみ）
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
RED = RGBColor(0xC0, 0x39, 0x2B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
PALE = RGBColor(0xF5, 0xF3, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


# ---- 骨格（本人のデッキのまま） ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 170, 40,
    [("機 材 購 入 の 相 談", {"size": 22, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

# ---- 外枠（既存と同じ・全幅） ----
BX, BY, BW, BH = 70, 132, 820, 300
rect(BX, BY, BW, BH, WHITE, line=LINE, lw=1.0)
rect(BX, BY, BW, 3.5, PUR)

# ================= 左：購入をお願いしたい品目（金額なし） =================
LX = BX + 24
txt(LX, BY + 22, 300, 22, [("購入をお願いしたい品目", {"size": 14, "bold": True})])
rect(LX, BY + 50, 300, 1.0, HAIR)

items = [
    ("Zoom H3-VR（4方向マイク）", "本体。右の条件を満たす唯一の機種"),
    ("騒音計（平均音量が出るもの）", "学内でお借りできれば購入は不要"),
    ("小物（ハーネス・変換金具ほか）", "自費でもまかなえる範囲"),
]
iy = BY + 70
for name, note in items:
    rect(LX + 2, iy + 6.5, 4, 4, PUR)
    txt(LX + 14, iy, 292, 34,
        [[(name, {"size": 11.5, "color": INK})],
         [(note, {"size": 9, "color": MUTED})]], line=1.3)
    iy += 46

txt(LX, BY + 246, 300, 18,
    [("※ 見積りは別紙でご用意しています", {"size": 9.5, "color": MUTED})])

# ---- 縦の細い仕切り ----
rect(BX + 352, BY + 26, 1.0, BH - 52, HAIR)

# ================= 右：この機種を選んだ理由（控えめ） =================
RX = BX + 386
RW = BW - 386 - 22

txt(RX, BY + 22, RW, 22, [("この機種を選んだ理由", {"size": 14, "bold": True})])
rect(RX, BY + 50, RW, 1.0, HAIR)
txt(RX, BY + 60, RW, 18,
    [("録音機に必要な条件を5つ決め、ZOOMの現行10機種を確認しました。",
      {"size": 10, "color": MUTED})])

conds = [
    ("上下を含む4方向で録れる",
     "ステレオでは高さの軸が無く、学習データと形式が合わない"),
    ("96kHzで録れる",
     "バックする車が出す40kHzの超音波を、同じ録音で拾うため"),
    ("音量の自動調整を切れる",
     "録音中に変わると、何デシベルだったかを計算できない"),
    ("身に着けて持ち歩ける重さ",
     "6か所を歩いて回る。公道では三脚を立てられない"),
    ("マイクが較正済み",
     "自作だと、誤差がモデルのせいかマイクのせいか分からない"),
]
cy = BY + 86
for a, b in conds:
    rect(RX + 2, cy + 6.5, 4, 4, PUR)
    txt(RX + 14, cy, RW - 14, 32,
        [[(a, {"size": 11, "color": INK})],
         [(b, {"size": 9, "color": MUTED})]], line=1.25)
    cy += 34

rect(RX, cy + 6, RW, 34, PALE)
txt(RX + 12, cy + 14, RW - 24, 20,
    [[("4方向で8機種が外れ、残る1機種は576gで重い。", {"size": 10.5, "color": SUB}),
      ("該当はH3-VRだけ。", {"size": 10.5, "bold": True, "color": INK})]])

# ---- フッター（本人のデッキと同形式・p14） ----
txt(58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "背景・目的")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    label(c, s, size=10.5, color=WHITE if active else MUTED, bold=active)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("14", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"修正_機材購入の相談_2026-09-15.pptx")
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
