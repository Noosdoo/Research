# -*- coding: utf-8 -*-
"""「検証の設計」スライドを本人テンプレ風デザインで1枚生成。

背景: 2026-08-21、夏ゼミでは実験結果・考察を載せない方針になり、結果ページを
削除した。その結果この枚が表1つだけで寂しくなったため、「どう検証するか」で
1枚を構成し直したもの。**数値は結果ではなくデータセットの規模のみ**。

旧版（表だけ）= _make_fig_dataset_slide.py / 図_データセット評価手順_2026-08-13.pptx
出力: md/seminar/図_検証の設計_2026-08-21.pptx（コピペ用・1枚）
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
GRAY_HDR = RGBColor(0xEF, 0xEF, 0xEF)
CREAM = RGBColor(0xFB, 0xF4, 0xDE)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def card(x, y, w, h, tag, title, fill=WHITE, accent=PUR):
    rect(x, y, w, h, fill, line=LINE, lw=1.0)
    rect(x, y, 4, h, accent)
    txt(x + 16, y + 10, w - 30, 18, [(tag, {"size": 11, "bold": True,
                                            "color": accent, "spc": 120})])
    txt(x + 16, y + 28, w - 30, 24, [(title, {"size": 15, "bold": True})])


# ---- テンプレ骨格 ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 160, 40, [("検証の設計 — 学習に使っていないデータで1回だけ測る",
                           {"size": 23, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

# ---- リード文 ----
txt(70, 100, W - 140, 24,
    [("データを3つに分け、", {"size": 13.5, "color": SUB}),
     ("評価は一度も見ていないデータで1回だけ採点する", {"size": 13.5, "bold": True}),
     ("。手順は測る前に凍結してある。", {"size": 13.5, "color": SUB})])

# ---- 3分割カード ----
CY, CH, CW = 132, 156, 250
xs = [70, 355, 640]

card(xs[0], CY, CW, CH, "STEP 1", "学習（微調整）")
txt(xs[0] + 16, CY + 58, CW - 30, 30,
    [("10,200", {"size": 26, "bold": True}), ("　クリップ", {"size": 12, "color": SUB})])
txt(xs[0] + 16, CY + 96, CW - 30, 54,
    [[("8クラス・約28時間の合成データ", {"size": 12, "color": SUB})],
     [("距離ヘッド込みで全体を微調整", {"size": 12, "color": SUB})]], line=1.35)

card(xs[1], CY, CW, CH, "STEP 2", "検証")
txt(xs[1] + 16, CY + 58, CW - 30, 30,
    [("別クリップ", {"size": 20, "bold": True}),
     ("　同一方式", {"size": 12, "color": SUB})])
txt(xs[1] + 16, CY + 96, CW - 30, 54,
    [[("モデルの選択と通知しきい値の決定", {"size": 12, "color": SUB})],
     [("ここで決めた設定は以後変えない", {"size": 12, "color": SUB})]], line=1.35)

card(xs[2], CY, CW, CH, "STEP 3", "評価", fill=CREAM, accent=GOLD)
txt(xs[2] + 16, CY + 58, CW - 30, 30,
    [("1,800", {"size": 26, "bold": True}), ("　クリップ", {"size": 12, "color": SUB})])
txt(xs[2] + 16, CY + 96, CW - 30, 54,
    [[("同一設計・新乱数のデータ", {"size": 12, "color": SUB})],
     [("学習にも検証にも未使用", {"size": 12, "color": SUB})],
     [("1回だけ採点する", {"size": 12, "bold": True, "color": GOLD})]], line=1.35)

for ax in (xs[0] + CW + 6, xs[1] + CW + 6):
    txt(ax, CY + 62, 24, 24, [("→", {"size": 20, "color": MUTED})],
        align=PP_ALIGN.CENTER)

# ---- 下段2枚 ----
RY, RH = 304, 122
rect(70, RY, 395, RH, WHITE, line=LINE, lw=1.0)
txt(86, RY + 12, 365, 20, [("何を測るか", {"size": 14, "bold": True})])
txt(86, RY + 38, 365, 78,
    [[("・種類が当たったか（検出率）", {"size": 12, "color": SUB})],
     [("・方向が何度ずれたか（方向誤差）", {"size": 12, "color": SUB})],
     [("・距離が何メートルずれたか（絶対誤差）", {"size": 12, "color": SUB})],
     [("・通知が届いたか、余計に鳴らなかったか", {"size": 12, "color": SUB})]], line=1.4)

rect(495, RY, 395, RH, WHITE, line=LINE, lw=1.0)
txt(511, RY + 12, 365, 20, [("なぜ1回だけなのか", {"size": 14, "bold": True})])
txt(511, RY + 38, 365, 78,
    [[("何度も測って一番良い回を選ぶと、", {"size": 12, "color": SUB})],
     [("偶然を実力と取り違えてしまう。", {"size": 12, "color": SUB})],
     [("そこで基準としきい値を先に凍結し、", {"size": 12, "color": SUB})],
     [("生成したてのデータで1回だけ採点する。", {"size": 12, "color": SUB})]], line=1.4)

# ---- 金色バーの強調枠 ----
rect(70, 440, 6, 52, GOLD)
rect(76, 440, W - 146, 52, WHITE, line=LINE, lw=1.0)
txt(94, 448, W - 180, 40, [
    ("基準を先に決める → 未使用データを新たに生成 → 1回だけ評価", {"bold": True, "size": 14}),
    ("　＝ 答案を見てから基準を変えられない手順", {"size": 13, "color": SUB}),
], anchor=MSO_ANCHOR.MIDDLE)

# ---- 下部フッター ----
txt(58, 506, 90, 20, [("2026/09", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "検証")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(10.5)
    r.font.bold = active
    r.font.color.rgb = WHITE if active else MUTED
    r.font.name = "Meiryo"
    meiryo(r)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("9", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/図_検証の設計_2026-08-21.pptx"
prs.save(OUT)
print("saved:", OUT)
