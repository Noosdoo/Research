# -*- coding: utf-8 -*-
"""付録用「実録の流れ」スライド3枚（6コマを2コマずつ・図入り）を生成。

Artifactのコマ割り（交差点編）をpptxネイティブ図形で再構成したもの。
配色・書体は本体デッキと同一（紺INK×アンバー×メイリオ）。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x16, 0x23, 0x3A)
INK2 = RGBColor(0x3A, 0x46, 0x58)
MUTED = RGBColor(0x66, 0x70, 0x7F)
PAPER2 = RGBColor(0xEF, 0xEF, 0xEA)
HAIR = RGBColor(0xD9, 0xDA, 0xD2)
AMBER = RGBColor(0xE8, 0xA2, 0x00)
AMBERD = RGBColor(0xB3, 0x7E, 0x00)
RED = RGBColor(0xC4, 0x43, 0x2B)
REDL = RGBColor(0xE0, 0x96, 0x85)
GREEN = RGBColor(0x2F, 0x7D, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY = RGBColor(0xE7, 0xEB, 0xF2)
ROAD = RGBColor(0xDD, 0xDF, 0xD8)
WALL = RGBColor(0xC2, 0xC6, 0xBA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def meiryo(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", "メイリオ")


def txt(sl, x, y, w, h, lines, size=12, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, mono=False, gap=3):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        if isinstance(ln, str):
            ln = (ln, {})
        t, o = ln
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.color.rgb = o.get("color", color)
        f.name = "Consolas" if mono else "Meiryo"
        if not mono:
            meiryo(r)
    return tb


def box(sl, x, y, w, h, fill=WHITE, line=HAIR, line_w=1.0, radius=0.08, shape=None):
    sh = sl.shapes.add_shape(shape or MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if shape is None:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def rect(sl, x, y, w, h, fill, line=None, line_w=1.0):
    return box(sl, x, y, w, h, fill=fill, line=line, line_w=line_w,
               shape=MSO_SHAPE.RECTANGLE)


def circ(sl, cx, cy, r, fill=None, line=None, line_w=1.5, dash=None):
    sh = box(sl, cx - r, cy - r, 2 * r, 2 * r, fill=fill, line=line,
             line_w=line_w, shape=MSO_SHAPE.OVAL)
    if dash and line is not None:
        ln = sh.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    return sh


def seg(sl, x1, y1, x2, y2, color=INK, w=2.0, dash=None):
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                 Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    if dash:
        ln = cn.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    cn.shadow.inherit = False
    return cn


def title(sl, text):
    txt(sl, 0.55, 0.3, W - 1.1, 0.6, [(text, {"size": 20, "bold": True})])
    bar = rect(sl, 0.55, 0.92, 1.6, 0.045, AMBER)


def panel(sl, x, no, name, time_s, scene_fill):
    """コマ枠。戻り値=シーン描画領域 (sx, sy, sw, sh)。"""
    box(sl, x, 1.12, 5.95, 5.95, fill=WHITE, line=HAIR, radius=0.05)
    txt(sl, x + 0.22, 1.26, 0.75, 0.55, [(no, {"size": 24, "bold": True, "color": AMBERD})])
    txt(sl, x + 0.78, 1.36, 3.9, 0.45, [(name, {"size": 14.5, "bold": True})])
    ch = box(sl, x + 4.62, 1.32, 1.15, 0.38, fill=INK, line=None, radius=0.5)
    txt(sl, x + 4.62, 1.36, 1.15, 0.32, [(time_s, {"size": 10.5, "bold": True, "color": WHITE})],
        align=PP_ALIGN.CENTER)
    sx, sy, sw, sh = x + 0.22, 1.86, 5.51, 3.55
    rect(sl, sx, sy, sw, sh, scene_fill)
    return sx, sy, sw, sh


def cap(sl, x, lines):
    txt(sl, 0.22 + x, 5.6, 5.51, 1.35, lines, size=11.5, color=INK2, gap=4)


def person_td(sl, cx, cy):
    """トップダウンの自分アイコン（点＋リング）。"""
    circ(sl, cx, cy, 0.2, fill=None, line=AMBER, line_w=2.25)
    circ(sl, cx, cy, 0.1, fill=INK, line=None)


def new_slide():
    return prs.slides.add_slide(BLANK)


# ================= スライド1: コマ1・2 =================
sl = new_slide()
title(sl, "付録：実録の流れ 1/3 — 準備は家で終わらせる（例：見通し不良の交差点）")

# --- コマ1 家で準備 ---
sx, sy, sw, sh = panel(sl, 0.55, "1", "前日・家で全部設定", "前日15分", PAPER2)
rect(sl, sx + 0.25, sy + 2.85, sw - 0.5, 0.07, HAIR)                # 机
rect(sl, sx + 0.45, sy + 1.7, 1.5, 1.05, SKY, line=HAIR)            # ノートPC画面
rect(sl, sx + 0.3, sy + 2.75, 1.8, 0.12, WALL)                      # キーボード
b = box(sl, sx + 2.35, sy + 1.95, 0.62, 0.88, fill=WHITE, line=INK, line_w=2, radius=0.12)
for i in range(2):
    for j in range(2):
        circ(sl, sx + 2.53 + j * 0.26, sy + 2.16 + i * 0.2, 0.035, fill=INK)
rect(sl, sx + 2.44, sy + 2.6, 0.44, 0.16, AMBER)
box(sl, sx + 3.25, sy + 0.35, 2.05, 2.5, fill=WHITE, line=AMBER, line_w=1.75)
txt(sl, sx + 3.4, sy + 0.5, 1.8, 2.2, [
    ("設定チェック", {"size": 11.5, "bold": True}),
    ("☑ AmbiX 4ch", {"size": 10.5, "color": INK2}),
    ("☑ 96kHz / 24bit", {"size": 10.5, "color": INK2}),
    ("☑ ゲイン固定", {"size": 10.5, "color": INK2}),
    ("☑ Lo Cut OFF", {"size": 10.5, "color": INK2}),
    ("☑ 地図にピン5〜8個", {"size": 10.5, "color": INK2}),
], gap=2)
txt(sl, sx + 0.3, sy + 0.35, 2.6, 1.1,
    [("レコーダー設定は家で固定。", {"size": 11, "color": MUTED}),
     ("現地で押すのはRECだけ", {"size": 11, "color": MUTED})], gap=1)
cap(sl, 0.55, "ロケハンは散歩1回（地図アプリにピンを立てるだけ）。持ち物はリュック1個：H3-VR＋風防・ハーネス・騒音計・ベル・スマホ・電池/SD・イヤホンの8点")

# --- コマ2 到着・立ち位置（トップダウン） ---
sx, sy, sw, sh = panel(sl, 6.85, "2", "到着 — 角の手前に立つだけ", "2分", SKY)
rect(sl, sx, sy + 1.3, sw, 0.85, ROAD)                               # 横道
rect(sl, sx + 2.35, sy, 0.85, sh, ROAD)                              # 縦道
rect(sl, sx, sy + 2.15, 2.35, sh - 2.15, WALL)                       # 塀ブロック左下
rect(sl, sx, sy + 2.15, sw, 0.1, HAIR)                               # 歩道ライン
txt(sl, sx + 0.25, sy + 2.65, 2.0, 0.4, [("塀・建物", {"size": 10.5, "color": INK2})])
person_td(sl, sx + 4.05, sy + 2.62)
txt(sl, sx + 4.35, sy + 2.42, 1.15, 0.7, [("自分", {"size": 10.5, "bold": True}),
                                          ("（歩道）", {"size": 9.5, "color": MUTED})], gap=0)
seg(sl, sx + 3.2, sy + 2.5, sx + 3.85, sy + 2.5, color=RED, w=1.75)
seg(sl, sx + 3.2, sy + 2.4, sx + 3.2, sy + 2.6, color=RED, w=1.75)
txt(sl, sx + 2.9, sy + 2.95, 1.6, 0.35, [("角から5〜10m", {"size": 9.5, "color": RED})])
cap(sl, 6.85, "ブロック塀のあるT字路・十字路の歩道が定位置。公道の歩道からの収録なので許可は不要（通行の邪魔にならない位置で）。協力者がいれば10m後ろ、いなくても成立")

# ================= スライド2: コマ3・4 =================
sl = new_slide()
title(sl, "付録：実録の流れ 2/3 — 現地は「儀式90秒 → 無言で待つ」だけ")

# --- コマ3 儀式 ---
sx, sy, sw, sh = panel(sl, 0.55, "3", "録音開始 — 90秒の儀式", "90秒", PAPER2)
cards = [(0.15, "① 手拍子1回", "時刻合わせの目印"),
         (1.95, "② 騒音計を読む", "音量のものさし・小声で1分"),
         (3.75, "③ ベル4方位", "マイクの向き合わせ・各1打")]
for dx, name, note in cards:
    box(sl, sx + dx, sy + 0.2, 1.62, 3.15, fill=WHITE, line=HAIR)
    txt(sl, sx + dx + 0.08, sy + 0.32, 1.46, 0.75, [(name, {"size": 11, "bold": True})],
        align=PP_ALIGN.CENTER)
    txt(sl, sx + dx + 0.08, sy + 2.55, 1.46, 0.75, [(note, {"size": 9, "color": MUTED})],
        align=PP_ALIGN.CENTER, gap=1)
# ①手拍子グリフ
gx, gy = sx + 0.15 + 0.81, sy + 1.7
seg(sl, gx - 0.22, gy - 0.22, gx + 0.22, gy + 0.22, color=INK, w=4)
seg(sl, gx + 0.22, gy - 0.22, gx - 0.22, gy + 0.22, color=INK, w=4)
seg(sl, gx + 0.3, gy - 0.34, gx + 0.44, gy - 0.46, color=AMBER, w=2.5)
seg(sl, gx + 0.36, gy - 0.1, gx + 0.54, gy - 0.14, color=AMBER, w=2.5)
# ②騒音計グリフ
mx = sx + 1.95 + 0.81
box(sl, mx - 0.3, sy + 1.15, 0.6, 1.15, fill=WHITE, line=INK, line_w=2, radius=0.15)
rect(sl, mx - 0.2, sy + 1.3, 0.4, 0.32, SKY)
txt(sl, mx - 0.35, sy + 1.31, 0.7, 0.3, [("52.3", {"size": 9, "bold": True})],
    align=PP_ALIGN.CENTER)
# ③ベルグリフ
bx, by = sx + 3.75 + 0.81, sy + 1.75
circ(sl, bx, by, 0.13, fill=INK)
for ddx, ddy, lab, lx, ly in [(0, -0.55, "前", -0.07, -0.95), (0.55, 0, "右", 0.62, -0.12),
                              (0, 0.55, "後", -0.07, 0.62), (-0.55, 0, "左", -0.78, -0.12)]:
    seg(sl, bx + ddx * 0.45, by + ddy * 0.45, bx + ddx, by + ddy, color=AMBER, w=3)
    txt(sl, bx + lx, by + ly, 0.3, 0.3, [(lab, {"size": 9.5, "color": MUTED})])
cap(sl, 0.55, "声を出すのはこの90秒だけ。マイクは胸元にあるのでささやき声で足りる。1人のときはベルを置いて自分が90°ずつ回っても同じ効果")

# --- コマ4 本番（トップダウン・見えないのに聞こえる） ---
sx, sy, sw, sh = panel(sl, 6.85, "4", "本番 — 無言で待つだけ", "10秒×台数", SKY)
rect(sl, sx, sy + 1.3, sw, 0.85, ROAD)
rect(sl, sx + 2.35, sy, 0.85, sh, ROAD)
# 音の波紋（車から広がる）→塀より先に描く
carx, cary = sx + 2.77, sy + 0.55
for r, col in [(0.45, RED), (0.8, RED), (1.2, REDL)]:
    circ(sl, carx, cary, r, fill=None, line=col, line_w=1.75)
box(sl, carx - 0.26, cary - 0.38, 0.52, 0.76, fill=RED, line=None, radius=0.3)
# 塀（右上=視線ブロック）
rect(sl, sx + 3.2, sy, sw - 3.2, 1.3, WALL)
rect(sl, sx, sy, 2.35, 1.3, WALL)
rect(sl, sx, sy + 2.15, sw, 0.1, HAIR)
txt(sl, sx + 3.5, sy + 0.35, 1.9, 0.75, [("塀", {"size": 10.5, "color": INK2}),
                                         ("（車はまだ見えない）", {"size": 9, "color": INK2})], gap=0)
person_td(sl, sx + 4.05, sy + 2.62)
# 視線（塀でブロック）
seg(sl, sx + 4.05, sy + 2.5, carx + 0.15, cary + 0.3, color=MUTED, w=1.5, dash="dash")
txt(sl, sx + 3.95, sy + 1.52, 1.5, 0.4, [("視線は塀でブロック", {"size": 9, "color": MUTED})])
txt(sl, sx + 0.2, sy + 1.42, 2.1, 0.6, [("音だけが先に", {"size": 10.5, "bold": True, "color": RED}),
                                        ("角を曲がって届く", {"size": 10.5, "bold": True, "color": RED})], gap=0)
txt(sl, sx + 4.35, sy + 2.75, 1.15, 0.4, [("無言で立つ", {"size": 9.5, "color": MUTED})])
cap(sl, 6.85, "録音は回しっぱなし。角の向こうの車は姿より先に音が届く——この10秒がそのまま1テイク。来た台数のぶんだけテイクが増える")

# ================= スライド3: コマ5・6 =================
sl = new_slide()
title(sl, "付録：実録の流れ 3/3 — 記録は1行、チェックは30秒")

# --- コマ5 スレート ---
sx, sy, sw, sh = panel(sl, 0.55, "5", "通過直後 — ひとことメモ", "30秒", PAPER2)
box(sl, sx + 0.25, sy + 0.3, 2.2, 2.95, fill=WHITE, line=INK, line_w=2, radius=0.1)
rect(sl, sx + 1.05, sy + 0.42, 0.6, 0.07, HAIR)
txt(sl, sx + 0.42, sy + 0.75, 1.9, 1.6, [
    ("「試行3終了。", {"size": 11, "bold": True}),
    ("車1台、右から、", {"size": 11, "bold": True}),
    ("徐行、横2mくらい」", {"size": 11, "bold": True}),
    ("と小声で吹き込む", {"size": 9.5, "color": MUTED}),
], gap=2)
tri = box(sl, sx + 2.6, sy + 1.55, 0.42, 0.42, fill=AMBER, line=None,
          shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
tri.rotation = 90
box(sl, sx + 3.15, sy + 0.85, 2.15, 1.85, fill=PAPER2, line=HAIR)
txt(sl, sx + 3.3, sy + 1.0, 1.9, 0.4, [("帰宅後、注釈CSVの1行に", {"size": 10, "bold": True})])
txt(sl, sx + 3.3, sy + 1.4, 1.9, 1.2, [
    ("trial=3, class=車,", {"size": 9}),
    ("象限=右, 徐行,", {"size": 9}),
    ("横2m, LAeq=52.3", {"size": 9}),
], mono=True, gap=1, color=INK2)
cap(sl, 0.55, "しゃべってよいのは試行と試行の間だけ（試行中は無言）。この一言が正解ラベルになる——フレーム単位のx,y,z・距離の手入力はしない設計")

# --- コマ6 現場QC ---
sx, sy, sw, sh = panel(sl, 6.85, "6", "撤収前チェック → 次の地点へ", "5分", PAPER2)
box(sl, sx + 0.25, sy + 0.3, 2.7, 2.95, fill=WHITE, line=HAIR)
hs = [0.14, 0.3, 0.2, 0.42, 0.26, 0.46, 0.3, 0.16, 0.36]
rows = [("W", INK), ("X", INK2), ("Y", AMBERD), ("Z", RED)]
for ri, (lab, col) in enumerate(rows):
    ry = sy + 0.62 + ri * 0.68
    txt(sl, sx + 0.36, ry - 0.1, 0.3, 0.3, [(lab, {"size": 10, "bold": True, "color": MUTED})])
    for bi, hv in enumerate(hs):
        rect(sl, sx + 0.72 + bi * 0.24, ry + 0.4 - hv * 0.75, 0.13, hv * 0.75, col)
txt(sl, sx + 3.15, sy + 0.5, 2.2, 2.6, [
    ("イヤホンで30秒だけ再生", {"size": 11.5, "bold": True}),
    ("✓ 4chすべて入っている", {"size": 11, "color": GREEN}),
    ("✓ 音割れしていない", {"size": 11, "color": GREEN}),
    ("✓ 手拍子・ベルが録れた", {"size": 11, "color": GREEN}),
    ("OK → 次のピンへ移動", {"size": 11, "bold": True, "color": GREEN}),
    ("NG → その場で撮り直し", {"size": 10.5, "color": MUTED}),
], gap=5)
cap(sl, 6.85, "「録れてなかった」で機材を借り直すのが一番高くつくので、その場で確認。1箇所15分×5箇所=半日で「角5本」完了。他のテイクもコマ2の場所が変わるだけ")

OUT = r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/図_付録_実録の流れ_2026-08-13.pptx"
prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
