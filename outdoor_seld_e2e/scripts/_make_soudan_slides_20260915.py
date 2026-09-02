# -*- coding: utf-8 -*-
"""9/15 相談会（夏ゼミ反省・相談会）用スライド 5枚。

目的は研究発表ではなく **相談**。議題は4件:
  ① 実録機材の購入（レンタルでなく購入）  ② 規程照会（このまま進めてよいか）
  ③ 騒音計の学内借用先                      ④ 卒論の提出期限・発表日程

材料の正本:
  md/design/ゼミ相談メモ_機材購入_2026-09-15用.md   （議題・論拠4つ・金額・想定問答）
  md/design/実録_収録前ゲート消化_2026-08-30.md §0/§1（逆算表・規程照会の文面）
  md/design/実録機材リサーチ_2026-08-12.md          （価格の出典・H3-VR一択の根拠）

表記の約束:
 - 初見の聴き手（他研究室の先生含む）向け。FOA/LAeq/ablation 等の内輪語は言い換える
 - 規模は「計210回（＋連続録音100分）」、期限は「12月末で全実験終了」で統一
 - 通知層の数値はこの資料では扱わない（相談会の議題ではないため）

出力: md/seminar/相談会スライド_機材購入_2026-09-15.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
PALE = RGBColor(0xF5, 0xF3, 0xF9)   # 薄紫（主張の枠）
PALEG = RGBColor(0xFD, 0xF7, 0xE8)  # 薄金（伺いたい点）
BAND = RGBColor(0xF2, 0xF3, 0xF5)   # 薄灰（補足帯）

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0

AGENDA = ["① 機材購入", "② 規程照会", "③ 騒音計", "④ 提出期限"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(sl, x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(sl, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER):
    """図形の中に1行だけ入れる。"""
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


def chip(sl, x, y, text, color, size=11.5, pad=16, h=22):
    """枠線だけの見出しチップ。幅は文字数から決める。"""
    c = rect(sl, x, y, pad + len(text) * 12.0, h, None, line=color, lw=1.2)
    label(c, text, size=size, color=color)
    return c


def bullet(sl, x, y, w, runs, size=11, mark=INK, line=1.34, h=40):
    rect(sl, x, y + 5.5, 6, 6, mark)
    return txt(sl, x + 16, y, w - 16, h, runs, size=size, line=line)


def new_slide(title, page, active=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 44, 14, 1.4, H - 28, PUR)
    rect(sl, 58, 46, 30, 6, PUR)
    txt(sl, 98, 30, W - 160, 40,
        [(title, {"size": 22, "bold": True, "spc": 250})])
    rect(sl, 44, 90, W - 90, 1.2, INK)

    txt(sl, 58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
    widths = [16 + len(s) * 12.0 for s in AGENDA]
    x = (W - (sum(widths) + 8 * (len(AGENDA) - 1))) / 2
    for s, wd in zip(AGENDA, widths):
        on = (s == active)
        c = rect(sl, x, 502, wd, 22, NAVY if on else CHIP)
        label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
        x += wd + 8
    txt(sl, W - 90, 506, 40, 20,
        [(str(page), {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)
    return sl


# ============================================================ P1 全体像
sl = new_slide("9/15 相談会でお願いしたいこと", 1)
txt(sl, 70, 100, 820, 20,
    [("テーマの変更はありません。今日は下の4件をまとめてご相談させてください。",
      {"size": 12, "color": MUTED})])

rect(sl, 70, 128, 820, 92, PALE)
rect(sl, 70, 128, 4, 92, PUR)
txt(sl, 92, 140, 780, 72,
    [[("卒業研究の実環境評価として、9月末〜11月に屋外での録音を", {"size": 14}),
      ("計210回", {"size": 14, "bold": True, "color": RED}),
      ("行います。", {"size": 14})],
     [("そこで使う4方向マイク（Zoom H3-VR・約2.6万円）を、", {"size": 14}),
      ("レンタルではなく研究費での購入", {"size": 14, "bold": True, "color": RED}),
      ("としてお願いできないでしょうか。", {"size": 14})]], line=1.45)

CY, CH, CW = 234, 170, 194
cards = [
    ("①", PUR, "機材の購入", "お願いしたいこと",
     ["4方向マイク Zoom H3-VR ほか", "合計 約3〜4万円",
      "→ 承認をいただきたい件"]),
    ("②", GOLD, "規程の照会", "確認したいこと",
     ["公道での録音に、倫理審査や", "届け出は必要でしょうか",
      "→ この場で伺えると助かります"]),
    ("③", GREEN, "騒音計の借用", "心当たりを伺いたい",
     ["音量の基準合わせに要る計測器", "学内で借りられる部署は",
      "→ 無ければ1万円以内で購入"]),
    ("④", NAVY, "卒論の期限", "教えていただきたい",
     ["正式な提出期限と発表日程", "収録・解析の逆算を確定したい",
      "→ 現在は12月末を自分の線に"]),
]
for i, (no, col, head, kind, lines) in enumerate(cards):
    x = 70 + i * (CW + 15)
    rect(sl, x, CY, CW, CH, WHITE, line=LINE, lw=1.0)
    rect(sl, x, CY, CW, 3.5, col)
    n = rect(sl, x + 14, CY + 18, 26, 26, col, shape=MSO_SHAPE.OVAL)
    label(n, no, size=13)
    txt(sl, x + 48, CY + 18, CW - 60, 34,
        [[(head, {"size": 14, "bold": True})],
         [(kind, {"size": 9.5, "color": MUTED})]], line=1.25)
    txt(sl, x + 16, CY + 66, CW - 30, 90,
        [[(t, {"size": 10.5, "color": SUB})] for t in lines], line=1.45)

rect(sl, 70, 418, 820, 40, BAND)
txt(sl, 86, 428, 790, 24,
    [[("承認が要るのは①だけです。", {"size": 12, "bold": True}),
      ("②③④は、この場で教えていただければそれで片づきます。",
       {"size": 12, "color": SUB})]])

# ============================================================ P2 議題① 購入
sl = new_slide("議題①　マイクは「レンタル」ではなく「購入」でお願いしたい", 2,
               active="① 機材購入")
txt(sl, 70, 100, 820, 20,
    [("録音は9月末〜11月、天候次第の日程で計210回。手元に1台ある方が、結果として"
      "安く・確実になります。", {"size": 12, "color": MUTED})])

reasons = [
    ("1", "天候で回数が読めない",
     "風速3m/s超と雨天は中止と決めています。台風と秋雨の9〜10月は、借りている期間が"
     "まるごと空振りになり得ます。晴れた日に録るには手元にある方が確実です。"),
    ("2", "総額が逆転しやすい",
     "レンタル3回で13,140円に対し新品は25,909円。天候で1〜2回借り直せばほぼ並び、"
     "それ以上は購入の方が安くなります。"),
    ("3", "測定の条件がそろう",
     "音量の基準合わせ（較正）は、同じ1台なら最初に1回やれば全収録に通ります。"
     "レンタルは毎回ちがう個体で、やり直すたびに測定条件が変わります。"),
    ("4", "修士の測定の基準器になる",
     "修士でデバイスを作ると自前のマイクが要りますが、その正しさは較正済みの機材と同じ場面を"
     "録って突き合わせる以外に確かめようがありません。2年半使う備品になります。"),
]
RY, RH = 132, 60
for i, (no, head, body) in enumerate(reasons):
    y = RY + i * (RH + 10)
    rect(sl, 70, y, 505, RH, WHITE, line=LINE, lw=1.0)
    n = rect(sl, 84, y + 8, 22, 22, PUR, shape=MSO_SHAPE.OVAL)
    label(n, no, size=12)
    txt(sl, 114, y + 7, 450, 18, [(head, {"size": 13, "bold": True})])
    txt(sl, 114, y + 27, 450, 30, [(body, {"size": 10, "color": SUB})], line=1.3)

TX, TY, TW, TH = 598, 132, 292, 270
rect(sl, TX, TY, TW, TH, WHITE, line=LINE, lw=1.0)
rect(sl, TX, TY, TW, 3.5, PUR)
txt(sl, TX + 16, TY + 14, TW - 32, 20,
    [("購入をお願いしたい品目", {"size": 13, "bold": True})])
rect(sl, TX + 16, TY + 38, TW - 32, 1.0, LINE)

rows = [
    ("Zoom H3-VR（4方向マイク）", "約26,000円",
     "要件を満たす現行機はこれだけ。中古なら約18,000円〜"),
    ("騒音計（平均音量が出るもの）", "0〜10,000円",
     "まず学内で借りられないか照会します"),
    ("小物（ハーネス・ブザー・SD等）", "約6,000円",
     "自費でもまかなえる範囲です"),
]
ry = TY + 46
for name, price, note in rows:
    txt(sl, TX + 16, ry, 178, 18, [(name, {"size": 10.5})])
    txt(sl, TX + 190, ry, 86, 18,
        [(price, {"size": 11, "bold": True})], align=PP_ALIGN.RIGHT)
    txt(sl, TX + 16, ry + 19, TW - 32, 26,
        [(note, {"size": 9, "color": MUTED})], line=1.25)
    ry += 52
rect(sl, TX + 16, ry - 6, TW - 32, 1.0, LINE)
txt(sl, TX + 16, ry + 4, 130, 20, [("合計", {"size": 12, "bold": True})])
txt(sl, TX + 146, ry + 3, 130, 22,
    [("約3〜4万円", {"size": 15, "bold": True, "color": RED})],
    align=PP_ALIGN.RIGHT)
txt(sl, TX + 16, ry + 34, TW - 32, 30,
    [("騒音計は学内で借りられれば0円、小物は自費でもまかなえるので、"
      "下がる余地があります。", {"size": 9, "color": MUTED})], line=1.3)

rect(sl, 70, 414, 820, 44, PALEG)
rect(sl, 70, 414, 4, 44, GOLD)
txt(sl, 88, 425, 790, 24,
    [[("比較の要点：", {"size": 12, "bold": True}),
      ("レンタル3回 13,140円 ＜ 新品 25,909円。ただし天候で1〜2回借り直すとほぼ並びます。"
       "実質の差額は1万円前後です。", {"size": 12, "color": SUB})]])

# ============================================================ P3 議題② 規程
sl = new_slide("議題②　この収録、規程のうえでそのまま進めてよいでしょうか", 3,
               active="② 規程照会")
txt(sl, 70, 100, 820, 20,
    [("公道での録音なので、通行人の声が偶然入る可能性があります。対策は先に決めてあります。",
      {"size": 12, "color": MUTED})])

rect(sl, 70, 128, 400, 168, WHITE, line=LINE, lw=1.0)
chip(sl, 86, 142, "収録の中身", PUR)
facts = [
    [("場所と内容：", {"bold": True}),
     ("公道と私有地（許可取得）での交通音の4方向録音", {"color": SUB})],
    [("規模：", {"bold": True}),
     ("計210回＋連続録音100分", {"color": SUB})],
    [("時期：", {"bold": True}), ("2026年9〜11月（天候により変動）", {"color": SUB})],
    [("収録者：", {"bold": True}),
     ("私1名（協力者の運転は私有地のみ・保険を事前確認）", {"color": SUB})],
]
fy = 176
for runs in facts:
    bullet(sl, 88, fy, 366, runs, size=10.5, mark=PUR, h=32)
    fy += 30

rect(sl, 490, 128, 400, 168, WHITE, line=LINE, lw=1.0)
chip(sl, 506, 142, "決めてある対策", GREEN)
guards = [
    "会話がはっきり聞き取れる区間は解析から除外する",
    "注釈用の動画は非公開。顔とナンバーはマスキングまたは削除",
    "原本は暗号化保存・卒論提出後30日以内に削除",
    "通学時間帯の学校周辺はルートから外す／説明カードを携行",
]
gy = 176
for t in guards:
    bullet(sl, 506, gy, 366, [(t, {"color": SUB})], size=10.5,
           mark=GREEN, h=32)
    gy += 30

rect(sl, 70, 312, 820, 116, PALEG)
rect(sl, 70, 312, 4, 116, GOLD)
txt(sl, 88, 324, 780, 20,
    [("特に伺いたいのは、この2点です", {"size": 14, "bold": True})])
asks = [
    ("1", "この収録と利用について、学内の倫理審査や届け出は必要でしょうか。"),
    ("2", "録音データの処理を Google Colab と学外の計算サーバーで行う予定です。"
          "人の声が偶然入り得るデータを学外で処理することは、情報管理の規程上"
          "問題ないでしょうか。"),
]
ay = 352
for no, t in asks:
    n = rect(sl, 90, ay + 1, 20, 20, GOLD, shape=MSO_SHAPE.OVAL)
    label(n, no, size=11)
    txt(sl, 118, ay, 752, 36, [(t, {"size": 11.5, "color": INK})], line=1.35)
    ay += 34

txt(sl, 70, 442, 820, 40,
    [[("この場で「進めてよい」と伺えれば、収録開始前の確認事項はこれで全部片づきます。",
       {"size": 11, "color": SUB})],
     [("委員会などへの確認が要ると分かった場合は、その結果を待ち、10月の収録日程を"
       "詰めて吸収します。", {"size": 11, "color": MUTED})]], line=1.4)

# ============================================================ P4 議題③④
sl = new_slide("議題③④　騒音計の借用先と、卒論の提出期限", 4, active="③ 騒音計")

rect(sl, 70, 106, 400, 158, WHITE, line=LINE, lw=1.0)
rect(sl, 70, 106, 400, 3.5, GREEN)
txt(sl, 88, 120, 366, 22, [("③ 騒音計を学内で借りられませんか",
                            {"size": 14, "bold": True})])
for i, runs in enumerate([
    [("何に使うか：", {"bold": True}),
     ("録音した音が実際に何デシベルだったかに換算するための基準。"
      "これが無いと音量の絶対値を言えません", {"color": SUB})],
    [("要件：", {"bold": True}),
     ("一定時間の平均音量（LAeq）が表示できるもの。できればJISクラス2相当",
      {"color": SUB})],
    [("お願い：", {"bold": True}),
     ("環境系・建築系・学生実験室など、心当たりはありますでしょうか",
      {"color": SUB})],
]):
    bullet(sl, 88, 150 + i * 36, 366, runs, size=10.5, mark=GREEN, h=34)
txt(sl, 88, 240, 366, 18,
    [("借りられなければ1万円以内で購入します", {"size": 9.5, "color": MUTED})])

rect(sl, 490, 106, 400, 158, WHITE, line=LINE, lw=1.0)
rect(sl, 490, 106, 400, 3.5, NAVY)
txt(sl, 508, 120, 366, 22, [("④ 卒論の提出期限と発表日程", {"size": 14, "bold": True})])
for i, runs in enumerate([
    [("伺いたいこと：", {"bold": True}),
     ("正式な提出期限と、発表の日程を教えてください", {"color": SUB})],
    [("いまの前提：", {"bold": True}),
     ("「12月末までに全実験を終える」を自分で決めた線として動いています",
      {"color": SUB})],
    [("理由：", {"bold": True}),
     ("1月からは執筆と資料作成だけにしたいので、正式な日程が分かれば"
      "収録と解析の逆算を確定できます", {"color": SUB})],
]):
    bullet(sl, 508, 150 + i * 36, 366, runs, size=10.5, mark=NAVY, h=34)

txt(sl, 70, 282, 500, 22,
    [("承認をいただけた場合の進み方", {"size": 14, "bold": True})])
steps = [
    ("9/15", "承認・発注", PUR),
    ("9月末", "着荷→動作確認・リハ", PUR),
    ("10〜11月中旬", "収録の本番", GOLD),
    ("11月中旬", "収録と注釈が完了", GOLD),
    ("12月末", "全実験を終了", RED),
    ("1月〜", "執筆に専念", NAVY),
]
CVW, STEP, CVY, CVH = 148, 134, 312, 54
for i, (when, what, col) in enumerate(steps):
    sh = rect(sl, 70 + i * STEP, CVY, CVW, CVH, col if i == 4 else WHITE,
              line=col, lw=1.2, shape=MSO_SHAPE.CHEVRON)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for pi, (t, sz, bd) in enumerate([(when, 10.5, True), (what, 9.5, False)]):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = WHITE if i == 4 else (col if bd else SUB)
        r.font.name = "Meiryo"
        meiryo(r)

rect(sl, 70, 390, 820, 44, BAND)
txt(sl, 86, 401, 790, 24,
    [[("承認後すぐ発注すれば数日で着荷し、9月末のリハに間に合います。",
       {"size": 11.5, "color": SUB}),
      ("録音データを処理する手順は9/1に通し確認まで済ませてあります。",
       {"size": 11.5, "color": INK, "bold": True})]])

# ============================================================ P5 想定問答
sl = new_slide("（参考）いただきそうな質問への答え", 5)
qa = [
    ("レンタルで足りるのではないですか",
     "天候で中止になると、借りている期間がまるごと無駄になります。回数が読めないうえ、"
     "毎回ちがう個体だと音量の基準合わせをやり直すことになり、測定条件がそろいません。"
     "実質の差額は1万円前後です。"),
    ("もっと安いマイクではだめですか",
     "上下方向まで含めて4方向を同時に録れて、学習に使ったデータと同じ形式で保存できる"
     "現行機は、調べた範囲でH3-VRだけでした（2026-08-12に出典つきで記録）。"
     "一段下の機種（H2n）は高さの軸が無く、学習済みモデルに入力できません。"),
    ("いつ買って、いつ録るのですか",
     "承認後すぐ発注（数日で着荷）→ 9月末に動作確認とリハーサル → 10〜11月中旬に本番 →"
     "12月末までに解析まで完了。1月からは執筆に専念します。"),
]
QY, QH = 112, 112
for i, (q, a) in enumerate(qa):
    y = QY + i * (QH + 16)
    rect(sl, 70, y, 820, QH, WHITE, line=LINE, lw=1.0)
    rect(sl, 70, y, 4, QH, GOLD)
    b = rect(sl, 90, y + 18, 24, 24, GOLD, shape=MSO_SHAPE.OVAL)
    label(b, "Q", size=12)
    txt(sl, 124, y + 18, 750, 24, [(q, {"size": 14, "bold": True})])
    b = rect(sl, 90, y + 54, 24, 24, None, line=PUR, lw=1.2, shape=MSO_SHAPE.OVAL)
    label(b, "A", size=12, color=PUR)
    txt(sl, 124, y + 52, 750, 48, [(a, {"size": 11, "color": SUB})], line=1.4)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"相談会スライド_機材購入_2026-09-15.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
