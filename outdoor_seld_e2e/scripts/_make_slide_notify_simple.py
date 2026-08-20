# -*- coding: utf-8 -*-
"""検証スライド「通知は間に合うようになったか」— 簡潔版（2026-08-20）。

## なぜ簡潔版を作るか

先に作った詳しい版（_make_slide_notify_results.py）は5行の表で、
「距離しきい値」「余裕が0.5秒以上あった割合」など**説明なしでは伝わらない語**が並び、
10分の発表では重すぎた。夏ゼミは先輩方のレジュメを見るかぎり提案段階の発表が主で、
結果を詳細に出す場ではない。

ただし 8/4 のゼミで「トラックが来る3秒前に検知できても使えない、そこを工夫するとか」
という指摘を受けている以上、**それに答えた事実は1枚だけ見せる**必要がある。

そこで **図1つ・数字1つ・代償1行**に絞る。詳しい版は付録に回して聞かれたら出す。

## 落とした行と理由

- 「余裕が0.5秒以上あった割合 0.6%→65.3%」…… 0.5秒という基準に外部の根拠がない
  （採点器がたまたま出す値を使っていた）
- 「安全な車への誤った至近警告 1.3%→3.7%」…… 代償は1行にまとめた
- 「実時間では75.0%」「キックボード88.7%」…… 付録・質疑へ

出力: md/seminar/図_検証_通知層_簡潔_2026-08-20.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

W, H = 960, 540
INK = RGBColor(0x16, 0x23, 0x3A)
INK2 = RGBColor(0x3A, 0x46, 0x58)
MUTED = RGBColor(0x66, 0x70, 0x7F)
HAIR = RGBColor(0xD9, 0xDA, 0xD2)
RED = RGBColor(0xC4, 0x43, 0x2B)
GREEN = RGBColor(0x2F, 0x7D, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHIP = RGBColor(0xF3, 0xF3, 0xEF)
JP = "游ゴシック"

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(int(W * 12700)), Emu(int(H * 12700))
sl = prs.slides.add_slide(prs.slide_layouts[6])


def _font(run, size, bold, color):
    f = run.font
    f.name, f.size, f.bold = JP, Pt(size), bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", JP)


def box(x, y, w, h):
    tb = sl.shapes.add_textbox(Emu(int(x * 12700)), Emu(int(y * 12700)),
                               Emu(int(w * 12700)), Emu(int(h * 12700)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb


def para(tf, runs, size=13, bold=False, color=INK2, align=PP_ALIGN.LEFT,
         first=False, line=1.3, space=2):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment, p.line_spacing, p.space_after = align, line, Pt(space)
    for text, kw in ([(runs, {})] if isinstance(runs, str) else runs):
        r = p.add_run()
        r.text = text
        _font(r, kw.get("size", size), kw.get("bold", bold), kw.get("color", color))
    return p


def rect(x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = sl.shapes.add_shape(shape, Emu(int(x * 12700)), Emu(int(y * 12700)),
                            Emu(int(w * 12700)), Emu(int(h * 12700)))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(12)
    tf.margin_top = tf.margin_bottom = Pt(9)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return s


def line(x1, y1, x2, y2, color, w=2.0, arrow=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1 * 12700)),
                                Emu(int(y1 * 12700)), Emu(int(x2 * 12700)),
                                Emu(int(y2 * 12700)))
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(w)
    if arrow:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return c


def label(x, y, w, text, size=11, color=MUTED, align=PP_ALIGN.LEFT, bold=False):
    tb = box(x, y, w, 20)
    tb.text_frame.word_wrap = False
    para(tb.text_frame, text, size=size, bold=bold, color=color, align=align,
         first=True, space=0)


# ============ 見出し ============
line(54, 96, W - 54, 96, INK, 1.5)
tb = box(54, 40, 860, 48)
para(tb.text_frame, "検 証 ― 通 知 は 「 間 に 合 う 」 よ う に な っ た か",
     size=26, bold=True, color=INK, first=True, space=0)

# ============ 図：いつ鳴るか ============
X0, X1 = 250, 800            # X1 = 最接近の位置
Y_OLD, Y_NEW = 186, 300
SEC = (X1 - X0) / 2.0        # 軸の全長を2.0秒とする
label(54, 126, 500, "車が一番近づく瞬間までに、どれだけ余裕があるか",
      size=13, color=INK2, bold=True)

# 最接近の縦線
line(X1, Y_OLD - 44, X1, Y_NEW + 52, MUTED, 1.2)
label(X1 - 60, Y_OLD - 68, 120, "車が一番近づく", size=12, color=INK,
      align=PP_ALIGN.CENTER, bold=True)

for y, tag, col, fire_x, note in (
        (Y_OLD, "前のやり方", RED, X1, "鳴った瞬間が、もう一番近い瞬間"),
        (Y_NEW, "新しいやり方", GREEN, X1 - 0.8 * SEC, "0.8秒前に鳴る")):
    line(X0, y, X1, y, col, 3.0)
    label(70, y - 11, 170, tag, size=14, color=col, bold=True)
    rect(fire_x - 11, y - 11, 22, 22, fill=col, shape=MSO_SHAPE.OVAL)
    if y == Y_OLD:
        label(fire_x - 240, y - 34, 210, "鳴る", size=13, color=col,
              align=PP_ALIGN.RIGHT, bold=True)
        label(fire_x - 330, y + 16, 300, note, size=12, color=col,
              align=PP_ALIGN.RIGHT)
    else:
        label(fire_x - 210, y - 34, 200, "鳴る", size=13, color=col,
              align=PP_ALIGN.RIGHT, bold=True)
        # 余裕の帯
        rect(fire_x, y + 14, X1 - fire_x, 22, fill=CHIP, line=GREEN, lw=1.0)
        label(fire_x, y + 17, X1 - fire_x, "0.8秒の余裕", size=12.5,
              color=GREEN, align=PP_ALIGN.CENTER, bold=True)

# ============ 数字 ============
c = rect(54, 372, 520, 76, fill=WHITE, line=HAIR)
para(c.text_frame, "1.5m以内まで近づく危険な車に、警告が届いた割合",
     size=12.5, color=INK2, first=True, space=6)
para(c.text_frame, [("69.9 %", {"size": 22, "bold": True, "color": MUTED}),
                    ("　→　", {"size": 16, "color": INK2}),
                    ("85.0 %", {"size": 26, "bold": True, "color": GREEN})],
     size=16, align=PP_ALIGN.CENTER)

c = rect(596, 372, 310, 76, fill=WHITE, line=RED, lw=1.5)
para(c.text_frame, "代わりに", size=12.5, bold=True, color=RED, first=True,
     space=4)
para(c.text_frame, "安全な車にも鳴りやすくなる（鳴らさずに済んだ割合 90.4→85.2%）",
     size=11.5, color=INK2, line=1.25)

label(54, 468, 860,
      "確定評価セット（学習にも検証にも使っていない1,800クリップ）で、"
      "危険な車508台を同じ採点方法で比較。判定の規則だけを差し替えた",
      size=11, color=MUTED)
label(54, 494, 860,
      "「3秒前に検知できても使えない」というご指摘に対する答え",
      size=13, color=INK, bold=True)

out = ROOT / "md/seminar/図_検証_通知層_簡潔_2026-08-20.pptx"
try:
    prs.save(out)
except PermissionError:
    out = out.with_name(out.stem + "_新" + out.suffix)
    prs.save(out)
    print("※ 元ファイルが開かれていたため別名で保存した")
print(f"saved: {out} ({out.stat().st_size // 1024} KB)")
