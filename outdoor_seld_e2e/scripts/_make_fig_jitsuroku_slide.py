# -*- coding: utf-8 -*-
"""「実録評価の測定方法」スライドを本人テンプレ風デザインで1枚生成。

背景: 2026-08-22、旧版(3枠)は 統制録音 / 機会捕捉 / 負例 の3分類で、
バック音を機会捕捉側に置いていた。しかし事前登録(実録スモーク計画書 R1)では
バック音は「統制4本（実物ブザー・S3幾何）」であり、分類が食い違っていた。
「自分で決められるか」で切り直して4分類にしたのがこの版。

分類の根拠 = md/design/実録スモーク計画書_2026-07.md R1 / 実録ハンドブック 6節
出力: md/seminar/図_実録の測定方法_2026-08-22.pptx（コピペ用・1枚）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
BAND = RGBColor(0xF2, 0xF3, 0xF5)
BANDR = RGBColor(0xFA, 0xEE, 0xEC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


# ---- テンプレ骨格 ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 160, 40, [("今後の方針②　実録評価の測定方法（9月）",
                           {"size": 23, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

txt(70, 100, W - 140, 22,
    [("合成データで作った通知を実環境で測り直す。共通100本を", {"size": 13.5, "color": SUB}),
     ("「自分で決められるか」で4つに分けて", {"size": 13.5, "bold": True}),
     ("録る。", {"size": 13.5, "color": SUB})])

# ---- 4分類カード ----
BY, BH, BW = 128, 142, 194
bxs = [70, 278, 486, 694]
cards = [
    ("① 統制録音", "64", PUR,
     ["幾何を自分で決めて再現", "・車の通過　20",
      "・弱点　20（EV・発進・", "　自転車・見通し不良の角）",
      "・キックボード　20", "・バック音　4（実物ブザー）"]),
    ("② 確実に録れる", "8", GREEN,
     ["時刻表どおりに待てば録れる", "・実在の踏切と列車　8",
      "", "自分で操作はしないが、", "計画に入れられる枠"]),
    ("③ 機会捕捉", "8", GOLD,
     ["意思では増やせない", "・サイレン・クラクション等",
      "", "上限目標。遭遇したら最優先。", "0本でも評価は成立する"]),
    ("④ 負例", "20", MUTED,
     ["鳴ってはいけない場面", "・静穏・繁華街/工事",
      "・雨上がり・風", "＋別枠で連続録音100分以上",
      "→ 誤警告率［回/時］"]),
]
for (tag, num, col, lines), bx in zip(cards, bxs):
    rect(bx, BY, BW, BH, WHITE, line=LINE, lw=1.0)
    rect(bx, BY, 4, BH, col)
    txt(bx + 14, BY + 9, BW - 70, 20, [(tag, {"size": 13, "bold": True})])
    txt(bx + BW - 62, BY + 5, 50, 26,
        [(num, {"size": 21, "bold": True, "color": col}),
         ("本", {"size": 10.5, "color": SUB})], align=PP_ALIGN.RIGHT)
    body = []
    for i, t in enumerate(lines):
        if not t:
            body.append([(" ", {"size": 5})])
        else:
            body.append([(t, {"size": 10.5,
                              "color": (INK if i == 0 else SUB),
                              "bold": (i == 0)})])
    txt(bx + 14, BY + 34, BW - 26, BH - 42, body, line=1.32)

# ---- 機材 ----
RY, RH = 284, 148
rect(70, RY, 300, RH, WHITE, line=LINE, lw=1.0)
txt(86, RY + 12, 270, 20, [("持っていく機材", {"size": 14, "bold": True})])
txt(86, RY + 40, 272, 100,
    [[("・4ch FOAマイク Zoom H3-VR ＋ 風防", {"size": 11.5, "color": SUB})],
     [("・チェストハーネス（全テイク統一）", {"size": 11.5, "color": SUB})],
     [("・騒音計（LAeq）・メジャー・コーン", {"size": 11.5, "color": SUB})],
     [("・自転車ベル（方位校正にも使う）", {"size": 11.5, "color": SUB})],
     [("・実物のバックブザー（①の4本用）", {"size": 11.5, "color": SUB})]], line=1.42)

# ---- 統制録音の幾何（横距離の帯） ----
rect(386, RY, 504, RH, WHITE, line=LINE, lw=1.0)
txt(402, RY + 12, 300, 20, [("統制録音の幾何（車）", {"size": 14, "bold": True})])
txt(700, RY + 14, 176, 18,
    [("速度：徐行〜約30km/h", {"size": 11, "color": SUB})], align=PP_ALIGN.RIGHT)

txt(402, RY + 34, 470, 16,
    [("歩行者は静止なので、最接近距離＝横距離になる", {"size": 10, "color": MUTED})])

bands = [
    ("≤1.5m", RED, "最接近予測が働く帯", "自転車5・キック20", "車では再現しない"),
    ("2.0〜3.2m", GOLD, "距離の保険が働く帯", "車10", None),
    ("5〜15m", GREEN, "鳴らないはずの帯", "車10", None),
]
byy = RY + 54
for rng, col, role, who, note in bands:
    rect(402, byy, 5, 28, col)
    rect(407, byy, 469, 28, BANDR if col is RED else BAND)
    txt(417, byy + 6, 76, 18, [(rng, {"size": 11.5, "bold": True})])
    txt(499, byy + 6, 150, 18, [(role, {"size": 11, "bold": True, "color": col})])
    txt(655, byy + 6, 130, 18, [(who, {"size": 10.5, "color": SUB})])
    if note:
        txt(781, byy + 6, 88, 18,
            [(note, {"size": 10, "color": RED})], align=PP_ALIGN.RIGHT)
    byy += 32

# ---- 金色バー ----
rect(70, 444, 6, 50, GOLD)
rect(76, 444, W - 146, 50, WHITE, line=LINE, lw=1.0)
txt(94, 452, W - 180, 38, [
    ("共通100本 ＋ 歩行の影響を見る対比100本（50ペア）＝ 計200本", {"bold": True, "size": 14}),
    ("　主要評価は「caution帯の通知到達率」1つに事前登録", {"size": 13, "color": SUB}),
], anchor=MSO_ANCHOR.MIDDLE)

# ---- フッター ----
txt(58, 506, 90, 20, [("2026/09", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "今後・まとめ")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(10.5)
    r.font.bold = active
    r.font.color.rgb = WHITE if active else MUTED
    r.font.name = "Meiryo"
    meiryo(r)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("11", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"図_実録の測定方法_2026-08-22.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
