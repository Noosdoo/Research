# -*- coding: utf-8 -*-
"""【修正スライド1枚だけ】Joy-conデモ — 「鳴る車と鳴らない車」の時間軸図。

本人の依頼（2026-09-03）: 発火のタイミング図を作る。動画が再生できないときの保険と、
Q2（すり抜けと対向）への回答の裏付けを兼ねる。

素材 = out/joycon_demo/ の custom_rei1_taikou_to_surinuke（Unityのクリップ 4/10）
 車① 正面から接近・横1.0m  → 最接近6.0s。3.6sで強域に入り、4.0sに振動
 車② 後方から追い越し・横4.5m → 最接近8.5s。一度も鳴らない
帯と距離は state.csv を読んで描く（手入力しない）。振動の印は cues.csv から。

⚠️ このクリップは**オラクル動作**（正解の軌跡＋v4.2）。モデル出力ではない。
   図にもその旨を明記してある。

（当初は fold31_room1_mix0142 を推したが、state.csv を読むと車が3台重なり
  語りにくいと判明したため rei1 に変更した。）

出力: md/seminar/修正_デモ時間軸_2026-09-15.pptx（1枚のみ）
"""
import csv
import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
DEMO = ROOT / "out/joycon_demo"
BASE = "custom_rei1_taikou_to_surinuke"

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
FAINT = RGBColor(0xE4, 0xE6, 0xEA)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF5, 0xF3, 0xF9)
PALEG = RGBColor(0xFD, 0xF7, 0xE8)

TIER = {"強": ("★ 至近警告", RED), "中": ("▲ 注意", GOLD),
        "安全": ("・ 抑制（鳴らさない）", GREEN), "計測中": ("計測中", MUTED)}

# ---------------- state.csv を読む ----------------
def tier_of(part):
    if "至近警告" in part:
        return "強"
    if "注意" in part:
        return "中"
    if "抑制" in part:
        return "安全"
    return "計測中"


lanes = {0: [], 1: []}          # [(t, tier, dist), ...]
for line in io.open(DEMO / (BASE + "_state.csv"), encoding="utf-8").read().splitlines():
    if "|" not in line:
        continue
    ts, body = line.split("|", 1)
    t = float(ts)
    parts = [p for p in body.split("／") if p.strip()]
    for i, p in enumerate(parts):
        if i > 1:
            continue
        m = re.search(r"([0-9.]+)m", p)
        lanes[i].append((t, tier_of(p), float(m.group(1)) if m else None))

# 帯（同じtierが続く区間）にまとめる
def bands(seq):
    out = []
    for t, tr, _d in seq:
        if out and out[-1][2] == tr:
            out[-1][1] = t
        else:
            out.append([t, t, tr])
    return out


# 最接近の時刻＝距離が最小になる時刻
def cpa_of(seq):
    ds = [(d, t) for t, _tr, d in seq if d is not None]
    d, t = min(ds)
    return t, d


# ---------------- cues.csv を読む ----------------
fires = []
for row in csv.DictReader(io.open(DEMO / (BASE + "_cues.csv"), encoding="utf-8")):
    if row["t_s"].startswith("#"):
        continue
    fires.append((float(row["t_s"]), row["tier"]))

# ---------------- スライド ----------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])
SECTIONS = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 170, 40,
    [("④　同 じ 場 面 で 、 鳴 る 車 と 鳴 ら な い 車",
      {"size": 22, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

txt(70, 100, 820, 20,
    [("助言どおりJoy-conとUnityで作り、9/2に実機で振動した。1本のクリップの中に、"
      "鳴る車と鳴らない車が両方入っている。", {"size": 12, "color": MUTED})])

# ---- 上段：デモの素性 ----
facts = [("使ったもの", "SwitchのJoy-con と Unity 6.6（ハンダごて不要・部品購入ゼロ）"),
         ("伝え方", "強＝4連打 ／ 中＝2発 ／ 警告音＝単発のパルス")]
fy = 128
for a, b in facts:
    rect(70, fy + 6, 6, 6, PUR)
    txt(88, fy, 100, 18, [(a, {"size": 10.5, "bold": True})])
    txt(190, fy, 700, 18, [(b, {"size": 10.5, "color": SUB})])
    fy += 22

# ---- 時間軸 ----
TX0, TX1 = 216.0, 872.0
TMAX = 10.0
tp = lambda t: TX0 + (TX1 - TX0) * t / TMAX

AX = 194.0                      # 目盛りの基線
rect(TX0, AX, TX1 - TX0, 0.8, HAIR)
for s in range(0, 11, 2):
    rect(tp(s), AX - 4, 0.8, 5, MUTED)
    txt(tp(s) - 20, AX - 20, 40, 14,
        [("%ds" % s, {"size": 8.5, "color": MUTED})], align=PP_ALIGN.CENTER)

BH = 26.0
LANES = [
    (0, "車①　正面から来る", "横1.0m・30km/h", RED, 214.0),
    (1, "車②　後ろから追い越す", "横4.5m・40km/h", GREEN, 316.0),
]
for idx, name, cond, col, ly in LANES:
    seq = lanes[idx]
    txt(70, ly + 2, 140, 18, [(name, {"size": 11.5, "bold": True})])
    txt(70, ly + 20, 140, 16, [(cond, {"size": 9, "color": MUTED})])
    for t0, t1, tr in bands(seq):
        _lab, c = TIER[tr]
        x, w = tp(t0), max(tp(t1) - tp(t0), 2.0)
        if tr == "計測中":
            rect(x, ly + 8, w, BH - 12, FAINT)
        elif tr == "安全":
            rect(x, ly, w, BH, WHITE, line=GREEN, lw=1.2)
        else:
            rect(x, ly, w, BH, c)
    # 段階の名前を帯の中／下に
    for t0, t1, tr in bands(seq):
        if tr in ("強", "中", "安全") and tp(t1) - tp(t0) > 40:
            nm = {"強": "至近警告（強）", "中": "注意（中）",
                  "安全": "抑制（鳴らさない）"}[tr]
            txt(tp(t0) + 4, ly + 5, tp(t1) - tp(t0) - 8, 18,
                [(nm, {"size": 9.5, "bold": True,
                       "color": WHITE if tr != "安全" else GREEN})])
    # 最接近
    ct, cd = cpa_of(seq)
    rect(tp(ct) - 0.7, ly - 14, 1.4, BH + 28, INK)
    txt(tp(ct) - 60, ly + BH + 16, 120, 16,
        [("最接近 %.1fm" % cd, {"size": 9.5, "bold": True})],
        align=PP_ALIGN.CENTER)

# ---- 振動の印（車①のレーンの上） ----
txt(70, 176, 140, 16, [("振動", {"size": 9.5, "color": MUTED})])
for t, tr in fires:
    rect(tp(t) - 5, 178, 10, 10, RED, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
txt(tp(fires[0][0]) - 66, 160, 132, 16,
    [("4.0s に最初の振動", {"size": 9.5, "bold": True, "color": RED})],
    align=PP_ALIGN.CENTER)

# ---- リードタイムの矢印 ----
ct0 = cpa_of(lanes[0])[0]
LY = 288.0
rect(tp(fires[0][0]), LY, tp(ct0) - tp(fires[0][0]), 1.4, RED)
txt(tp(fires[0][0]), LY - 16, tp(ct0) - tp(fires[0][0]), 16,
    [("最接近の %.1f 秒前" % (ct0 - fires[0][0]),
      {"size": 10, "bold": True, "color": RED})], align=PP_ALIGN.CENTER)

# ---- 読み取り ----
rect(70, 398, 400, 56, PALEG)
rect(70, 398, 4, 56, GOLD)
txt(88, 406, 370, 44,
    [[("車①：正面から来る車", {"size": 11, "bold": True})],
     [("方位が変わらないまま近づくので、最接近の2秒前に強く鳴る。",
       {"size": 10.5, "color": SUB})]], line=1.4)
rect(490, 398, 400, 56, WHITE, line=GREEN, lw=1.0)
txt(508, 406, 370, 44,
    [[("車②：後ろから追い越す車", {"size": 11, "bold": True})],
     [("横4.5mを通り過ぎるだけなので、一度も鳴らない。",
       {"size": 10.5, "color": SUB})]], line=1.4)

txt(70, 462, 820, 16,
    [("※ このクリップは正解の軌跡を入力にした動作確認（オラクル）。モデル出力ではない。",
      {"size": 9, "color": MUTED})])

# ---- フッター ----
txt(58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
ws = [16 + len(s) * 12.0 for s in SECTIONS]
x = (W - (sum(ws) + 8 * (len(SECTIONS) - 1))) / 2
for s, wd in zip(SECTIONS, ws):
    on = (s == "振り返り")
    c = rect(x, 502, wd, 22, NAVY if on else CHIP)
    label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("11", {"size": 11, "color": MUTED})],
    align=PP_ALIGN.RIGHT)

OUT = ROOT / "md/seminar/修正_デモ時間軸_2026-09-15.pptx"
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ PowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
print("車①", cpa_of(lanes[0]), "／ 車②", cpa_of(lanes[1]), "／ 振動", fires)
