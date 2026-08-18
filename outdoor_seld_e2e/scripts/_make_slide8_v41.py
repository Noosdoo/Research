# -*- coding: utf-8 -*-
"""中間発表スライド8枚目「提案手法 — 2層の仕組み」を通知v4.1に合わせて作り直す。

本人が編集している手作りデッキ（游ゴシック・25枚）に貼り込めるよう、
**この1枚だけ**の pptx を出す。既存デッキは触らない。

右の図を同心円から差し替えたのが要点。同心円は「距離で決まる」という図であり、
v4.1（最接近を予測して鳴らす）を表せない。代わりに
**正面に来る相手（鳴る）と横を通り過ぎる相手（鳴らない）**を並べ、
両者を分けているのが「方位が流れるかどうか」であることを示す。

出力: md/seminar/図_提案手法2層_v4.1_2026-08-18.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

W, H = 960, 540
INK = RGBColor(0x16, 0x23, 0x3A)
INK2 = RGBColor(0x3A, 0x46, 0x58)
MUTED = RGBColor(0x66, 0x70, 0x7F)
HAIR = RGBColor(0xD9, 0xDA, 0xD2)
AMBERD = RGBColor(0xB3, 0x7E, 0x00)
RED = RGBColor(0xC4, 0x43, 0x2B)
GREEN = RGBColor(0x2F, 0x7D, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xBF, 0xC6, 0xCC)
JP = "游ゴシック"          # 本人のデッキに合わせる

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(int(W * 12700)), Emu(int(H * 12700))
sl = prs.slides.add_slide(prs.slide_layouts[6])


def _font(run, size, bold, color):
    f = run.font
    f.name = JP
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", JP)


def box(x, y, w, h):
    tb = sl.shapes.add_textbox(Emu(int(x * 12700)), Emu(int(y * 12700)),
                               Emu(int(w * 12700)), Emu(int(h * 12700)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb


def para(tf, runs, size=13, bold=False, color=INK2, align=PP_ALIGN.LEFT,
         first=False, line=1.3, space=2):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    p.space_after = Pt(space)
    for text, kw in ([(runs, {})] if isinstance(runs, str) else runs):
        r = p.add_run()
        r.text = text
        _font(r, kw.get("size", size), kw.get("bold", bold), kw.get("color", color))
    return p


def rect(x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
         dash=None):
    s = sl.shapes.add_shape(shape, Emu(int(x * 12700)), Emu(int(y * 12700)),
                            Emu(int(w * 12700)), Emu(int(h * 12700)))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
        if dash:
            ln = s.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    s.text_frame.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(s.text_frame, m, Pt(10 if "left" in m or "right" in m else 8))
    s.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    return s


def conn(x1, y1, x2, y2, color, w=2.0, dash=None, arrow=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1 * 12700)),
                                Emu(int(y1 * 12700)), Emu(int(x2 * 12700)),
                                Emu(int(y2 * 12700)))
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    if arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return c


def label(x, y, w, text, size=10.5, color=MUTED, align=PP_ALIGN.CENTER,
          bold=False, wrap=False):
    tb = box(x, y, w, 18)
    tb.text_frame.word_wrap = wrap
    para(tb.text_frame, text, size=size, bold=bold, color=color, align=align,
         first=True, space=0)


def card(x, y, w, h, title, border=HAIR, lw=1.0, tcolor=INK):
    c = rect(x, y, w, h, fill=WHITE, line=border, lw=lw)
    para(c.text_frame, title, size=14, bold=True, color=tcolor, first=True, space=4)
    return c


# ============ 見出し ============
conn(54, 96, W - 54, 96, INK, 1.5)
tb = box(54, 40, 700, 48)
para(tb.text_frame, "提案手法 ― ２層の仕組み", size=27, bold=True, color=INK,
     first=True, space=0)

# ============ 左：3段の流れ ============
LX, LW = 54, 452
c = card(LX, 112, LW, 58, "知覚層(SELD＋SDE)")
para(c.text_frame, "8クラスの検出・方向・距離を0.1秒ごとに推定", size=12.5)
label(LX, 174, LW, "▼", size=11, color=INK2)

c = card(LX, 192, LW, 176, "通知層(3役割) ― 最接近の予測で出し分け",
         border=AMBERD, lw=1.5)
para(c.text_frame, [("近距離", {"bold": True, "color": RED}),
                    ("：このまま進むと約1m以内・2.5秒以内に到達 … ", {}),
                    ("4フレーム連続", {"bold": True})], size=12.5)
para(c.text_frame, [("中距離", {"bold": True, "color": AMBERD}),
                    ("：約2m以内・4秒以内 … ", {}),
                    ("4フレーム連続", {"bold": True})], size=12.5)
para(c.text_frame, [("遠距離", {"bold": True, "color": GREEN}),
                    ("：横を通り過ぎるだけの対象は鳴らさない", {})], size=12.5)
para(c.text_frame, "※ 既に1.5m/3.2m以内なら距離だけでも鳴らす(保険・2フレーム連続)",
     size=11.5, color=MUTED)
para(c.text_frame, [("対象は", {}), ("車・キックボード・バイク", {"bold": True}),
                    ("。サイレン等5クラスは距離を使わず検出したら通知", {})],
     size=11.5, color=MUTED)
label(LX, 372, LW, "▼", size=11, color=INK2)

c = card(LX, 390, LW, 56, "首元振動デバイス")
para(c.text_frame, "方向 × 振動でユーザに伝える", size=12.5)

label(LX, 458, LW, "全部鳴らすと使えない → 通知の頻度そのものを設計対象に",
      size=12, color=INK, bold=True, align=PP_ALIGN.LEFT)

# ============ 右：正面に来る相手 vs 横を通り過ぎる相手 ============
BX, BY, BW, BH = 526, 112, W - 54 - 526, 300
rect(BX, BY, BW, BH, fill=WHITE, line=HAIR)
label(BX, BY + 10, BW, "同じ「近づいてくる」でも、鳴らすべき相手は片方だけ",
      size=11.5, color=INK, bold=True)

PX, PY = BX + 62, BY + 208           # 歩行者の位置
GY = BY + 78                         # 横を通り過ぎる相手の経路

# 距離しきい値（保険）は薄い点線で残す
rect(PX - 46, PY - 46, 92, 92, fill=None, line=FAINT, lw=1.0,
     shape=MSO_SHAPE.OVAL, dash="dash")
rect(PX - 22, PY - 22, 44, 44, fill=None, line=FAINT, lw=1.0,
     shape=MSO_SHAPE.OVAL, dash="dash")

# --- 横を通り過ぎる相手（上）: 方位が流れる ---
conn(BX + BW - 22, GY, PX - 34, GY, GREEN, 2.5, arrow=True)
for gx in (BX + BW - 40, BX + 150):
    conn(PX, PY, gx, GY, FAINT, 1.0, dash="sysDot")
    rect(gx - 5, GY - 5, 10, 10, fill=GREEN, shape=MSO_SHAPE.OVAL)
label(BX + 12, GY - 30, BW - 26, "横を通り過ぎる → 鳴らさない", size=12,
      color=GREEN, bold=True, align=PP_ALIGN.RIGHT)
label(BX + 30, GY + 12, 150, "方位が横に流れる", size=10.5, color=GREEN,
      align=PP_ALIGN.LEFT)

# 予測した最接近距離（ここまでしか近づかない）
conn(PX, PY - 46, PX, GY, MUTED, 1.0, dash="dash")
label(PX - 96, (GY + PY - 46) / 2 - 9, 88, "最接近", size=10, color=MUTED,
      align=PP_ALIGN.RIGHT)

# --- 正面に来る相手（下）: 方位が変わらない ---
conn(BX + BW - 22, PY, PX + 56, PY, RED, 2.5, arrow=True)
for rx in (BX + BW - 40, BX + 186):
    rect(rx - 5, PY - 5, 10, 10, fill=RED, shape=MSO_SHAPE.OVAL)
label(BX + 12, PY - 32, BW - 26, "正面に来る → 至近警告", size=12, color=RED,
      bold=True, align=PP_ALIGN.RIGHT)
label(BX + 12, PY + 14, BW - 26, "方位が変わらないまま近づく", size=10.5,
      color=RED, align=PP_ALIGN.RIGHT)

rect(PX - 7, PY - 7, 14, 14, fill=INK, shape=MSO_SHAPE.OVAL)
label(PX - 50, PY + 52, 100, "歩行者", size=10.5, color=INK2)
label(BX + 8, PY + 72, BW - 16, "点線の円＝保険として残した距離しきい値(1.5m / 3.2m)",
      size=9.5, color=MUTED)

label(BX, BY + BH + 8, BW,
      "船の見張りと同じ原理 ― 方位が変わらず近づく相手が危ない",
      size=11, color=INK2)
label(BX, BY + BH + 26, BW,
      "※「変わらない」の許容幅は固定値ではなく距離と速さで決まる",
      size=9.5, color=MUTED)
label(BX, BY + BH + 40, BW,
      "（時速50kmの車なら 20m先で毎秒2度以内 / 10m先で毎秒8度以内）",
      size=9.5, color=MUTED)

out = ROOT / "md/seminar/図_提案手法2層_v4.1_2026-08-18.pptx"
prs.save(out)
print(f"saved: {out} ({out.stat().st_size // 1024} KB)")
