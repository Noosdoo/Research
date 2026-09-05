# -*- coding: utf-8 -*-
"""「①すり抜けと対向を、どう区別するか」の図を直し、「①その結果」に型の分け方を足す。

背景（2026-09-06・監査 R03 と同じ取り違え）:
  図のキャプション「判定に使うのは、最接近の2.5〜1.5秒前の区間」は誤り。
  通知層の判定は毎フレーム、直近5フレーム＝0.5秒の方位の傾きで行い、
  最接近がいつ来るかは知らない（因果処理）。
    step12_notify_v4_ttc.py azimuth_rate(win=VEL_WIN=5)、FPS=10
    step12_notify_v43.py L60-64（route_c: |adot|≤0.10 ∧ 接近 ∧ d≤15m）
  「最接近の2.5〜1.5秒前」は、①その結果の型別表を作るときに正解の軌跡から
  対向型／すり抜け型を事後に分類した窓（_notify_v42_q2_table.py gt_adot_before_cpa）。
  判定の窓ではないので、説明を次ページの型別表の注記へ移す。

  もう1つ: 対向型の赤い矢頭が装着者から離れる向きに読めた（向かってくる車なのに）。
  矢頭を装着者向きに、緑は進行方向（右）に揃える。

本人のデッキから2枚を複製して直すので、書式は元のまま変わらない。
出力: md/seminar/修正_すり抜けと対向_2026-09-15.pptx（2枚: p7 / p8）
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = ROOT / "md/seminar/修正_すり抜けと対向_2026-09-15.pptx"

T7 = "すり抜けと対向を、どう区別するか"
T8 = "その結果"

CAPTION_OLD = "判定に使うのは、最接近の2.5〜1.5秒前の区間"
CAPTION_NEW = "判定は毎フレーム、直近0.5秒の方位の傾きで行う（最接近がいつかは知らない）"
NOTE_NEW = ("※型の分け方＝正解の軌跡で、最接近の2.5〜1.5秒前の方位変化率が "
            "0.10 rad/s 未満なら対向型、以上ならすり抜け型。")

RED, GREEN = (0xC0, 0x39, 0x2B), (0x2E, 0x7D, 0x5B)


def find(prs, title):
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and title in sh.text_frame.text \
                    and sh.top / 12700 < 70:
                return sl
    raise SystemExit("見出し「%s」のスライドが見つからない" % title)


def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(
            copy.deepcopy(shp._element), "p:extLst")
    return dst


def settext(sp, lines):
    """段落ごとに本文を差し替える。1段落目の書式を雛形として引き継ぐ。"""
    tf = sp.text_frame
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, line in enumerate(lines):
        p = tf.paragraphs[i]
        if not p.runs:
            src = tf.paragraphs[0].runs[0]
            r = p.add_run()
            r.font.size, r.font.bold, r.font.name = src.font.size, src.font.bold, src.font.name
            try:
                r.font.color.rgb = src.font.color.rgb
            except Exception:
                pass
            for ea in src._r.findall(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr/"
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}ea"):
                r._r.get_or_add_rPr().append(copy.deepcopy(ea))
        p.runs[0].text = line
        for r in list(p.runs[1:]):
            r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[len(lines):]):
        p._p.getparent().remove(p._p)


def pick_text(sl, needle):
    hits = [s for s in sl.shapes
            if s.has_text_frame and s.text_frame.text.strip().startswith(needle)]
    assert len(hits) == 1, "「%s」が %d 個ある" % (needle, len(hits))
    return hits[0]


def pick_triangle(sl, rgb):
    hits = []
    for s in sl.shapes:
        try:
            if "ISOSCELES_TRIANGLE" in str(s.auto_shape_type) and s.fill.type == 1 \
                    and tuple(s.fill.fore_color.rgb) == rgb:
                hits.append(s)
        except Exception:
            pass
    assert len(hits) == 1, "三角 %s が %d 個ある" % (rgb, len(hits))
    return hits[0]


def fix_footer(sl):
    """複製で日付・ページ番号が空レイアウトの位置に戻るので置き直す。"""
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == "2026/09/15":
            sh.left, sh.top, sh.width, sh.height = (int(Pt(66.0)), int(Pt(500.5)),
                                                    int(Pt(216.0)), int(Pt(28.8)))
        elif t.isdigit() and sh.top / 12700 > 480:
            sh.left, sh.top, sh.width, sh.height = (int(Pt(678.0)), int(Pt(500.5)),
                                                    int(Pt(216.0)), int(Pt(28.8)))


def fix_p7(sl):
    cap = pick_text(sl, CAPTION_OLD)
    settext(cap, [CAPTION_NEW])
    pick_triangle(sl, RED).rotation = 180.0     # 矢頭を装着者向きに（向かってくる）
    pick_triangle(sl, GREEN).rotation = 90.0    # 進行方向（右）に
    fix_footer(sl)


def fix_p8(sl):
    note = pick_text(sl, "※合成データ")
    old = note.text_frame.paragraphs[0].runs[0].text
    settext(note, [old, NOTE_NEW])
    note.top, note.height = int(Pt(458.0)), int(Pt(40.0))   # 2行ぶんに広げる
    fix_footer(sl)


def main() -> int:
    src = Presentation(str(SRC))
    out = Presentation()
    out.slide_width, out.slide_height = src.slide_width, src.slide_height
    fix_p7(clone(find(src, T7), out))
    fix_p8(clone(find(src, T8), out))

    dst = OUT
    try:
        out.save(str(dst))
    except PermissionError:
        dst = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
        out.save(str(dst))
    print("保存:", dst)
    print("p7: キャプション →", CAPTION_NEW)
    print("p7: 赤三角 180° / 緑三角 90°")
    print("p8: 注記に追加 →", NOTE_NEW)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
