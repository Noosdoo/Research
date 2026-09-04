# -*- coding: utf-8 -*-
"""「関連研究のまとめ ― どこまでできていて、どこから踏み出すか」スライドを1枚生成。

背景: 2026-08-25の中間発表リハで3人の教員から同じ指摘を受けた。
  「先行研究がどこまで実現していて、どういう課題なのかを示した方がいい」
  「研究目的のスコープが明確にならない」「自分がどこの一歩に踏み出すのか」
現行デッキの関連研究(p5 PSELDNets / p6 DynamicSound)は**道具の研究**だけで、
同じ問題に取り組んだ先行研究(PAWS等)が1枚も出ていなかった。その穴を埋める1枚。

内容の出典 = md/research/新規性サーベイ更新_2026-07-16.md（must-cite 3件）
2026-08-25に一次情報で照合（本人「正しい？根拠は？」）。その結果:
  - PAWS 80m/99%/97%/30mphで約6秒前 = Columbia ICSL 公式ページと一致（確認済み）
  - PAWS の対象は難聴者ではなく「通話・音楽で聞こえにくい聴者」= 公式ページの明記（追加した）
  - Sensors2023 = 7音（**感情音声4＋警告音3**）・振動＋OLED・モデル97.05%／エッジ95.22%
    （08-25に「97%は誤り」と書いたのは誤り。モデル値97.05%とエッジ値95.22%の2つがある）
  - 6DoF SELD の「歩行者安全を動機に」は abstract に無い → 記述から削除した
  - 「当事者調査で方向が最重視」は公開範囲で裏が取れず → ※行ごと削除した
挿入位置: 現p4(SELDとは)の後（＝新p5）。道具の研究(PSELDNets/DynamicSound)より**前**に
          置く。足りないものを先に言うことで、道具を使う必要性が立つため（2026-08-25 本人指摘）。
          以降のページ番号が1つずつ繰り下がる
出力: md/seminar/図_先行研究の到達点_2026-08-25.pptx（コピペ用・1枚）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
HEAD = RGBColor(0x3A, 0x41, 0x55)
ZEBRA = RGBColor(0xF6, 0xF7, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


# ---- テンプレ骨格 ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 160, 40, [("先行研究　―　どこまでできていて、どこから踏み出すか",
                           {"size": 22, "bold": True, "spc": 220})])
rect(44, 90, W - 90, 1.2, INK)

# ---- 表 ----
X0, TW = 70, W - 140
C1, C2 = 236.0, 342.0
C3 = TW - C1 - C2
TY, RH, HH = 116.0, 74.0, 30.0

rect(X0, TY, TW, HH, HEAD)
for cx, cw, t in [(X0, C1, "先行研究"), (X0 + C1, C2, "どこまで到達しているか"),
                  (X0 + C1 + C2, C3, "足りないもの")]:
    txt(cx + 14, TY + 6, cw - 20, 20,
        [(t, {"size": 12.5, "bold": True, "color": WHITE})])

rows = [
    ("PAWS （2018）", "歩行者向け 車両接近警告",
     ["接近する車を80m先から検出し、", "時速48kmなら約6秒前に警告"],
     ["車だけ（サイレン・踏切などは対象外）／複数台は扱えない",
      "主な想定は通話・音楽聴取中の歩行者（難聴者を明示的には扱わない）"]),
    ("サイレン認識ウェアラブル （2023）", "難聴者向け",
     ["7つの音を高い精度で聞き分け、", "振動とディスプレイで伝える"],
     ["方向がない（種類だけ）", "同時に1つの音しか扱えない"]),
    ("6DoF SELD （2024）", "装着マイクでのSELD",
     ["歩行中の人に装着したマイクでSELDを実現", "（自己移動をセンサで補償）"],
     ["距離がない", "通知の仕組みがない"]),
]
y = TY + HH
for i, (name, sub, mid, gap) in enumerate(rows):
    if i % 2 == 1:
        rect(X0, y, TW, RH, ZEBRA)
    rect(X0, y + RH, TW, 0.8, LINE)
    txt(X0 + 14, y + 12, C1 - 20, 40,
        [[(name, {"size": 12, "bold": True})],
         [(sub, {"size": 10.5, "color": MUTED})]], line=1.3)
    txt(X0 + C1 + 14, y + 14, C2 - 20, 44,
        [[(t, {"size": 11.5, "color": SUB})] for t in mid], line=1.32)
    txt(X0 + C1 + C2 + 14, y + 14, C3 - 20, 44,
        [[("・", {"size": 11.5, "color": RED}),
          (t, {"size": 11.5, "bold": True})] for t in gap], line=1.32)
    y += RH

# ---- 本研究の一歩 ----
BY = y + 26
rect(X0, BY, 6, 62, GOLD)
rect(X0 + 6, BY, TW - 6, 62, WHITE, line=LINE, lw=1.0)
txt(X0 + 24, BY + 10, TW - 60, 46,
    [[("本研究の一歩：", {"size": 14, "bold": True, "color": PUR}),
      ("8クラス", {"size": 14, "bold": True}),
      (" × ", {"size": 14, "color": SUB}),
      ("方向", {"size": 14, "bold": True}),
      (" × ", {"size": 14, "color": SUB}),
      ("距離", {"size": 14, "bold": True}),
      (" を1つにまとめ、", {"size": 14, "color": SUB})],
     [("　　　　　　　　その上に", {"size": 14, "color": SUB}),
      ("「危険度へ変える規則」", {"size": 14, "bold": True, "color": GOLD}),
      ("を載せる", {"size": 14, "color": SUB})]], line=1.5)

# ---- フッター ----
txt(58, 506, 90, 20, [("2026/08/30", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "関連研究")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(10.5)
    r.font.bold = active
    r.font.color.rgb = WHITE if active else MUTED
    r.font.name = "Meiryo"
    meiryo(r)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("5", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"図_先行研究の到達点_2026-08-25.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
