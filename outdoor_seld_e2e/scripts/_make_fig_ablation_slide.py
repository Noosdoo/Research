# -*- coding: utf-8 -*-
"""「今後の方針① ablation実験計画」スライドを本人テンプレ風デザインで1枚生成。

背景: 2026-08-23、旧版は3列（規模／測り方／位置付け）の箇条書きで、右の位置付け列が
上の「動機」行と同じ内容を繰り返していた。動機行だけ削ったところ上部が空いて寂しく
なったため、5条件を**図**にして空間を埋める構成に作り直したもの。

構成: 問い → 5条件を横並びの箱で図示 → 同じ評価セットで比較 → 何で測るか2枚 → 締め
出力: md/seminar/図_ablation実験計画_2026-08-23.pptx（コピペ用・1枚）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
FAINT = RGBColor(0xF4, 0xF5, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


# ---- テンプレ骨格 ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 160, 40, [("今後の方針①　ablation実験計画",
                           {"size": 23, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

# ---- 問い ----
txt(70, 102, W - 140, 26,
    [("問い：", {"size": 14, "bold": True, "color": PUR}),
     ("合成データに入れた物理のうち、", {"size": 14, "color": SUB}),
     ("どれを外すと、どの性能が落ちるのか？", {"size": 14, "bold": True})])

# ---- 5条件の箱 ----
BY, BH, BW = 142, 96, 154
bxs = [70, 236, 402, 568, 734]

# 基準
rect(bxs[0], BY, BW, BH, NAVY)
txt(bxs[0], BY + 20, BW, 22, [("基準", {"size": 13, "bold": True, "color": WHITE})],
    align=PP_ALIGN.CENTER)
txt(bxs[0], BY + 48, BW, 22,
    [("全物理あり", {"size": 15, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER)

offs = ["距離減衰", "ドップラー", "地面反射", "大気吸収"]
for name, bx in zip(offs, bxs[1:]):
    rect(bx, BY, BW, BH, WHITE, line=LINE, lw=1.0)
    rect(bx, BY, BW, 4, GOLD)
    txt(bx, BY + 22, BW, 22,
        [("この1つだけ外す", {"size": 10.5, "color": MUTED})], align=PP_ALIGN.CENTER)
    txt(bx, BY + 46, BW, 24,
        [("− ", {"size": 15, "bold": True, "color": RED}),
         (name, {"size": 15, "bold": True})], align=PP_ALIGN.CENTER)

txt(bxs[1], BY + BH + 6, bxs[4] + BW - bxs[1], 20,
    [("物理を1つだけ外した4条件", {"size": 11.5, "color": SUB})], align=PP_ALIGN.CENTER)

# ---- 中央の帯（同じ土俵で比べる） ----
rect(70, 272, W - 140, 30, FAINT)
txt(70, 278, W - 140, 22,
    [("計5条件を", {"size": 12.5, "color": SUB}),
     ("同じ評価セットで採点し、基準からの性能差を測る", {"size": 12.5, "bold": True})],
    align=PP_ALIGN.CENTER)

# ---- 何で測るか2枚 ----
CY, CH = 318, 108
rect(70, CY, 395, CH, WHITE, line=LINE, lw=1.0)
rect(70, CY, 4, CH, PUR)
txt(88, CY + 12, 360, 22, [("合成データで測る", {"size": 14, "bold": True})])
txt(88, CY + 40, 360, 60,
    [[("・SELD性能（検出・方向）", {"size": 12, "color": SUB})],
     [("・距離推定の性能", {"size": 12, "color": SUB})]], line=1.45)

rect(495, CY, 395, CH, WHITE, line=LINE, lw=1.0)
rect(495, CY, 4, CH, GREEN)
txt(513, CY + 12, 360, 22, [("実録で測る", {"size": 14, "bold": True})])
txt(513, CY + 40, 360, 60,
    [[("・通知性能（鳴るべきときに鳴るか）", {"size": 12, "color": SUB})],
     [("・共通110本を全条件に通して比べる", {"size": 12, "color": SUB})]], line=1.45)

# ---- 締め ----
rect(70, 448, 6, 46, GOLD)
rect(76, 448, W - 146, 46, WHITE, line=LINE, lw=1.0)
txt(94, 456, W - 180, 32,
    [("「合成データのどの物理が、実環境の性能に必要か」に", {"size": 14, "color": SUB}),
     ("9月に直接答える", {"size": 14, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)

# ---- フッター ----
txt(58, 506, 90, 20, [("2026/08", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "今後・まとめ")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(10.5)
    r.font.bold = active
    r.font.color.rgb = WHITE if active else MUTED
    r.font.name = "Meiryo"
    meiryo(r)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("10", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"図_ablation実験計画_2026-08-23.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
