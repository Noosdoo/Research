# -*- coding: utf-8 -*-
"""カード上端の「色帯」を、左端の縦線へ移す。

本人の指示（2026-09-05）:
 「線で囲むだけは少し寂しい」→ 装飾は減らさず、形だけ変える。

根拠: このデッキは元から**左端の縦線**で箱の種類を示している。
 - 3枚目の下の枠   : 左端に金色の縦線（幅6.5pt）
 - 4枚目の3つの箱  : 左端に縦線（幅5pt・紫／金／灰）
上端の帯を足したのは後から作った枚だけで、そこだけ浮いていた。
帯を左へ移すと、デッキ全体が元の作法で揃う。

やること（この2つだけ。ほかは触らない）:
 1. 高さ 3.5pt・幅128pt以上の塗りつぶし長方形 ＝ 色帯 を見つける
 2. その帯を **幅5pt・カードの高さいっぱいの縦線** に作り替える（位置は左上のまま）
    カードが見つからない帯は、安全のためそのまま残す

見出しの色は変えない（4枚目と同じく、色は線が持ち、見出しは黒のまま）。

使い方:
  python scripts/_move_bars_to_left.py <入力.pptx> [出力.pptx]
出力を省略すると <入力>_左線.pptx を同じ場所に作る。原本は変更しない。
"""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

BAR_H_EMU = int(Pt(3.5))
MIN_W_PT = 128.0
NEW_W_EMU = int(Pt(5))          # 4枚目の縦線と同じ幅
TOL_EMU = int(Pt(6))            # カード照合の許容（手で数pt動かされていても拾う）


def fill_rgb(sh):
    try:
        if sh.fill.type == 1:
            return sh.fill.fore_color.rgb
    except Exception:
        pass
    return None


def find_card(slide, bar):
    """帯と左上・幅が一致し、背が高い図形＝カード本体を探す。"""
    best = None
    for sh in slide.shapes:
        if sh is bar or sh.height < int(Pt(20)):
            continue
        if abs(sh.left - bar.left) <= TOL_EMU and \
           abs(sh.top - bar.top) <= TOL_EMU and \
           abs(sh.width - bar.width) <= TOL_EMU:
            if best is None or sh.height > best.height:
                best = sh
    return best


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 2
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2]) if len(sys.argv) > 2 else \
        src.with_name(src.stem + "_左線" + src.suffix)
    shutil.copyfile(src, dst)

    prs = Presentation(str(dst))
    moved = kept = 0
    for i, slide in enumerate(prs.slides, 1):
        bars = [sh for sh in slide.shapes
                if sh.height == BAR_H_EMU
                and sh.width / 12700.0 >= MIN_W_PT
                and fill_rgb(sh) is not None]
        for bar in bars:
            card = find_card(slide, bar)
            rgb = fill_rgb(bar)
            if card is None:
                kept += 1
                print("p%-3d 帯 w=%6.1f  → カードが見つからず、そのまま残した"
                      % (i, bar.width / 12700.0))
                continue
            h_pt = card.height / 12700.0
            bar.width = NEW_W_EMU
            bar.height = card.height
            moved += 1
            print("p%-3d 帯 w=%6.1f → 左の縦線 幅5pt × 高さ%.0fpt  #%02X%02X%02X"
                  % (i, card.width / 12700.0, h_pt, *tuple(rgb)))
    prs.save(str(dst))
    print("\n帯 %d 本を左の縦線へ。%d 本はそのまま。" % (moved, kept))
    print("保存:", dst)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
