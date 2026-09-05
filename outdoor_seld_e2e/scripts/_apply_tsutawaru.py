# -*- coding: utf-8 -*-
"""伝わるデザインの2ルールを、デッキ全体に当てる。

参考: 伝わるデザイン（千葉大・高橋佑磨）https://tsutawarudesign.com/
  「囲い・枠が多くなると全体が煩雑になるので、濫用は避けましょう」
  「使う色は、背景や文字の色を含めて４色ぐらいに」
14枚目で試作して本人が了承（2026-09-05）。同じ扱いを全枚へ。

やること:
 1. 装飾の帯・縦線を削除（高さ3.5pt横帯 / 幅2〜9pt縦線）
    ただしカードが見つからないもの＝図の一部は絶対に消さない
 2. その帯が付いていた「カード」の塗りと枠線を消す（囲みをやめる）
 3. 代わりに、見出しの直下に細い横罫を1本引く（見出し＋罫線＋中身）
 4. 同じ枚で左右に並ぶカードの間に、細い縦罫を1本引く

色: 罫線は薄い灰（#DDDFE4）のみ。装飾に色を使わない。
 赤・金・緑は「危険度を表す文字・図」に残っているものだけが生き残る。

使い方: python scripts/_apply_tsutawaru.py <入力.pptx> [出力.pptx]
"""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

BAR_H = int(Pt(3.5))
MIN_W_PT = 128.0
TOL = int(Pt(6))
HAIR = RGBColor(0xDD, 0xDF, 0xE4)


def fill_rgb(sh):
    try:
        if sh.fill.type == 1:
            return sh.fill.fore_color.rgb
    except Exception:
        pass
    return None


def match_card(slide, bar, vertical):
    best = None
    for sh in slide.shapes:
        if sh is bar:
            continue
        if abs(sh.left - bar.left) > TOL or abs(sh.top - bar.top) > TOL:
            continue
        ok = (abs(sh.height - bar.height) <= TOL and sh.width > bar.width * 5) \
            if vertical else \
            (abs(sh.width - bar.width) <= TOL and sh.height >= int(Pt(20)))
        if ok and (best is None or sh.width * sh.height > best.width * best.height):
            best = sh
    return best


# E案で色を付けた囲みの塗り色（これだけを狙う。表の縞 #F2F3F5 や図 #ECEEF1 は除外）
E_TINTS = {(0xF5, 0xF4, 0xF7), (0xFC, 0xF8, 0xEF), (0xF6, 0xF7, 0xF7),
           (0xF0, 0xF5, 0xF3), (0xFA, 0xF1, 0xF0)}


def tinted_cards(slide):
    out = []
    for sh in slide.shapes:
        if sh.width / 12700.0 < 100 or sh.height / 12700.0 < 30:
            continue
        rgb = fill_rgb(sh)
        if rgb is not None and tuple(rgb) in E_TINTS:
            out.append(sh)
    return out


def pairs_of(slide):
    out = []
    for bar in list(slide.shapes):
        w, h = bar.width / 12700.0, bar.height / 12700.0
        vert = 2.0 <= w <= 9.0 and h >= 20.0
        horz = bar.height == BAR_H and w >= MIN_W_PT
        if not (vert or horz):
            continue
        rgb = fill_rgb(bar)
        if rgb is None:
            continue
        card = match_card(slide, bar, vert)
        if card is not None:                       # 図の一部は対象外
            out.append((bar, card))
    return out


def heading_of(slide, card):
    """カードの上のほうにある最初のテキスト枠＝見出し。"""
    best, by = None, None
    cx, cy = card.left, card.top
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if cx - Pt(6) <= sh.left <= cx + Pt(60) and cy < sh.top <= cy + Pt(46):
            if by is None or sh.top < by:
                best, by = sh, sh.top
    return best


def rule(slide, x, y, w, h):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = HAIR
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def main() -> int:
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2]) if len(sys.argv) > 2 else \
        src.with_name(src.stem + "_伝わる" + src.suffix)
    shutil.copyfile(src, dst)

    prs = Presentation(str(dst))
    n_box = n_rule = n_vrule = 0
    for i, slide in enumerate(prs.slides, 1):
        pairs = pairs_of(slide)
        cards = []
        for card in tinted_cards(slide):           # E案の色付き囲みを外す
            pairs.append((None, card))
        for bar, card in pairs:
            if bar is not None:
                bar._element.getparent().remove(bar._element)
            card.fill.background()                 # 囲みをやめる
            card.line.fill.background()
            n_box += 1
            cards.append(card)
            hd = heading_of(slide, card)
            if hd is not None:
                y = hd.top + hd.height + Pt(2)
                rule(slide, card.left, y, card.width, Pt(0.8))
                n_rule += 1
        # 左右に並ぶカードの間に縦罫
        cards.sort(key=lambda c: c.left)
        for a, b in zip(cards, cards[1:]):
            if abs(a.top - b.top) <= Pt(10) and b.left > a.left + a.width - Pt(2):
                gap = b.left - (a.left + a.width)
                if Pt(8) <= gap <= Pt(80):
                    rule(slide, a.left + a.width + gap / 2,
                         a.top + Pt(8), Pt(0.8), a.height - Pt(16))
                    n_vrule += 1
        if pairs:
            print("p%-3d 囲みを外した %d 件" % (i, len(pairs)))
    prs.save(str(dst))
    print("\n囲み %d 個を外し、見出し下の横罫 %d 本・仕切りの縦罫 %d 本を引いた"
          % (n_box, n_rule, n_vrule))
    print("保存:", dst)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
