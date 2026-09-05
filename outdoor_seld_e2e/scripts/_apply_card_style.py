# -*- coding: utf-8 -*-
"""カードの「種類の示し方」を、デッキ全体に一括で当てる。

経緯（2026-09-05）: 上端の色帯も左端の縦線も「箱に色の棒を貼る」同じ発想で、
本人から Claude の癖だと指摘された。棒を使わない方法を6案つくって比較し、
選ばれた1つをデッキ全体（16か所）に当てるのがこのスクリプト。

案（_make_card_variants.py と同じ定義）:
  A 枠線そのものを色にする
  B 見出し行だけ薄く色を敷く
  C 箱全体を薄いグレーで塗る（枠線なし）
  D 右と下の枠線を消す（上と左だけのかぎ形）        出典: パワポ研「枠線デザイン9選」⑦
  E 薄い背景＋同系色の少し濃い枠線                  出典: Webクリエイターボックス
  F 薄い色の背景だけ（枠線なし）
  DE Dの形 ＋ Eの面（薄い背景＋上と左のかぎ形）

対象の見つけ方（2026-09-05 追加: 縦線にも対応）:
 - 横の色帯 : 高さ 3.5pt・幅128pt以上。左上と幅が一致する背の高い図形がカード
 - 縦の色線 : 幅 2〜9pt・高さ20pt以上。左上と高さが一致する横長の図形がカード
   （3枚目の下の枠・4枚目の3つの箱がこれ。中間発表から引き継いだ形）
 帯は必ず消し、色はカードの塗り／枠線／かぎ形が引き継ぐ。
 道路の帯(h=6.0)・白線(h=3.0)・軌跡の線(h=2.4)は高さが違うので巻き込まない。

使い方:
  python scripts/_apply_card_style.py <入力.pptx> <案> [出力.pptx]
  例) python scripts/_apply_card_style.py deck.pptx E
出力を省略すると <入力>_<案>.pptx を同じ場所に作る。原本は変更しない。
"""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

BAR_H_EMU = int(Pt(3.5))
MIN_W_PT = 128.0
TOL = int(Pt(6))
GREY = RGBColor(0xF2, 0xF3, 0xF5)


def light(rgb, k):
    return RGBColor(*[int(255 - (255 - c) * k) for c in tuple(rgb)])


def fill_rgb(sh):
    try:
        if sh.fill.type == 1:
            return sh.fill.fore_color.rgb
    except Exception:
        pass
    return None


def _match(slide, bar, vertical):
    card = None
    for sh in slide.shapes:
        if sh is bar:
            continue
        if abs(sh.left - bar.left) > TOL or abs(sh.top - bar.top) > TOL:
            continue
        if vertical:                       # 縦線: 高さが一致し、ずっと横長
            ok = abs(sh.height - bar.height) <= TOL and sh.width > bar.width * 5
        else:                              # 横帯: 幅が一致し、ずっと背が高い
            ok = abs(sh.width - bar.width) <= TOL and sh.height >= int(Pt(20))
        if ok and (card is None or sh.width * sh.height >
                   (card.width * card.height if card else 0)):
            card = sh
    return card


def cards_of(slide):
    out = []
    # --- 縦の色線（3・4枚目の形）---
    for bar in list(slide.shapes):
        w, h = bar.width / 12700.0, bar.height / 12700.0
        if not (2.0 <= w <= 9.0 and h >= 20.0):
            continue
        rgb = fill_rgb(bar)
        if rgb is None:
            continue
        out.append((bar, _match(slide, bar, True), rgb))
    # --- 横の色帯 ---
    for bar in list(slide.shapes):
        if bar.height != BAR_H_EMU or bar.width / 12700.0 < MIN_W_PT:
            continue
        rgb = fill_rgb(bar)
        if rgb is None:
            continue
        card = None
        for sh in slide.shapes:
            if sh is bar or sh.height < int(Pt(20)):
                continue
            if abs(sh.left - bar.left) <= TOL and abs(sh.top - bar.top) <= TOL \
               and abs(sh.width - bar.width) <= TOL:
                if card is None or sh.height > card.height:
                    card = sh
        out.append((bar, card, rgb))
    return out


def drop(sh):
    sh._element.getparent().remove(sh._element)


def bracket(slide, card, rgb, w=Pt(1.6)):
    """上と左だけの「かぎ形」を描く。"""
    for x, y, ww, hh in [(card.left, card.top, card.width, w),
                         (card.left, card.top, w, card.height)]:
        ln = slide.shapes.add_shape(1, x, y, ww, hh)
        ln.fill.solid()
        ln.fill.fore_color.rgb = rgb
        ln.line.fill.background()
        ln.shadow.inherit = False


def apply(style, slide, card, rgb):
    if style == "A":
        card.line.color.rgb = rgb
        card.line.width = Pt(1.5)
    elif style == "B":
        band = slide.shapes.add_shape(1, card.left, card.top, card.width, Pt(42))
        band.fill.solid()
        band.fill.fore_color.rgb = light(rgb, 0.14)
        band.line.fill.background()
        band.shadow.inherit = False
        sp = band._element
        sp.getparent().remove(sp)
        card._element.addnext(sp)
    elif style == "C":
        card.fill.solid()
        card.fill.fore_color.rgb = GREY
        card.line.fill.background()
    elif style == "D":
        card.line.fill.background()
        bracket(slide, card, rgb)
    elif style == "E":
        card.fill.solid()
        card.fill.fore_color.rgb = light(rgb, 0.07)
        card.line.color.rgb = light(rgb, 0.30)
        card.line.width = Pt(1.0)
    elif style == "F":
        card.fill.solid()
        card.fill.fore_color.rgb = light(rgb, 0.10)
        card.line.fill.background()
    elif style == "DE":
        card.fill.solid()
        card.fill.fore_color.rgb = light(rgb, 0.07)
        card.line.fill.background()
        bracket(slide, card, rgb)
    else:
        raise SystemExit("知らない案: %s" % style)


def main() -> int:
    if len(sys.argv) < 3:
        print(__doc__)
        return 2
    src, style = Path(sys.argv[1]), sys.argv[2].upper()
    dst = Path(sys.argv[3]) if len(sys.argv) > 3 else \
        src.with_name(src.stem + "_" + style + src.suffix)
    shutil.copyfile(src, dst)

    prs = Presentation(str(dst))
    n = miss = 0
    for i, slide in enumerate(prs.slides, 1):
        for bar, card, rgb in cards_of(slide):
            if card is None:
                # カードが無い＝図の一部（軌跡の線・目盛りなど）。絶対に消さない。
                miss += 1
                print("p%-3d 図の一部とみなして残した w=%.1f h=%.1f #%02X%02X%02X"
                      % (i, bar.width / 12700.0, bar.height / 12700.0, *tuple(rgb)))
                continue
            apply(style, slide, card, rgb)
            drop(bar)
            n += 1
            print("p%-3d %6.0f x %-4.0f pt  #%02X%02X%02X"
                  % (i, card.width / 12700.0, card.height / 12700.0, *tuple(rgb)))
    prs.save(str(dst))
    print("\n案%s を %d か所に適用（帯は全て削除。%d か所はカード未検出）"
          % (style, n, miss))
    print("保存:", dst)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
