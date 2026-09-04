# -*- coding: utf-8 -*-
"""【修正スライド1枚だけ】収録の流れ — 「1テイクごと」の誤りを直した版。

2026-09-04の確認: 事前登録（実録スモーク計画書 収録規約2）は
「**セッション冒頭に毎回**」、記録紙テンプレートも「冒頭の儀式チェック（毎セッション）」。
儀式はもともとテイクごとではない。ハンドブック§4の表題「1テイクの標準手順」が
紛らわしかっただけで、規約は変えていない（同日ハンドブックの表題も訂正済み）。

直した3点:
 1. 「1テイクあたり6つの手順」→「1地点あたり6つの手順。④だけをテイクごとに繰り返す」
 2. ④のコマに「×n」を明記
 3. 儀式の見出しに「セッション冒頭に1回」を明記

出力: md/seminar/修正_収録の流れ_2026-09-15.pptx（1枚のみ・p14）
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF5, 0xF3, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])
SECTIONS = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


def panel(x, y, w, h, accent, title, tsize=14):
    rect(x, y, w, h, WHITE, line=LINE, lw=1.0)
    rect(x, y, w, 3.5, accent)
    txt(x + 16, y + 14, w - 32, 22, [(title, {"size": tsize, "bold": True})])
    return y + 40


rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 170, 40,
    [("収 録 の 流 れ", {"size": 22, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

txt(70, 100, 820, 20,
    [[("1地点あたり6つの手順。", {"size": 12, "color": MUTED}),
      ("④だけをテイクごとに繰り返します", {"size": 12, "bold": True}),
      ("（1通過＝1テイク）。", {"size": 12, "color": MUTED})]])

steps = [
    ("① 前日", "15分", "設定を家で全部固定する", PUR),
    ("② 到着", "2分", "定位置に立つ", PUR),
    ("③ 儀式", "90秒", "手拍子・騒音計・ベル4方位", GOLD),
    ("④ 本番", "20〜30秒 ×n", "無言。1通過＝1テイク", RED),
    ("⑤ 直後", "30秒", "5秒数えてから声で記録", GOLD),
    ("⑥ 確認", "5分", "その場で再生して品質確認", GREEN),
]
SW, SGAP = 128, 8
for i, (a, t, b, col) in enumerate(steps):
    x = 70 + i * (SW + SGAP)
    rect(x, 132, SW, 108, WHITE, line=LINE, lw=1.0)
    rect(x, 132, SW, 3.5, col)
    txt(x + 12, 146, SW - 24, 20, [(a, {"size": 12.5, "bold": True, "color": col})])
    txt(x + 12, 166, SW - 24, 18, [(t, {"size": 10.5, "color": MUTED})])
    txt(x + 12, 188, SW - 24, 46, [(b, {"size": 10.5, "color": SUB})], line=1.3)
    if i < len(steps) - 1:
        txt(x + SW - 2, 176, 14, 20, [("›", {"size": 15, "color": MUTED})])

# ④ が繰り返しであることを図で示す
rect(70 + 3 * (SW + SGAP), 244, SW, 22, PALE)
txt(70 + 3 * (SW + SGAP), 248, SW, 18,
    [("↻ 通るたびに繰り返す", {"size": 9.5, "color": PUR, "bold": True})],
    align=PP_ALIGN.CENTER)

y0 = panel(70, 278, 400, 122, GOLD, "③ 儀式の中身（90秒・セッション冒頭に1回）")
for i, (a, b) in enumerate([
    ("手拍子を1回", "録音と動画の時刻を合わせる"),
    ("騒音計と並べて無言で60秒", "何デシベルだったかの基準を取る"),
    ("ベルを前後左右から1打ずつ", "マイクの正面と体の正面のズレを測る"),
]):
    yy = y0 + i * 26
    n = rect(88, yy + 2, 16, 16, GOLD, shape=MSO_SHAPE.OVAL)
    label(n, str(i + 1), size=9.5)
    txt(112, yy, 350, 22,
        [[(a + "　", {"size": 10.5, "bold": True}),
          (b, {"size": 10, "color": SUB})]])

y0 = panel(490, 278, 400, 122, GREEN, "帰宅後、その日のうちに")
for i, t in enumerate([
    "録音を学習データと同じ形式に変換する",
    "1件ずつの10秒クリップに切り出す",
    "紙の記録を表に打ち込む（注釈）",
    "検査スクリプトで形式ミスを機械的に見つける",
]):
    rect(508, y0 + i * 22 + 5.5, 6, 6, GREEN)
    txt(524, y0 + i * 22, 350, 20, [(t, {"size": 10.5, "color": SUB})])

rect(70, 414, 820, 44, PALE)
rect(70, 414, 4, 44, PUR)
txt(88, 425, 790, 24,
    [[("この流れは9/1に手元で通し確認まで済ませてあります。", {"size": 11.5, "bold": True}),
      ("　機材が届いたその日から本番に入れます。", {"size": 11.5, "color": SUB})]])

txt(58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
ws = [16 + len(s) * 12.0 for s in SECTIONS]
x = (W - (sum(ws) + 8 * (len(SECTIONS) - 1))) / 2
for s, wd in zip(SECTIONS, ws):
    on = (s == "実録の計画")
    c = rect(x, 502, wd, 22, NAVY if on else CHIP)
    label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("14", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"修正_収録の流れ_2026-09-15.pptx")
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ PowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
