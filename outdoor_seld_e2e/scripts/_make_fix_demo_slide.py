# -*- coding: utf-8 -*-
"""【修正スライド1枚だけ】④ 触覚デモを動かした — 震え方をパルス図にした版。

本人の依頼（2026-09-05）: 「単発のパルス」が分かりにくい。スライドを分かりやすくしたい。

直した点:
 1. 「強＝4連打 / 中＝2発 / 警告音＝単発のパルス」という文字を、
    **実際の長さと間隔で描いたパルス図**に置き換えた。
    数値は out/joycon_demo_v2/unity/JoyconDemoPlayer.cs の SetRumble 実引数そのまま:
      強  : 320/640Hz 振幅1.0 110ms を 0.18秒間隔で4回
      中  :  80/160Hz 振幅0.4  90ms を 0.39秒間隔で2回
      警告: 120/240Hz 振幅0.5 300ms を 1回
    棒の長さ＝1発の長さ、棒の高さ＝振幅、間隔＝実際の間隔。同じ物差しで並べた。
 2. 「作った3本の歩道シナリオ（2.0/2.8/3.5m）」を削除。
    あれは v1（out/joycon_demo/）の内容で、いま動いている v2 には無い。
    代わりに v2 の実態（自作28場面＋評価用4本・すべてモデル出力）を書いた。

出力: md/seminar/修正_触覚デモ_2026-09-15.pptx（1枚のみ・p10）
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF5, 0xF3, 0xF9)
ROAD = RGBColor(0xE9, 0xE9, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])
SECTIONS = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


def panel(x, y, w, h, accent, title, sub=""):
    rect(x, y, w, h, WHITE, line=LINE, lw=1.0)
    rect(x, y, w, 3.5, accent)
    runs = [(title, {"size": 13.5, "bold": True})]
    if sub:
        runs.append(("　" + sub, {"size": 10, "color": MUTED}))
    txt(x + 16, y + 14, w - 32, 22, [runs])
    return y + 44


rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 170, 40,
    [("④　触 覚 デ モ を 動 か し た", {"size": 22, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

txt(70, 100, 820, 20,
    [("助言どおり Joy-con と Unity で作り、9/2に実機で振動しました。",
      {"size": 12, "color": MUTED})])

# ---------------- 左：実機で振動した ----------------
y = panel(70, 128, 400, 194, GREEN, "実機で振動した", "9/2")
for a, b in [("使ったもの",
              "SwitchのJoy-con と Unity 6.6\nハンダごて不要・部品の購入ゼロ"),
             ("鳴らす中身",
              "自作の場面28本＋評価用データ4本\n通知はすべて本物のモデル出力"),
             ("動かし方",
              "←→で場面を切り替え、Space で再生\n画面には俯瞰の街と判定が出る")]:
    rect(88, y + 6, 6, 6, GREEN)
    txt(104, y - 1, 350, 18, [(a, {"size": 11.5, "bold": True})])
    txt(104, y + 17, 350, 34,
        [[(t, {"size": 9.5, "color": SUB})] for t in b.split("\n")], line=1.3)
    y += 52

# ---------------- 右：3つの震え方（パルス図） ----------------
PX, PY, PW, PH = 490, 128, 400, 194
py = panel(PX, PY, PW, PH, PUR, "3つの震え方", "実際の長さと間隔")

BX0 = PX + 118          # 波形の左端
SEC = 340.0             # 1秒あたりの長さ(pt)
patterns = [
    ("強", "至近警告", RED, 4, 0.110, 0.18, 1.00, "短く強く4回"),
    ("中", "注意", GOLD, 2, 0.090, 0.39, 0.40, "弱く2回"),
    ("警告", "サイレン等", GREEN, 1, 0.300, 0.00, 0.50, "長めに1回"),
]
for i, (nm, kind, col, n, dur, gap, amp, note) in enumerate(patterns):
    base = py + 30 + i * 50            # 波形の底
    rect(BX0, base, 262, 0.8, HAIR)
    txt(PX + 16, base - 26, 96, 18, [(nm, {"size": 12.5, "bold": True, "color": col})])
    txt(PX + 16, base - 9, 96, 16, [(kind, {"size": 9, "color": MUTED})])
    for k in range(n):
        x = BX0 + k * gap * SEC
        w = dur * SEC
        h = 4 + 20 * amp
        rect(x, base - h, w, h, col)
    txt(BX0 + 190, base - 22, 76, 16,
        [(note, {"size": 9, "color": MUTED})], align=PP_ALIGN.RIGHT)

txt(PX + 16, PY + PH - 30, PW - 32, 20,
    [[("回数で危険度を、長さで種類を", {"size": 10.5, "bold": True}),
      ("表しています。", {"size": 10.5, "color": SUB})]])

# ---------------- 下：首元デバイス ----------------
DY = 334
y = panel(70, DY, 820, 124, PUR, "この先に作る首元デバイス", "構想・未実装")
CX, CY = 232, DY + 68
rect(CX - 74, CY - 30, 148, 60, None, line=NAVY, lw=7, shape=MSO_SHAPE.OVAL)
rect(CX - 36, CY - 16, 72, 32, ROAD, shape=MSO_SHAPE.OVAL)
txt(CX - 36, CY - 8, 72, 16, [("首", {"size": 10, "color": SUB})],
    align=PP_ALIGN.CENTER)
for dx, dy, on in [(-70, 8, 0), (-46, 26, 0), (0, 32, 0), (46, 26, 1), (70, 8, 0)]:
    rect(CX + dx - 6, CY + dy - 6, 12, 12,
         GOLD if on else CHIP, line=RED if on else MUTED,
         lw=1.6 if on else 0.8, shape=MSO_SHAPE.OVAL)
txt(CX - 90, CY - 52, 180, 16,
    [("歩く向き（前）", {"size": 9.5, "color": MUTED})], align=PP_ALIGN.CENTER)
txt(CX - 110, CY + 40, 220, 16,
    [("振動子5個＋4chマイクを同じバンドに", {"size": 9.5, "color": SUB})],
    align=PP_ALIGN.CENTER)
txt(400, DY + 52, 480, 56,
    [[("鳴っている位置で方向を、鳴り方で危険度を伝えます。",
       {"size": 12, "bold": True})],
     [("マイクと振動子を同じバンドに載せるので、「マイクから見た右」と"
       "「体から見た右」が一致します。", {"size": 10.5, "color": SUB})]], line=1.45)

# ---------------- フッター ----------------
txt(58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
ws = [16 + len(s) * 12.0 for s in SECTIONS]
x = (W - (sum(ws) + 8 * (len(SECTIONS) - 1))) / 2
for s, wd in zip(SECTIONS, ws):
    on = (s == "振り返り")
    c = rect(x, 502, wd, 22, NAVY if on else CHIP)
    label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("10", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"修正_触覚デモ_2026-09-15.pptx")
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ PowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
