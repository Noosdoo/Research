# -*- coding: utf-8 -*-
"""11枚目「④ 同じ場面で、鳴る車と鳴らない車」のタイムラインを実データに直す。

背景（2026-09-06）:
  この図はオラクル版（9_参考_正解の位置から_モデル不使用）の数字で描かれていた。
  スライド本文は「実機で振動した」＝モデル版の話なので、図と中身が食い違っていた。
  デモの正本は out/joycon_demo_v2/場面/1_路地・住宅街/前から対向車と、遠くを追い越す車。

  実データ（_cues.csv / _scene.csv より）:
    車①(obj0) 正面から来る    最接近 t=5.0s / 1.20m / 約34km/h
               通知 4.2s 強 → 5.0s 中 → 5.1s 中（束ねモードONで1回に）
               リード = 5.0 - 4.2 = 0.8秒
    車②(obj1) 後ろから追い越す 最接近 t=8.5s / 4.50m / 約46km/h  通知ゼロ＝抑制
  （オラクル版は 3.0s 強 → リード2.0秒。旧図はこちらの数字だった）

直すもの:
  - 車①: 強を先・中を後に（実際の順）。帯の幅も4連打/2発の実長に合わせる
  - 最接近 6.0s/1.0m → 5.0s/1.20m、リード 2.0秒 → 0.8秒
  - 速度 30/40km/h → 約34/約46km/h（_scene.csv の実測中央値）
  - 「通知の元＝本物の検出層の出力」の行を追加（オラクルと取り違えられないように）
  - 実体のない6.2sの金の細片を削除
  車②の最接近(8.5s)・4.5m・抑制は元から正しいので据え置き。

本人のデッキから11枚目を複製して直すので、書式は元のまま変わらない。
出力: md/seminar/修正_鳴る車と鳴らない車_2026-09-15.pptx（1枚のみ・p11）
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = ROOT / "md/seminar/修正_鳴る車と鳴らない車_2026-09-15.pptx"

# タイムラインの目盛り: 0s=257.6pt, 10s=913.6pt → 1秒 = 65.6pt
T0, PPS = 257.6, 65.6
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 実データ
CUE_STRONG, CUE_MID = 4.2, 5.0          # 強の発火 / 中の発火（5.1の再発火は束ねる）
DUR_STRONG, DUR_MID = 0.7, 0.7          # 4連打・2発を束ねた長さ（伝え方の定義どおり）
CPA1, DIST1 = 5.0, 1.20
CAR1_EXIST = (0.2, 9.9)


def X(t):
    return T0 + PPS * t


def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(
            copy.deepcopy(shp._element), "p:extLst")
    return dst


def put(sp, left=None, top=None, width=None, height=None):
    if left is not None:
        sp.left = int(Pt(left))
    if top is not None:
        sp.top = int(Pt(top))
    if width is not None:
        sp.width = int(Pt(width))
    if height is not None:
        sp.height = int(Pt(height))


def settext(sp, lines, size=None, color=None, bold=None, align=None):
    """段落ごとに本文を差し替える。1段落目の書式を雛形として引き継ぐ。"""
    tf = sp.text_frame
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, line in enumerate(lines):
        p = tf.paragraphs[i]
        if not p.runs:
            src = tf.paragraphs[0].runs
            r = p.add_run()
            if src:
                r.font.size = src[0].font.size
                r.font.bold = src[0].font.bold
                r.font.name = src[0].font.name
                try:
                    r.font.color.rgb = src[0].font.color.rgb
                except Exception:
                    pass
        p.runs[0].text = line
        for r in list(p.runs[1:]):
            r._r.getparent().remove(r._r)
        if size is not None:
            p.runs[0].font.size = Pt(size)
        if color is not None:
            p.runs[0].font.color.rgb = color
        if bold is not None:
            p.runs[0].font.bold = bold
        if align is not None:
            p.alignment = align
    for p in list(tf.paragraphs[len(lines):]):
        p._p.getparent().remove(p._p)


def fix(sl):
    s = list(sl.shapes)

    # 取り違え防止: 文言と目盛りで本人確認してから触る
    assert "同じ場面" in s[7].text_frame.text, "11枚目ではない"
    assert s[19].text_frame.text.strip() == "0s", "目盛りの並びが想定と違う"
    assert s[29].text_frame.text.strip() == "10s", "目盛りの並びが想定と違う"
    assert "30km/h" in s[31].text_frame.text, "既に修正済みか、図形の並びが違う"

    # 複製すると日付・ページ番号のプレースホルダが空レイアウトの位置に戻るので置き直す
    put(s[5], left=66.0, top=500.5, width=216.0, height=28.8)    # 2026/09/15
    put(s[6], left=678.0, top=500.5, width=216.0, height=28.8)   # 11

    # ---------- 使ったもの: 「通知の元」の行を足す ----------
    put(s[15], height=63)
    put(s[16], height=63)
    settext(s[15], ["使ったもの", "伝え方", "通知の元"])
    settext(s[16], ["SwitchのJoy-conとUnity 6.6",
                    "強＝4連打 / 中＝2発 / 警告音＝単発のパルス",
                    "本物の検出層の出力（正解の位置から作った版ではない）"])

    # ---------- 車① ----------
    settext(s[31], ["横1.2m・約34km/h"])

    a, b = X(CAR1_EXIST[0]), X(CUE_STRONG)
    put(s[32], left=a, width=b - a)                         # 通知前の灰色帯

    put(s[34], left=X(CUE_STRONG), width=PPS * DUR_STRONG)  # 赤＝強（先に出る）
    put(s[33], left=X(CUE_MID), width=PPS * DUR_MID)        # 金＝中（後に下がる）

    put(s[38], left=X(CUE_STRONG), top=234.9, width=PPS * DUR_STRONG)
    settext(s[38], ["強"], size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    put(s[37], left=X(CUE_MID), top=234.9, width=PPS * DUR_MID)
    settext(s[37], ["中"], size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    s[35]._element.getparent().remove(s[35]._element)       # 6.2sの金の細片は実体なし

    a = X(CUE_MID + DUR_MID)
    put(s[36], left=a, width=X(CAR1_EXIST[1]) - a)          # 通知後の灰色帯

    # 発火の三角（4.2 / 5.0 / 5.1。後ろ2つは束ねられる再発火）
    for sp, t in ((s[50], 4.2), (s[51], 5.0), (s[52], 5.1)):
        put(sp, left=X(t) - 5.0)
    put(s[53], left=X(CUE_STRONG) - 66.0)
    settext(s[53], ["4.2s に最初の振動"])

    # 最接近とリード
    put(s[39], left=X(CPA1) - 0.7)
    put(s[40], left=X(CPA1) - 60.0)
    settext(s[40], ["最接近 %.2fm" % DIST1])
    lead = CPA1 - CUE_STRONG
    put(s[54], left=X(CUE_STRONG), width=PPS * lead)
    put(s[55], left=X(CUE_STRONG) + PPS * lead / 2 - 71.4)
    settext(s[55], ["最接近の %.1f 秒前" % lead])

    # ---------- 車② ----------
    settext(s[42], ["横4.5m・約46km/h"])
    settext(s[48], ["最接近 4.50m"])                        # 8.5sの縦線は既に正しい

    # ---------- 下の説明 ----------
    settext(s[56], ["車①：正面から来る車",
                    "方位が変わらないまま近づくので、最接近の0.8秒前に強く鳴る。"
                    "通過後は中に下がる。"])
    settext(s[57], ["車②：後ろから追い越す車",
                    "横4.5mを通り過ぎるだけなので、最接近(8.5s)まで一度も鳴らない。"])
    return lead


def main() -> int:
    src = Presentation(str(SRC))
    out = Presentation()
    out.slide_width = src.slide_width
    out.slide_height = src.slide_height
    sl = clone(list(src.slides)[10], out)
    lead = fix(sl)

    dst = OUT
    try:
        out.save(str(dst))
    except PermissionError:
        dst = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
        out.save(str(dst))
    print("保存:", dst)
    print("車①: 強 %.1fs → 中 %.1fs / 最接近 %.1fs %.2fm / リード %.1f秒"
          % (CUE_STRONG, CUE_MID, CPA1, DIST1, lead))
    print("車②: 最接近 8.5s 4.50m / 通知なし")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
