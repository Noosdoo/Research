# -*- coding: utf-8 -*-
"""「どんな音を、どれだけ録るのか」の区分表を直す（2026-09-06 本人決定）。

決めたこと:
  1. **本数は1本も動かさない**（計210本・共通110本・D区分20本のまま）
  2. **D「固定・機会」→「警告音」に改名**。中身は不変
     元の名前は取り方（確実／運任せ）の軸で、A・E・Fのクラス名、B・Cの目的名と
     軸が混ざっていた。Dの中身は crossing / backup_beep / siren / horn の
     警告音4クラスなので、クラス名の軸に揃う
  3. **「サイレン等の遭遇待ち8」→「サイレンの遭遇待ち4 / クラクションの統制4」**
     合計8は不変。クラクションを機会枠から私有地の統制に移すことで、
     8クラスで唯一「実録0本もあり得る」状態だった horn の穴を塞ぐ
     （旧: md/design/実録_再設計の論点_2026-08-22.md「クラクションは配分ゼロ」）

本人のデッキから該当スライドを複製して直すので、書式は元のまま変わらない。
デッキの枚数が動くので、位置ではなく**見出しの文言**でスライドを探す。
出力: md/seminar/修正_どんな音を録るのか_2026-09-15.pptx（1枚のみ）
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = ROOT / "md/seminar/修正_どんな音を録るのか_2026-09-15.pptx"
TITLE = "どんな音を、どれだけ録るのか"

D_NAME = "D 警告音"
D_BODY = "実在の踏切8 / バック音の統制4 / サイレンの遭遇待ち4 / クラクションの統制4"


def find(prs, title):
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and title in sh.text_frame.text:
                return sl
    raise SystemExit("見出し「%s」のスライドが見つからない" % title)


def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(
            copy.deepcopy(shp._element), "p:extLst")
    return dst


def settext(sp, text):
    """1つ目のrunの書式を保ったまま本文を差し替える。"""
    p = sp.text_frame.paragraphs[0]
    p.runs[0].text = text
    for r in list(p.runs[1:]):
        r._r.getparent().remove(r._r)


def pick(sl, needle):
    """本文が needle で始まるテキスト枠を1つだけ取る（複数あれば止める）。"""
    hits = [s for s in sl.shapes
            if s.has_text_frame and s.text_frame.text.strip().startswith(needle)]
    assert len(hits) == 1, "「%s」が %d 個ある" % (needle, len(hits))
    return hits[0]


def fix(sl):
    # 取り違え防止: 本数が元のままであることを確かめてから触る
    assert pick(sl, "検出対象の8種類"), "見出し下の説明が違う"
    d_name = pick(sl, "D 固定・機会")
    d_body = pick(sl, "実在の踏切8")
    assert "サイレン等の遭遇待ち8" in d_body.text_frame.text, "D区分の中身が想定と違う"

    settext(d_name, D_NAME)
    settext(d_body, D_BODY)

    # 日付・ページ番号は複製で空レイアウトの位置に戻るので置き直す
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "2026/09/15":
            sh.left, sh.top = int(Pt(66.0)), int(Pt(500.5))
            sh.width, sh.height = int(Pt(216.0)), int(Pt(28.8))
        elif sh.has_text_frame and sh.text_frame.text.strip().isdigit() \
                and sh.top / 12700 > 480:
            sh.left, sh.top = int(Pt(678.0)), int(Pt(500.5))
            sh.width, sh.height = int(Pt(216.0)), int(Pt(28.8))


def main() -> int:
    src = Presentation(str(SRC))
    out = Presentation()
    out.slide_width = src.slide_width
    out.slide_height = src.slide_height
    sl = clone(find(src, TITLE), out)
    fix(sl)

    dst = OUT
    try:
        out.save(str(dst))
    except PermissionError:
        dst = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
        out.save(str(dst))
    print("保存:", dst)
    print("D:", D_NAME, "/", D_BODY)
    print("本数は不変（A20 B20 C20 D20 E20 F10 歩行対比100 ＝ 計210）")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
