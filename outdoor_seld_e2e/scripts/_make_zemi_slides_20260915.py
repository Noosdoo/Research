# -*- coding: utf-8 -*-
"""9/15 ゼミ（夏ゼミ反省・相談会）スライド 全10枚。

構成は 2026-09-03 に本人と合意した骨組み（→ md/seminar/作業記録_9-15準備_2026-09-02〜03.md §6 の改訂版）:
  1 今日お話しすること / 2 いただいた意見と、その後 / 3 約束したこと（すり抜けと対向）
  4 その結果（案A=型別の表） / 5 触覚デモを動かした
  6 どんな音を録るのか / 7 どこで・どんな条件で録るのか / 8 収録の流れ
  9 なぜ回数が読めないのか / 10 お願い：マイクは購入で

本人の指示で外したもの: 規程照会の議題（口頭で1行言う）・騒音計と提出期限の議題
（先生から話が出る見込み）・日程スライド（9月収録は既知）・想定問答・「ほかに進めたこと」
（10月のゼミへ回す）。

文字サイズ: 18pt以上の規定は中間発表・卒論発表など投影する本番だけ。この資料は対象外
（2026-09-03本人訂正）。既存の図と同じ10〜15pt級で情報密度を優先する。

材料の正本:
  md/seminar/中間発表_質疑と宿題_2026-08-30.md      （質疑5件・Slack3件）
  out/notify_v42_sweep2/q2_anzen/q2_table.md        （4枚目の型別表）
  md/design/通知v4.2_選定手順の事前宣言_2026-08-30.md §6（3枚目の仕組み）
  md/design/実録ハンドブック_2026-08-13.md §4・§6   （6〜9枚目の内訳と流れ）
  md/design/ゼミ相談メモ_機材購入_2026-09-15用.md   （10枚目）

出力: md/seminar/ゼミ相談スライド_2026-09-15.pptx
旧5枚版 md/seminar/相談会スライド_機材購入_2026-09-15.pptx は上書きせず残す。
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
PALE = RGBColor(0xF5, 0xF3, 0xF9)
PALEG = RGBColor(0xFD, 0xF7, 0xE8)
BAND = RGBColor(0xF2, 0xF3, 0xF5)
ROAD = RGBColor(0xE9, 0xE9, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0

SECTIONS = ["振り返り", "実録の計画", "お願い"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(sl, x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(sl, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


def chip(sl, x, y, text, color, size=11.5, pad=16, h=22):
    c = rect(sl, x, y, pad + len(text) * 12.0, h, None, line=color, lw=1.2)
    label(c, text, size=size, color=color)
    return c


def bullet(sl, x, y, w, runs, size=11, mark=INK, line=1.34, h=40):
    rect(sl, x, y + 5.5, 6, 6, mark)
    return txt(sl, x + 16, y, w - 16, h, runs, size=size, line=line)


def panel(sl, x, y, w, h, accent=None, title=None, tsize=14):
    rect(sl, x, y, w, h, WHITE, line=LINE, lw=1.0)
    if accent is not None:
        rect(sl, x, y, w, 3.5, accent)
    if title:
        txt(sl, x + 16, y + 14, w - 32, 22, [(title, {"size": tsize, "bold": True})])
    return y + (40 if title else 12)


def new_slide(title, page, active=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 44, 14, 1.4, H - 28, PUR)
    rect(sl, 58, 46, 30, 6, PUR)
    txt(sl, 98, 30, W - 170, 40,
        [(title, {"size": 22, "bold": True, "spc": 250})])
    rect(sl, 44, 90, W - 90, 1.2, INK)
    txt(sl, 58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
    widths = [16 + len(s) * 12.0 for s in SECTIONS]
    x = (W - (sum(widths) + 8 * (len(SECTIONS) - 1))) / 2
    for s, wd in zip(SECTIONS, widths):
        on = (s == active)
        c = rect(sl, x, 502, wd, 22, NAVY if on else CHIP)
        label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
        x += wd + 8
    txt(sl, W - 90, 506, 40, 20,
        [(str(page), {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)
    return sl


# ==================================================== P1 今日お話しすること
sl = new_slide("9/15 ゼミ — 夏ゼミの振り返りと、お願いが1件", 1)
txt(sl, 70, 100, 820, 20,
    [("夏ゼミでいただいた意見にどう応えたかをご報告し、最後に機材のお願いを1件させてください。",
      {"size": 12, "color": MUTED})])

y0 = panel(sl, 70, 130, 400, 190, PUR, "1. 振り返り")
for i, (h_, b) in enumerate([
    ("いただいた意見と、その後", "質疑5件・Slackコメント3件のすべてに対応しました"),
    ("約束したことを果たしました", "すり抜けと対向を方位の変化率で区別する（2〜4枚目）"),
    ("触覚デモを動かしました", "Joy-conで実際に振動させました（5枚目）"),
]):
    txt(sl, 88, y0 + i * 46, 366, 40,
        [[(h_, {"size": 12.5, "bold": True})],
         [(b, {"size": 10.5, "color": SUB})]], line=1.3)

y0 = panel(sl, 490, 130, 400, 190, GOLD, "2. お願い（1件だけ）")
txt(sl, 508, y0, 366, 46,
    [[("実環境の評価で使う4方向マイクを、", {"size": 12.5})],
     [("レンタルではなく購入", {"size": 14, "bold": True, "color": RED}),
      ("でお願いしたい", {"size": 12.5})]], line=1.35)
txt(sl, 508, y0 + 54, 366, 22,
    [[("Zoom H3-VR ほか　合計 ", {"size": 12}),
      ("約3〜4万円", {"size": 15, "bold": True, "color": RED})]])
txt(sl, 508, y0 + 84, 366, 40,
    [("9月末〜11月に屋外で計210回の録音を行います。その理由を6〜9枚目でご説明します。",
      {"size": 10.5, "color": SUB})], line=1.35)

rect(sl, 70, 336, 820, 40, BAND)
txt(sl, 86, 346, 790, 24,
    [[("承認をいただきたいのは2番だけです。", {"size": 12, "bold": True}),
      ("1番はご報告です。", {"size": 12, "color": SUB})]])

rect(sl, 70, 392, 820, 66, PALE)
rect(sl, 70, 392, 4, 66, PUR)
txt(sl, 92, 402, 780, 48,
    [[("研究の目的（再掲）", {"size": 11, "bold": True, "color": PUR})],
     [("屋外の危険な音を8種類検出し、種類・方向・距離から危険度を3段階に分けて、"
       "首元の振動で難聴の方に伝えるシステムの構築。行動の判断は本人に委ねます。",
       {"size": 11.5, "color": SUB})]], line=1.4)

# ==================================================== P2 いただいた意見と、その後
sl = new_slide("いただいた意見と、その後", 2, active="振り返り")
txt(sl, 70, 100, 820, 20,
    [("質疑5件とSlackコメント3件、すべてに手を動かしました。", {"size": 12, "color": MUTED})])

hdr = 128
rect(sl, 70, hdr, 820, 26, NAVY)
for cx, cw, t in [(84, 300, "いただいた意見"), (396, 380, "やったこと"), (788, 90, "状態")]:
    txt(sl, cx, hdr + 5, cw, 18, [(t, {"size": 11, "bold": True, "color": WHITE})])

rows = [
    ("Q 方向まで分かるのに、なぜ方向を教えないのか",
     "伝える設計です。付録でご説明しました", "説明済", GREEN),
    ("Q すり抜けと対向が同じ警告になり区別できないのでは",
     "方位の変化率で区別する仕組みを実装（3〜4枚目）", "完了", RED),
    ("Q 学習データでの評価は。どういう結果か",
     "全数値を記録と照合。複数車を低く言っていたと判明", "訂正済", GREEN),
    ("Q デバイス試作は修士を見越した想定か",
     "はい。振動デバイスを研究の範囲に入れました", "決定", GREEN),
    ("Q 未踏事業に挑戦してもよいのでは",
     "応募時期を調査（2027年度は11月ごろエントリー）", "調査済", GREEN),
    ("S 全盲・弱視向けの振動デバイスにノウハウがある",
     "先行研究を調査。首元は4〜6方向が現実的と判明", "調査済", GREEN),
    ("S 研究上はとりあえず全方位を録ってみては",
     "記録紙に前方の欄を追加。前方も記録します", "反映済", GREEN),
    ("S Joy-conとUnityで簡単に作ってみては",
     "作りました。9/2に実機で振動成功（5枚目）", "完了", RED),
]
ry = hdr + 26
for i, (q, a, st, col) in enumerate(rows):
    rect(sl, 70, ry, 820, 32, WHITE if i % 2 == 0 else BAND, line=LINE, lw=0.6)
    txt(sl, 84, ry + 7, 306, 20, [(q, {"size": 10.5})])
    txt(sl, 396, ry + 7, 386, 20, [(a, {"size": 10.5, "color": SUB})])
    b = rect(sl, 788, ry + 6, 62, 20, None, line=col, lw=1.1)
    label(b, st, size=10, color=col)
    ry += 32

rect(sl, 70, ry + 12, 820, 34, PALEG)
rect(sl, 70, ry + 12, 4, 34, GOLD)
txt(sl, 88, ry + 20, 790, 20,
    [[("赤枠の2件", {"size": 11.5, "bold": True, "color": RED}),
      ("＝この場でご説明したいものです。", {"size": 11.5, "color": SUB})]])

# ==================================================== P3 約束したこと
sl = new_slide("約束したこと：すり抜けと対向を、どう区別するか", 3, active="振り返り")
txt(sl, 70, 100, 820, 20,
    [("夏ゼミで「三角関数で区別する」とお答えした件です。", {"size": 12, "color": MUTED})])

# --- 左: 図 ---
FX, FY, FW, FH = 70, 128, 430, 250
rect(sl, FX, FY, FW, FH, WHITE, line=LINE, lw=1.0)
txt(sl, FX + 16, FY + 12, 300, 20,
    [("装着者から見た「方向」の変わり方", {"size": 12.5, "bold": True})])

OX, OY = FX + 215, FY + 210          # 観測者の位置
rect(sl, OX - 8, OY - 8, 16, 16, INK, shape=MSO_SHAPE.OVAL)
txt(sl, OX - 70, OY + 12, 140, 18,
    [("装着者", {"size": 10.5, "color": SUB})], align=PP_ALIGN.CENTER)

# 対向型：正面からまっすぐ
rect(sl, OX - 1.2, FY + 58, 2.4, 138, RED)
rect(sl, OX - 7, FY + 52, 14, 14, RED, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
txt(sl, OX - 108, FY + 40, 100, 34,
    [[("対向型", {"size": 11.5, "bold": True, "color": RED})],
     [("方位が変わらない", {"size": 9.5, "color": SUB})]],
    align=PP_ALIGN.RIGHT, line=1.25)

# すり抜け型：横を通過
rect(sl, FX + 40, FY + 118, 350, 2.4, GREEN)
rect(sl, FX + 356, FY + 112, 14, 14, GREEN, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
txt(sl, FX + 150, FY + 78, 176, 34,
    [[("すり抜け型", {"size": 11.5, "bold": True, "color": GREEN})],
     [("方位が速く変わる", {"size": 9.5, "color": SUB})]], line=1.25)
for dx in (-150, -60, 60, 150):
    rect(sl, OX, OY, 1.0, 1.0, GREEN)
    ln = sl.shapes.add_connector(1, Pt(OX), Pt(OY),
                                 Pt(OX + dx), Pt(FY + 119))
    ln.line.color.rgb = GREEN
    ln.line.width = Pt(0.8)
    ln.line.dash_style = 4

txt(sl, FX + 16, FY + 222, FW - 32, 20,
    [("判定に使うのは、最接近の2.5〜1.5秒前の区間", {"size": 10, "color": MUTED})])

# --- 右: 仕組み ---
y0 = panel(sl, 520, 128, 370, 250, PUR, "救済ルートを1本足した")
txt(sl, 538, y0, 334, 34,
    [("次の3つが4フレーム続いたら、強い警告の候補として救い上げます。",
      {"size": 11, "color": SUB})], line=1.35)
conds = [
    ("方位がほとんど変わらない", "|dθ/dt| ≤ 0.10 rad/s"),
    ("距離が縮んでいる", "頑健な傾きがマイナス"),
    ("15m以内にいる", "距離ゲート"),
]
for i, (a, b) in enumerate(conds):
    yy = y0 + 44 + i * 40
    n = rect(sl, 538, yy + 2, 20, 20, PUR, shape=MSO_SHAPE.OVAL)
    label(n, str(i + 1), size=11)
    txt(sl, 566, yy, 310, 34,
        [[(a, {"size": 11.5, "bold": True})],
         [(b, {"size": 9.5, "color": MUTED})]], line=1.25)

rect(sl, 538, y0 + 166, 334, 48, PALEG)
txt(sl, 550, y0 + 174, 312, 36,
    [("船の衝突回避と同じ原則です。方位が変わらないまま近づいてくる相手は、"
      "衝突コースにいます。", {"size": 10.5, "color": SUB})], line=1.35)

rect(sl, 70, 392, 820, 66, PALE)
rect(sl, 70, 392, 4, 66, PUR)
txt(sl, 92, 400, 780, 52,
    [[("なぜ15mで切るのか", {"size": 11.5, "bold": True, "color": PUR})],
     [("方位の変化率は距離の2乗に反比例します。遠くをまっすぐ通り過ぎる安全な車も"
       "「方位が変わらない」ように見えるので、距離で切らないと遠い車まで強い警告になります。",
       {"size": 11.5, "color": SUB})]], line=1.4)

# ==================================================== P4 その結果
sl = new_slide("その結果", 4, active="振り返り")
txt(sl, 70, 100, 820, 20,
    [("車を「対向型」と「すり抜け型」に分けて、前の版と比べました。",
      {"size": 12, "color": MUTED})])

hdr = 130
cols = [(84, 250, "どんな車か"), (340, 110, "型"), (456, 90, "件数"),
        (552, 150, "前（v4.1）"), (708, 150, "今（v4.2）")]
rect(sl, 70, hdr, 820, 28, NAVY)
for cx, cw, t in cols:
    txt(sl, cx, hdr + 6, cw, 18, [(t, {"size": 11.5, "bold": True, "color": WHITE})])

trows = [
    ("危ない（1.5m以内）", "対向", "460", "強い警告 81%", "90%", RED, True),
    ("危ない（1.5m以内）", "すり抜け", "43", "強い警告 77%", "93%", RED, True),
    ("鳴らすべきでない（3.2m超）", "対向", "52", "黙れた 33%", "71%", GREEN, True),
    ("鳴らすべきでない（3.2m超）", "すり抜け", "501", "黙れた 62%", "64%", MUTED, False),
]
ry = hdr + 28
for i, (a, b, n, old, new, col, good) in enumerate(trows):
    rect(sl, 70, ry, 820, 44, WHITE if i % 2 == 0 else BAND, line=LINE, lw=0.6)
    txt(sl, 84, ry + 13, 250, 20, [(a, {"size": 11.5})])
    txt(sl, 340, ry + 13, 110, 20, [(b, {"size": 11.5, "bold": True, "color": col})])
    txt(sl, 456, ry + 13, 90, 20, [(n, {"size": 11, "color": MUTED})])
    txt(sl, 552, ry + 13, 150, 20, [(old, {"size": 11.5, "color": SUB})])
    txt(sl, 708, ry + 10, 150, 26,
        [(new, {"size": 17 if good else 14, "bold": True,
                "color": col if good else MUTED})])
    ry += 44

rect(sl, 70, 322, 400, 62, PALEG)
rect(sl, 70, 322, 4, 62, GOLD)
txt(sl, 88, 330, 370, 48,
    [[("読み方", {"size": 11.5, "bold": True})],
     [("危ない車には強く鳴るようになり、鳴らすべきでない対向車では黙るようになりました。",
       {"size": 11, "color": SUB})]], line=1.35)

rect(sl, 490, 322, 400, 62, BAND)
rect(sl, 490, 322, 4, 62, MUTED)
txt(sl, 508, 330, 370, 48,
    [[("残っている弱点（正直に）", {"size": 11.5, "bold": True})],
     [("すり抜けていく安全な車での無駄鳴りは、62→64%とほぼ改善していません。",
       {"size": 11, "color": SUB})]], line=1.35)

txt(sl, 70, 398, 820, 56,
    [[("※ 合成データ・調整用の新しい検証セット（1,800本）での値です。",
       {"size": 10.5, "color": MUTED})],
     [("※ 卒論に載せる確定評価とは別枠です。また、中間発表で述べた数値とは"
       "測り方が違うため、比較できません。", {"size": 10.5, "color": MUTED})],
     [("※ 何を試すか・どう選ぶか・何をもって採用とするかは、結果を見る前に文章で固定しました。",
       {"size": 10.5, "color": PUR, "bold": True})]], line=1.5)

# ==================================================== P5 触覚デモ
sl = new_slide("触覚デモを動かしました", 5, active="振り返り")
txt(sl, 70, 100, 820, 20,
    [("「Joy-conとUnityで簡単に作ってみては」というSlackの助言をそのまま実行しました。",
      {"size": 12, "color": MUTED})])

y0 = panel(sl, 70, 130, 400, 176, GREEN, "9/2 実機で振動しました")
for i, runs in enumerate([
    [("使ったもの：", {"bold": True}),
     ("Joy-con と Unity 6.6。ハンダごて不要・部品の購入ゼロ", {"color": SUB})],
    [("鳴らす中身：", {"bold": True}),
     ("本物のモデル出力5本と、自作した歩道のシナリオ3本", {"color": SUB})],
    [("伝え方：", {"bold": True}),
     ("強＝4連打／中＝2発／警告音＝単発のパルス", {"color": SUB})],
]):
    bullet(sl, 88, y0 + i * 42, 366, runs, size=10.5, mark=GREEN, h=38)

y0 = panel(sl, 490, 130, 400, 176, PUR, "作った3本の歩道シナリオ")
for i, (d, t, c) in enumerate([
    ("2.0m", "強く鳴る", RED), ("2.8m", "中くらい", GOLD), ("3.5m", "鳴らさない", GREEN)]):
    yy = y0 + i * 42
    b = rect(sl, 508, yy, 62, 26, c)
    label(b, d, size=12)
    txt(sl, 582, yy + 4, 290, 20, [(t, {"size": 12, "color": SUB})])
txt(sl, 508, y0 + 130, 366, 20,
    [("同じ場面で距離だけを変え、鳴り方の違いを体で確かめられます。",
      {"size": 10, "color": MUTED})])

# 首元デバイスの完成イメージ（簡略）
DX, DY, DW, DH = 70, 322, 820, 136
rect(sl, DX, DY, DW, DH, WHITE, line=LINE, lw=1.0)
txt(sl, DX + 16, DY + 12, 400, 20,
    [("この先に作る首元デバイス（構想）", {"size": 12.5, "bold": True})])
CX2, CY2 = DX + 190, DY + 84
rect(sl, CX2 - 62, CY2 - 40, 124, 80, None, line=NAVY, lw=9, shape=MSO_SHAPE.OVAL)
rect(sl, CX2 - 30, CY2 - 20, 60, 40, ROAD, shape=MSO_SHAPE.OVAL)
txt(sl, CX2 - 30, CY2 - 9, 60, 18, [("首", {"size": 10.5, "color": SUB})],
    align=PP_ALIGN.CENTER)
for dx, dy, on in [(-62, 0, 0), (-44, -30, 0), (0, -40, 0), (44, -30, 1), (62, 0, 0)]:
    rect(sl, CX2 + dx - 7, CY2 + dy - 7, 14, 14,
         GOLD if on else CHIP, line=RED if on else MUTED,
         lw=1.6 if on else 0.8, shape=MSO_SHAPE.OVAL)
txt(sl, CX2 - 100, CY2 - 66, 200, 18,
    [("歩く向き（前）", {"size": 9.5, "color": MUTED})], align=PP_ALIGN.CENTER)
txt(sl, CX2 - 110, CY2 + 46, 220, 18,
    [("振動子5個＋4chマイクを同じバンドに載せる", {"size": 10, "color": SUB})],
    align=PP_ALIGN.CENTER)
txt(sl, DX + 380, DY + 44, 420, 76,
    [[("鳴っている位置が、そのまま音の方向になります。", {"size": 12, "bold": True})],
     [("マイクと振動子を同じバンドに載せるので、「マイクから見た右」と"
       "「体から見た右」が一致します。頭の向きを測るセンサーが要りません。",
       {"size": 11, "color": SUB})]], line=1.45)

# ==================================================== P6 どんな音を録るのか
sl = new_slide("どんな音を、どれだけ録るのか", 6, active="実録の計画")
txt(sl, 70, 100, 820, 20,
    [("検出対象の8種類を、屋外で計210回ぶん録ります。内訳は録る前に固定してあります。",
      {"size": 12, "color": MUTED})])

x = 70.0
for c in ["サイレン", "クラクション", "バック音", "自転車ベル", "車",
          "踏切・列車", "キックボード", "バイク"]:
    wd = 16 + len(c) * 12.0
    b = rect(sl, x, 128, wd, 26, PALE, line=PUR, lw=1.0)
    label(b, c, size=11, color=PUR)
    x += wd + 8

hdr = 170
rect(sl, 70, hdr, 820, 26, NAVY)
for cx, cw, t in [(84, 130, "区分"), (222, 70, "本数"), (300, 550, "中身")]:
    txt(sl, cx, hdr + 5, cw, 18, [(t, {"size": 11, "bold": True, "color": WHITE})])

grid = [
    ("A 走行車", "20", "横2.0〜3.2m を10本（うち後方から6本以上）／横5〜15m を10本"),
    ("B 弱点", "20", "EV・HVの接近5／停車→発進5／自転車の至近通過5／見通し不良の角5"),
    ("C 負例", "20", "静穏6／繁華街・工事6／雨上がり4／風あり4　＝鳴ってはいけない場面"),
    ("D 固定・機会", "20", "実在の踏切8／バック音の統制4／サイレン等の遭遇待ち8"),
    ("E キックボード", "20", "横1.0/1.2/1.5m × 速度2段 × 左右＝12／後方4／前方4"),
    ("F バイク", "10", "交通量の多い道での待ち受け（自分では走らせない）"),
    ("歩行対比", "100", "5条件×10反復を、止まって録る／歩いて録るの対で比較"),
]
ry = hdr + 26
for i, (a, n, b) in enumerate(grid):
    rect(sl, 70, ry, 820, 30, WHITE if i % 2 == 0 else BAND, line=LINE, lw=0.6)
    txt(sl, 84, ry + 6, 130, 20, [(a, {"size": 11, "bold": True})])
    txt(sl, 222, ry + 6, 70, 20, [(n, {"size": 12, "bold": True, "color": RED})])
    txt(sl, 300, ry + 6, 550, 20, [(b, {"size": 10.5, "color": SUB})])
    ry += 30

rect(sl, 70, ry + 10, 820, 36, PALEG)
rect(sl, 70, ry + 10, 4, 36, GOLD)
txt(sl, 88, ry + 18, 790, 22,
    [[("これに加えて、鳴ってはいけない場面の連続録音を100分以上。", {"size": 11.5}),
      ("「1時間あたり何回の誤警告か」を出すために必要です。",
       {"size": 11.5, "color": SUB})]])

# ==================================================== P7 どこで・どんな条件で
sl = new_slide("どこで、どんな条件で録るのか", 7, active="実録の計画")
txt(sl, 70, 100, 820, 20,
    [("場所を変えないと録れない音があり、天候の条件も決めてあります。",
      {"size": 12, "color": MUTED})])

y0 = panel(sl, 70, 128, 400, 168, PUR, "6か所を回ります")
places = ["静かな道（交通量の少ない住宅路）", "交通量の多い道（車・バイク）",
          "繁華街（雑踏＝鳴ってはいけない場面）", "工事現場の近く（高い騒音）",
          "見通しの悪い角（塀のある交差点）", "実在の踏切（時刻表どおり＝確実）"]
for i, p in enumerate(places):
    bullet(sl, 88, y0 + i * 21, 366, [(p, {"color": SUB})], size=10.5,
           mark=PUR, h=20, line=1.2)

y0 = panel(sl, 490, 128, 400, 168, GOLD, "録るときの決めごと")
for i, runs in enumerate([
    [("中止の基準：", {"bold": True}),
     ("風速3m/sを超える日と雨の日は録らない", {"color": SUB})],
    [("装着：", {"bold": True}),
     ("全テイク、止まった状態で胸に装着して統一", {"color": SUB})],
    [("録音設定：", {"bold": True}),
     ("家で全部固定し、現地で押すのは録音ボタンだけ", {"color": SUB})],
    [("記録：", {"bold": True}),
     ("1テイクごとに車種・速度・横距離・騒音計の値を紙に書く", {"color": SUB})],
]):
    bullet(sl, 508, y0 + i * 32, 366, runs, size=10.5, mark=GOLD, h=30)

GX, GY, GW, GH = 210, 312, 540, 146
rect(sl, GX, GY, GW, GH, WHITE, line=LINE, lw=1.0)
txt(sl, GX + 16, GY + 10, 300, 20,
    [("車を記録するときの位置関係", {"size": 12.5, "bold": True})])
txt(sl, GX + 330, GY + 12, 194, 18,
    [("速度：徐行〜約30km/h", {"size": 10, "color": SUB})], align=PP_ALIGN.RIGHT)
rect(sl, GX + 22, GY + 38, GW - 44, 32, ROAD)
c = rect(sl, GX + 350, GY + 44, 52, 20, GOLD)
label(c, "車", size=11)
txt(sl, GX + 408, GY + 42, 40, 22, [("→", {"size": 17})])
rect(sl, GX + 168, GY + 70, 1.6, 34, RED)
rect(sl, GX + 160, GY + 104, 16, 16, INK, shape=MSO_SHAPE.OVAL)
txt(sl, GX + 182, GY + 76, 330, 20,
    [("横距離 2.0〜3.2m ／ 5〜15m", {"size": 12, "bold": True, "color": RED})])
txt(sl, GX + 96, GY + 122, 260, 18,
    [("装着者（胸に4方向マイク・静止）", {"size": 10, "color": SUB})],
    align=PP_ALIGN.CENTER)

# ==================================================== P8 収録の流れ
sl = new_slide("1回の収録の流れ", 8, active="実録の計画")
txt(sl, 70, 100, 820, 20,
    [("1テイクあたり6つの手順。この骨格を全210回で共通にします。",
      {"size": 12, "color": MUTED})])

steps = [
    ("① 前日", "15分", "設定を家で全部固定する", PUR),
    ("② 到着", "2分", "定位置に立つ", PUR),
    ("③ 儀式", "90秒", "手拍子・騒音計・ベル4方位", GOLD),
    ("④ 本番", "20〜30秒", "無言。1通過＝1テイク", RED),
    ("⑤ 直後", "30秒", "5秒数えてから声で記録", GOLD),
    ("⑥ 確認", "5分", "その場で再生して品質確認", GREEN),
]
SW, SGAP = 128, 8
for i, (a, t, b, col) in enumerate(steps):
    x = 70 + i * (SW + SGAP)
    rect(sl, x, 132, SW, 108, WHITE, line=LINE, lw=1.0)
    rect(sl, x, 132, SW, 3.5, col)
    txt(sl, x + 12, 146, SW - 24, 20, [(a, {"size": 12.5, "bold": True, "color": col})])
    txt(sl, x + 12, 166, SW - 24, 18, [(t, {"size": 10.5, "color": MUTED})])
    txt(sl, x + 12, 188, SW - 24, 46, [(b, {"size": 10.5, "color": SUB})], line=1.3)
    if i < len(steps) - 1:
        txt(sl, x + SW - 2, 176, 14, 20, [("›", {"size": 15, "color": MUTED})])

y0 = panel(sl, 70, 256, 400, 122, GOLD, "③ 儀式の中身（90秒）")
for i, (a, b) in enumerate([
    ("手拍子を1回", "録音と動画の時刻を合わせる"),
    ("騒音計と並べて無言で60秒", "何デシベルだったかの基準を取る"),
    ("ベルを前後左右から1打ずつ", "マイクの正面と体の正面のズレを測る"),
]):
    yy = y0 + i * 26
    n = rect(sl, 88, yy + 2, 16, 16, GOLD, shape=MSO_SHAPE.OVAL)
    label(n, str(i + 1), size=9.5)
    txt(sl, 112, yy, 350, 22,
        [[(a + "　", {"size": 10.5, "bold": True}),
          (b, {"size": 10, "color": SUB})]])

y0 = panel(sl, 490, 256, 400, 122, GREEN, "帰宅後、その日のうちに")
for i, t in enumerate([
    "録音を学習データと同じ形式に変換する",
    "1件ずつの10秒クリップに切り出す",
    "紙の記録を表に打ち込む（注釈）",
    "検査スクリプトで形式ミスを機械的に見つける",
]):
    bullet(sl, 508, y0 + i * 22, 366, [(t, {"color": SUB})], size=10.5,
           mark=GREEN, h=20, line=1.2)

rect(sl, 70, 392, 820, 40, PALE)
rect(sl, 70, 392, 4, 40, PUR)
txt(sl, 88, 402, 790, 24,
    [[("この一連の流れは、9/1に手元で通し確認まで済ませてあります。", {"size": 11.5, "bold": True}),
      ("機材が届いたその日から本番に入れます。", {"size": 11.5, "color": SUB})]])

# ==================================================== P9 なぜ回数が読めないのか
sl = new_slide("なぜ「何回で終わるか」が読めないのか", 9, active="実録の計画")
txt(sl, 70, 100, 820, 20,
    [("210回のうち、自分の都合だけでは録れないものが揃っています。",
      {"size": 12, "color": MUTED})])

cards = [
    ("天候を待つ", "8本", RED,
     ["雨上がり4本・風のある日4本を", "録ると決めてあります",
      "その天気になるまで録れません"]),
    ("遭遇を待つ", "8本", GOLD,
     ["サイレン・クラクションは", "自分では鳴らせません",
      "「0本でも成立」と決めた枠です"]),
    ("人の予定が要る", "10本", PUR,
     ["EV・HVの接近と停車→発進は", "協力者に運転してもらいます",
      "大学駐車場の許可も要ります"]),
]
for i, (t, n, col, lines) in enumerate(cards):
    x = 70 + i * 278
    rect(sl, x, 130, 264, 168, WHITE, line=LINE, lw=1.0)
    rect(sl, x, 130, 264, 3.5, col)
    txt(sl, x + 16, 146, 150, 22, [(t, {"size": 14, "bold": True})])
    txt(sl, x + 170, 142, 80, 28,
        [(n, {"size": 20, "bold": True, "color": col})], align=PP_ALIGN.RIGHT)
    txt(sl, x + 16, 182, 232, 100,
        [[(l, {"size": 10.5, "color": SUB})] for l in lines], line=1.5)

rect(sl, 70, 314, 820, 66, BAND)
txt(sl, 88, 324, 790, 48,
    [[("さらに、風速3m/sを超える日と雨の日は録らないと決めています。",
       {"size": 12, "bold": True})],
     [("9月から10月は台風と秋雨の時期です。借りている期間がまるごと空振りになることが、"
       "十分に起こり得ます。", {"size": 11.5, "color": SUB})]], line=1.4)

rect(sl, 70, 394, 820, 64, PALEG)
rect(sl, 70, 394, 4, 64, GOLD)
txt(sl, 88, 404, 790, 48,
    [[("つまり、210回は「まとめて3日で片づける」ことができません。",
       {"size": 12.5, "bold": True})],
     [("晴れた日・協力者の空いた日・雨上がりの日を、9月末から11月中旬にかけて"
       "拾い集める形になります。", {"size": 11.5, "color": SUB})]], line=1.4)

# ==================================================== P10 お願い
sl = new_slide("お願い：マイクは「レンタル」ではなく「購入」で", 10, active="お願い")
txt(sl, 70, 100, 820, 20,
    [("前の3枚のとおり、天候と人の都合に合わせて何度も現場に出ることになります。",
      {"size": 12, "color": MUTED})])

reasons = [
    ("1", "総額が逆転しやすい",
     "レンタル3回で13,140円に対し新品は25,909円。天候で1〜2回借り直せばほぼ並び、"
     "それ以上は購入の方が安くなります。"),
    ("2", "測定の条件がそろう",
     "音量の基準合わせは、同じ1台なら最初に1回やれば全収録に通ります。"
     "レンタルは毎回ちがう個体で、やり直すたびに測定条件が変わります。"),
    ("3", "修士の測定の基準器になる",
     "修士でデバイスを作ると自前のマイクが要りますが、その正しさは較正済みの機材と"
     "同じ場面を録って突き合わせる以外に確かめようがありません。"),
]
RY, RH = 132, 74
for i, (no, head, body) in enumerate(reasons):
    y = RY + i * (RH + 10)
    rect(sl, 70, y, 505, RH, WHITE, line=LINE, lw=1.0)
    n = rect(sl, 84, y + 10, 22, 22, PUR, shape=MSO_SHAPE.OVAL)
    label(n, no, size=12)
    txt(sl, 114, y + 9, 450, 20, [(head, {"size": 13, "bold": True})])
    txt(sl, 114, y + 31, 450, 38, [(body, {"size": 10.5, "color": SUB})], line=1.35)

TX, TY, TW, TH = 598, 132, 292, 242
rect(sl, TX, TY, TW, TH, WHITE, line=LINE, lw=1.0)
rect(sl, TX, TY, TW, 3.5, PUR)
txt(sl, TX + 16, TY + 14, TW - 32, 20,
    [("購入をお願いしたい品目", {"size": 13, "bold": True})])
rect(sl, TX + 16, TY + 38, TW - 32, 1.0, LINE)
rows2 = [
    ("Zoom H3-VR（4方向マイク）", "約26,000円",
     "要件を満たす現行機はこれだけ。中古なら約18,000円〜"),
    ("騒音計（平均音量が出るもの）", "0〜10,000円",
     "学内で借りられるなら0円です"),
    ("小物（ハーネス・変換金具ほか）", "約6,000円",
     "自費でもまかなえる範囲です"),
]
ry = TY + 46
for name, price, note in rows2:
    txt(sl, TX + 16, ry, 178, 18, [(name, {"size": 10.5})])
    txt(sl, TX + 190, ry, 86, 18,
        [(price, {"size": 11, "bold": True})], align=PP_ALIGN.RIGHT)
    txt(sl, TX + 16, ry + 19, TW - 32, 26,
        [(note, {"size": 9, "color": MUTED})], line=1.25)
    ry += 52
rect(sl, TX + 16, ry - 6, TW - 32, 1.0, LINE)
txt(sl, TX + 16, ry + 4, 130, 20, [("合計", {"size": 12, "bold": True})])
txt(sl, TX + 146, ry + 3, 130, 22,
    [("約3〜4万円", {"size": 15, "bold": True, "color": RED})], align=PP_ALIGN.RIGHT)

rect(sl, 70, 386, 820, 72, PALEG)
rect(sl, 70, 386, 4, 72, GOLD)
txt(sl, 88, 394, 790, 58,
    [[("承認をいただけたら、すぐ発注して9月末に受け取り、動作確認とリハーサルを経て"
       "10月から本番に入ります。", {"size": 11.5, "color": SUB})],
     [("12月末までに解析まで終える計画です。1月からは執筆に専念します。",
       {"size": 11.5, "bold": True})]], line=1.45)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"ゼミ相談スライド_2026-09-15.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT, f"({len(prs.slides.__iter__.__self__._sldIdLst)}枚)")
