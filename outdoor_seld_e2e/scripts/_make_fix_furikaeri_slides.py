# -*- coding: utf-8 -*-
"""【修正スライド2枚】夏ゼミの振り返りを「番号つき一覧＋1件ずつ」に組み替える。

本人の依頼（2026-09-03）:
 7枚目に①〜④の番号を振り、後ろのスライドで1件ずつ見ていく形にする。
 ②と③は中身が薄いので1枚にまとめる。デザインは元のまま。

出す2枚:
 p7  ①〜④の一覧（コメント＋やったこと1行＋どの枚に詳細があるか）
 p10 ②③ 触覚の先行調査 ／ 前方も記録する

参照した先の枚:
 ① → 8・9枚目（既存）
 ④ → 11枚目（= 修正_デモ時間軸_2026-09-15.pptx）

数値の正本:
 ② md/research/触覚方向提示_先行調査_2026-08-30.md
    （8振動子・間隔107mm→定位98% / 12振動子・間隔72mm→74% /
      首回り35〜40cmで8方向は間隔45〜50mm / 首元は4〜6方向が現実的）
    ⚠️ レビュー値。卒論掲載前に原典で再確認すること
 ③ md/design/実録ハンドブック_2026-08-13.md §8（2026-08-30決定）

出力: md/seminar/修正_振り返り_一覧と②③_2026-09-15.pptx（2枚）
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF5, 0xF3, 0xF9)
PALEG = RGBColor(0xFD, 0xF7, 0xE8)
BAND = RGBColor(0xF2, 0xF3, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
SECTIONS = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(sl, x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(sl, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, text, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)
    return sh


def new_slide(title, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 44, 14, 1.4, H - 28, PUR)
    rect(sl, 58, 46, 30, 6, PUR)
    txt(sl, 98, 30, W - 170, 40,
        [(title, {"size": 22, "bold": True, "spc": 250})])
    rect(sl, 44, 90, W - 90, 1.2, INK)
    txt(sl, 58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
    ws = [16 + len(s) * 12.0 for s in SECTIONS]
    x = (W - (sum(ws) + 8 * (len(SECTIONS) - 1))) / 2
    for s, wd in zip(SECTIONS, ws):
        on = (s == "振り返り")
        c = rect(sl, x, 502, wd, 22, NAVY if on else CHIP)
        label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
        x += wd + 8
    txt(sl, W - 90, 506, 40, 20,
        [(str(page), {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)
    return sl


# ==================================================== p7 一覧
sl = new_slide("夏 ゼ ミ の 振 り 返 り", 7)
txt(sl, 70, 100, 820, 20,
    [[("いただいた4件すべてに手を動かしました。", {"size": 12.5, "bold": True}),
      ("　この後、1件ずつご説明します。", {"size": 12, "color": SUB})]])

rows = [
    ("①", "狭いところをすり抜けていく車と対向してくる車が、"
     "同じ至近警告になって区別できないのではないか",
     "方位の変化率で区別する仕組みを実装した", "8・9枚目", RED),
    ("②", "振動による方向指示のデバイスは、全盲や弱視の方向けの"
     "製品やノウハウもいろいろある",
     "先行研究を調べ、首元は4〜6方向が現実的と分かった", "10枚目", GOLD),
    ("③", "前方を検出対象外としているが、研究上はとりあえず"
     "全方位録ってみるでいいのでは",
     "記録紙に前方の欄を追加した（追加費用はゼロ）", "10枚目", GOLD),
    ("④", "振動デバイスを一から作るとハードルが高そう。"
     "Joy-conとUnityとかで簡単に作ってみてもいいかも",
     "作った。9/2に実機で振動した", "11枚目", GREEN),
]
y = 140.0
for no, quote, ans, where, col in rows:
    n = rect(sl, 70, y + 2, 28, 28, col, shape=MSO_SHAPE.OVAL)
    label(n, no, size=14)
    txt(sl, 112, y, 640, 34, [(quote, {"size": 11, "color": SUB})], line=1.35)
    b = rect(sl, 786, y + 4, 104, 22, None, line=col, lw=1.1)
    label(b, where, size=10, color=col)
    txt(sl, 112, y + 38, 30, 20,
        [("→", {"size": 13, "bold": True, "color": PUR})])
    txt(sl, 140, y + 38, 750, 22, [(ans, {"size": 12.5, "bold": True})])
    y += 72
    if no != "④":
        rect(sl, 70, y - 8, 820, 0.8, HAIR)

# ==================================================== p10 ②③
sl = new_slide("②③　触 覚 の 先 行 調 査 と 、 前 方 の 記 録", 10)
txt(sl, 70, 100, 820, 20,
    [("どちらもご助言のとおりに動きました。②は調査、③は収録前に決めました。",
      {"size": 12, "color": MUTED})])

# --- ② 左 ---
BX, BY, BW, BH = 70, 130, 400, 328
rect(sl, BX, BY, BW, BH, WHITE, line=LINE, lw=1.0)
rect(sl, BX, BY, BW, 3.5, GOLD)
txt(sl, BX + 18, BY + 14, BW - 36, 22,
    [("② 触覚で方向は、どこまで伝わるのか", {"size": 13.5, "bold": True})])
txt(sl, BX + 18, BY + 40, BW - 36, 18,
    [("腰の振動ベルトで測った先行研究の値", {"size": 10, "color": MUTED})])

hy = BY + 62
rect(sl, BX + 18, hy, BW - 36, 24, BAND)
txt(sl, BX + 26, hy + 4, 150, 16, [("並べ方", {"size": 10, "bold": True})])
txt(sl, BX + 250, hy + 4, 130, 16,
    [("方向を当てられた割合", {"size": 10, "bold": True})], align=PP_ALIGN.RIGHT)
for i, (a, v, col) in enumerate([("8個・間隔 107mm", "98%", GREEN),
                                 ("12個・間隔 72mm", "74%", RED)]):
    ry = hy + 28 + i * 28
    txt(sl, BX + 26, ry, 200, 18, [(a, {"size": 11})])
    txt(sl, BX + 250, ry - 2, 130, 20,
        [(v, {"size": 14, "bold": True, "color": col})], align=PP_ALIGN.RIGHT)
txt(sl, BX + 18, hy + 86, BW - 36, 18,
    [("間隔が狭くなると急に落ちる", {"size": 10, "color": MUTED})])

rect(sl, BX + 18, BY + 176, BW - 36, 1.0, HAIR)
txt(sl, BX + 18, BY + 188, BW - 36, 52,
    [[("首回りは35〜40cmしかないので、8方向にすると間隔45〜50mm。",
       {"size": 10.5, "color": SUB})],
     [("これは腰で74%に落ちた間隔より狭く、しかも首は腰より鈍い。",
       {"size": 10.5, "color": SUB})]], line=1.4)

rect(sl, BX + 18, BY + 244, BW - 36, 40, PALEG)
txt(sl, BX + 30, BY + 253, BW - 60, 24,
    [[("→ 首元は ", {"size": 12}),
      ("4〜6方向", {"size": 14, "bold": True, "color": RED}),
      (" が現実的", {"size": 12})]])
txt(sl, BX + 18, BY + 292, BW - 36, 30,
    [[("いまの構想は振動子5個。この範囲に収まっている。",
       {"size": 10.5, "bold": True})],
     [("※ レビュー論文からの引用。卒論に載せる前に原典で確認する。",
       {"size": 9, "color": MUTED})]], line=1.35)

# --- ③ 右 ---
CX = 490
rect(sl, CX, BY, BW, BH, WHITE, line=LINE, lw=1.0)
rect(sl, CX, BY, BW, 3.5, GOLD)
txt(sl, CX + 18, BY + 14, BW - 36, 22,
    [("③ 前方も記録する（8/30に決定）", {"size": 13.5, "bold": True})])

rect(sl, CX + 18, BY + 44, BW - 36, 44, PALE)
txt(sl, CX + 30, BY + 52, BW - 60, 30,
    [("マイクは最初から全方位を録っているので、追加の費用も手間もゼロ。",
      {"size": 11, "bold": True})], line=1.35)

cy = BY + 100
for a, b in [("変えるのは記録紙だけ",
              "前方のイベントも1行書く。それ以外は何も変えない"),
             ("採点では既定で除外する",
              "別の列に入れて別集計。前方は通知の対象外という前提は保つ"),
             ("110本の内訳は動かさない",
              "録る前に決めた配分（A〜F）はそのまま"),
             ("録らなければ二度と取れない",
              "解析の段階なら後から外せるが、録音は撮り直せない")]:
    rect(sl, CX + 18, cy + 6, 6, 6, GOLD)
    txt(sl, CX + 34, cy - 1, BW - 56, 18, [(a, {"size": 11.5, "bold": True})])
    txt(sl, CX + 34, cy + 17, BW - 56, 30,
        [(b, {"size": 9.5, "color": MUTED})], line=1.3)
    cy += 52

rect(sl, CX + 18, BY + 300, BW - 36, 1.0, HAIR)
txt(sl, CX + 18, BY + 310, BW - 36, 18,
    [("データを見て思いがけないことが分かる可能性を、残しました。",
      {"size": 10, "color": SUB})])

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"修正_振り返り_一覧と②③_2026-09-15.pptx")
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ PowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
