# -*- coding: utf-8 -*-
"""P13（実録評価の測定方法）の下半分に貼る「収録→採点タイムライン」部品。

本人テンプレ風（白箱・細グレー枠・紺の番号・金の矢印）。1枚のpptxとして出力し、
Ctrl+A→コピー→P13の下半分へ貼り付けて使う。数値は計画値（事前登録）表記。
出力: md/seminar/図_実録タイムライン_2026-08-13.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)


def para(tf, runs, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for t, o in runs:
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(o.get("size", 11))
        f.bold = o.get("bold", False)
        f.color.rgb = o.get("color", INK)
        f.name = "Meiryo"
        meiryo(r)
    return p


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


# ---- 見出し行 ----
tb = sl.shapes.add_textbox(Pt(70), Pt(320), Pt(400), Pt(24))
tb.text_frame.word_wrap = False
para(tb.text_frame, [("収録から採点までの流れ（計画）", {"size": 13, "bold": True})],
     first=True)

# ---- 5ステップのタイムライン ----
steps = [
    ("①", "機材レンタル", "H3-VR・7泊8日\n（この中で全工程）"),
    ("②", "初日リハ", "設定を固定し\n試し録りを機械チェック"),
    ("③", "統制録音", "大学駐車場で\n台本どおり反復（半日）"),
    ("④", "機会捕捉＋負例", "徒歩ルートで随時\n（サイレンは最優先）"),
    ("⑤", "変換・注釈・採点", "現地メモ→注釈CSV\n→自動採点（半日）"),
]
x, y, bw, bh, gap = 66.0, 352.0, 148.0, 92.0, 22.0
for i, (no, name, body) in enumerate(steps):
    rect(x, y, bw, bh, WHITE, line=LINE, lw=1.0)
    tb = sl.shapes.add_textbox(Pt(x + 6), Pt(y + 6), Pt(bw - 12), Pt(bh - 12))
    tf = tb.text_frame
    tf.word_wrap = True
    para(tf, [(no + " ", {"size": 12, "bold": True, "color": GOLD}),
              (name, {"size": 11.5, "bold": True, "color": NAVY})], first=True)
    for ln in body.split("\n"):
        para(tf, [(ln, {"size": 9.5, "color": SUB})])
    if i < 4:
        ar = rect(x + bw + 3, y + bh / 2 - 8, 16, 16, GOLD,
                  shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        ar.rotation = 90
    x += bw + gap

# ---- 下の注記 ----
tb = sl.shapes.add_textbox(Pt(70), Pt(456), Pt(830), Pt(22))
tb.text_frame.word_wrap = False
para(tb.text_frame,
     [("実働3日＋注釈1日・計120本＝共通100＋歩行対比20（計画値）", {"size": 11, "bold": True, "color": INK}),
      ("　※1テイクの具体的な手順は付録「実録の流れ」を参照", {"size": 10.5, "color": MUTED})],
     first=True)

OUT = r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/図_実録タイムライン_2026-08-13.pptx"
prs.save(OUT)
print("saved:", OUT)
