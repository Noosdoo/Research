# -*- coding: utf-8 -*-
"""「実録評価の測定方法」スライド・本人レイアウト維持版（3枠＋幾何図）。

v1(_make_fig_jitsuroku_slide.py)は4分類に組み替えたが、本人が3枠のまま
直すことを希望したためレイアウトを戻し、事実誤りだけを潰した版。

直した点（すべて事前登録との食い違い）:
 1. 踏切を統制から外した（時刻表どおり＝待ち受け側）
 2. バイクを削除（A〜Eのどこにも配分がない＝統制で録る計画が存在しない）
 3. 車の走行音を統制から外した（A20は公道の受動観測・横距離は目測）
 4. コーン(▲)を削除（公道ではコーン設置に警察・道路管理者の事前確認が要る。
    コーンを使うのは私有地の統制試行のみ）
 5. 横距離 2・3・5m → 2.0〜3.2m ／ 5〜15m（目測）＝事前登録R1の帯
 6. 機材にチェストハーネスと風防を追加（図に「胸部4chマイク」とあるのに欠けていた）

出典: md/design/実録スモーク計画書_2026-07.md R1 / 実録収録計画v2追補_2026-07-28.md
      / 実録ハンドブック_2026-08-13.md
出力: md/seminar/図_実録の測定方法_2026-08-22_v2.pptx（コピペ用・1枚）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
ROAD = RGBColor(0xE9, 0xE9, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


# ---- テンプレ骨格 ----
rect(44, 14, 1.4, H - 28, PUR)
rect(58, 46, 30, 6, PUR)
txt(98, 30, W - 160, 40, [("今後の方針②　実録評価の測定方法（9月予定）",
                           {"size": 23, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)

# ---- 3枠 ----
BY, BH, BW = 118, 124, 270
bxs = [70, 355, 640]
boxes = [
    ("① 統制録音", INK, "自分で音を出す・走らせる（私有地）",
     ["バック音（実物ブザー）／自転車ベル",
      "自転車の至近通過／キックボード",
      "EV・HVの接近／停車→発進"]),
    ("② 待ち受け録音", GOLD, "通るのを待って記録する（公道）",
     ["車の走行音（静穏路・交通量路 半々）",
      "踏切・列車（時刻表どおり＝確実）",
      "サイレン・クラクションは遭遇時のみ"]),
    ("③ 負例", GREEN, "鳴ってはいけない場面",
     ["静穏・繁華街／工事・雨上がり・風",
      "連続録音100分以上を別枠で確保",
      "→ 誤警告率［回/時］を評価"]),
]
for (tag, col, lead, lines), bx in zip(boxes, bxs):
    rect(bx, BY, BW, BH, WHITE, line=LINE, lw=1.0)
    c = rect(bx + 14, BY + 12, 16 + len(tag) * 12.0, 22, None, line=col, lw=1.2)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = tag
    r.font.size = Pt(11.5)
    r.font.bold = True
    r.font.color.rgb = col
    r.font.name = "Meiryo"
    meiryo(r)
    txt(bx + 14, BY + 40, BW - 26, 18,
        [(lead, {"size": 10.5, "color": MUTED})])
    txt(bx + 14, BY + 60, BW - 26, 60,
        [[(t, {"size": 11, "color": SUB})] for t in lines], line=1.38)

# ---- 箇条書き ----
BUY = 258
items = [
    [("機材：", {"bold": True}),
     ("4ch FOAマイク(Zoom H3-VR)＋風防・チェストハーネス", {"color": SUB})],
    [("計測具：", {"bold": True}),
     ("騒音計(LAeq)・メジャー・音源用の自転車ベルとバックブザー", {"color": SUB})],
    [("実環境性能：", {"bold": True}),
     ("静止100本 ／ 歩行影響：静止50本 ↔ 歩行50本の対応比較", {"color": SUB})],
]
for i, runs in enumerate(items):
    rect(72, BUY + i * 26 + 5, 8, 8, INK)
    txt(92, BUY + i * 26, 800, 22, runs, size=13)

# ---- 幾何図 ----
FX, FY, FW, FH = 210, 344, 540, 148
rect(FX, FY, FW, FH, WHITE, line=LINE, lw=1.0)
txt(FX + 16, FY + 10, 300, 20,
    [("車の記録のしかた（公道・待ち受け）", {"size": 13, "bold": True})])
txt(FX + 340, FY + 12, 184, 18,
    [("速度：徐行〜約30km/h", {"size": 10.5, "color": SUB})], align=PP_ALIGN.RIGHT)

rect(FX + 22, FY + 38, FW - 44, 32, ROAD)
ch = rect(FX + 350, FY + 44, 52, 20, GOLD)
tf = ch.text_frame
tf.margin_left = tf.margin_right = Pt(2)
tf.margin_top = tf.margin_bottom = Pt(0)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "車"
r.font.size = Pt(11)
r.font.bold = True
r.font.color.rgb = WHITE
r.font.name = "Meiryo"
meiryo(r)
txt(FX + 408, FY + 42, 40, 22, [("→", {"size": 17, "color": INK})])

rect(FX + 168, FY + 70, 1.6, 34, RED)
rect(FX + 160, FY + 104, 16, 16, INK, shape=MSO_SHAPE.OVAL)
txt(FX + 182, FY + 76, 320, 20,
    [("横距離 2.0〜3.2m ／ 5〜15m（目測）",
      {"size": 12, "bold": True, "color": RED})])
txt(FX + 96, FY + 122, 260, 18,
    [("装着者（胸部4chマイク・静止）", {"size": 10.5, "color": SUB})],
    align=PP_ALIGN.CENTER)

# ---- フッター ----
txt(58, 506, 90, 20, [("2026/09", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "今後・まとめ")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(10.5)
    r.font.bold = active
    r.font.color.rgb = WHITE if active else MUTED
    r.font.name = "Meiryo"
    meiryo(r)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("11", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)

OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/"
           r"図_実録の測定方法_2026-08-22_v2.pptx")
try:                                   # PowerPointで開いたままだと上書きできない
    prs.save(OUT)
except PermissionError:
    OUT = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
    prs.save(OUT)
    print("※ 元のファイルがPowerPointで開かれていたため別名で保存した")
print("saved:", OUT)
