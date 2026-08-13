# -*- coding: utf-8 -*-
"""「データセットと評価手順」スライドを本人テンプレ風デザインで1枚生成。

本人の夏ゼミデッキ（20260804系テンプレ）に合わせた要素:
紫のタイトルダッシュ／見出し下の全幅罫線／左端の細い縦線／字間広めの見出し／
金色左バーの強調枠／下部ナビチップ（検証=アクティブ）
出力: md/seminar/図_データセット評価手順_2026-08-13.pptx（コピペ用・1枚）
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)      # 本文の濃紺黒
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
PUR = RGBColor(0x7E, 0x6F, 0x98)      # テンプレの紫アクセント
GOLD = RGBColor(0xE0, 0xA5, 0x26)     # 強調枠の金色バー
GRAY_HDR = RGBColor(0xEF, 0xEF, 0xEF)
CREAM = RGBColor(0xFB, 0xF4, 0xDE)    # 評価行の淡い強調
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)     # アクティブチップ
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = 960.0, 540.0
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spc=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for t, o in runs:
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", False)
        f.color.rgb = o.get("color", INK)
        f.name = "Meiryo"
        meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(x, y, w, h, fill, line=None, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


# ---- テンプレ骨格 ----
rect(44, 14, 1.4, H - 28, PUR)                      # 左端の細い縦線
rect(58, 46, 30, 6, PUR)                            # タイトル前の紫ダッシュ
txt(98, 30, W - 160, 40, [("検証の前に — データセットと評価手順",
                           {"size": 23, "bold": True, "spc": 250})])
rect(44, 90, W - 90, 1.2, INK)                      # 見出し下の全幅罫線

# ---- 表 ----
rows = [
    ("区分", "中身", "本研究での役割"),
    ("事前学習", "1,167時間・170クラス（PSELDNets付属）", "出発点の基盤モデル（既存を利用）"),
    ("学習（微調整）", "本研究の合成10,200クリップ（8クラス・約28時間）",
     "距離ヘッド込みで全体をファインチューニング"),
    ("検証", "学習と同一方式の別クリップ", "モデル選択・通知閾値の決定（発表数値には不使用）"),
    ("評価", "同一設計・新乱数の1,800クリップ（学習・検証に未使用）",
     "最終結果の測定 — 1回だけ採点"),
]
gt = sl.shapes.add_table(5, 3, Pt(70), Pt(118), Pt(W - 140), Pt(250)).table
gt.columns[0].width = Pt(140)
gt.columns[1].width = Pt(360)
gt.columns[2].width = Pt(320)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gt.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = (GRAY_HDR if ri == 0 else
                                    (CREAM if ri == 4 else WHITE))
        cell.margin_top = cell.margin_bottom = Pt(5)
        cell.margin_left = Pt(8)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = val
        f = r.font
        f.size = Pt(12.5)
        f.bold = (ri == 0 or ci == 0)
        f.color.rgb = INK if (ri == 0 or ci == 0) else SUB
        f.name = "Meiryo"
        meiryo(r)

# ---- 金色バーの強調枠（テンプレのまとめ枠風） ----
rect(70, 396, 6, 62, GOLD)
rect(76, 396, W - 146, 62, WHITE, line=LINE, lw=1.0)
txt(94, 404, W - 180, 46, [
    ("基準を先に決める → 未使用データを新たに生成 → 1回だけ評価", {"bold": True, "size": 14}),
    ("　＝ 答案を見てから基準を変えられない手順", {"size": 13, "color": SUB}),
], anchor=MSO_ANCHOR.MIDDLE)

# ---- 下部フッター（ナビチップ・検証アクティブ） ----
txt(58, 506, 90, 20, [("2026/09", {"size": 11, "color": MUTED})])
steps = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
x = 300.0
for s in steps:
    wd = 16 + len(s) * 12.0
    active = (s == "検証")
    c = rect(x, 502, wd, 22, NAVY if active else CHIP)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(10.5)
    r.font.bold = active
    r.font.color.rgb = WHITE if active else MUTED
    r.font.name = "Meiryo"
    meiryo(r)
    x += wd + 8
txt(W - 90, 506, 40, 20, [("9", {"size": 11, "color": MUTED})],
    align=PP_ALIGN.RIGHT)

OUT = r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/図_データセット評価手順_2026-08-13.pptx"
prs.save(OUT)
print("saved:", OUT)
