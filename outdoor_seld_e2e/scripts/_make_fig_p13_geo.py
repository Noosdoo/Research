# -*- coding: utf-8 -*-
"""P13下半分用（タイムライン案の代替）: 統制録音の幾何俯瞰図＋確かめること3点。

本人テンプレ風。左=俯瞰図（歩行者・通過車・横距離コーン・速度）、
右=「実録で確かめること」カード。1枚pptxで出力しCtrl+A→P13下半分へ貼る。
出力: md/seminar/図_実録幾何図_2026-08-13.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
GOLD = RGBColor(0xE0, 0xA5, 0x26)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xCB, 0xD2)
ROAD = RGBColor(0xE9, 0xE9, 0xE4)
RED = RGBColor(0xC0, 0x50, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
sl = prs.slides.add_slide(prs.slide_layouts[6])


def meiryo(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)


def txt(x, y, w, h, runs, align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    first = True
    if isinstance(runs, tuple):
        runs = [runs]
    for item in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        rs = item if isinstance(item, list) else [item]
        for t, o in rs:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", 11))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r)
    return tb


def shp(x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
        dash=None):
    s = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
        if dash:
            ln = s.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    s.shadow.inherit = False
    return s


Y0 = 318.0   # 下半分の開始

# ================= 左: 俯瞰図 =================
shp(70, Y0, 520, 158, WHITE, line=LINE)
txt(84, Y0 + 6, 300, 20, ("統制録音の幾何（上から見た図）",
                          {"size": 12, "bold": True, "color": NAVY}))

# 道路帯
shp(84, Y0 + 34, 492, 44, ROAD)
# 車（金）＋進行矢印
car = shp(300, Y0 + 44, 64, 26, GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tf = car.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = tf.margin_right = Pt(1)
tf.margin_top = tf.margin_bottom = Pt(0)
p0 = tf.paragraphs[0]
p0.alignment = PP_ALIGN.CENTER
r0 = p0.add_run()
r0.text = "車"
r0.font.size = Pt(11)
r0.font.bold = True
r0.font.color.rgb = WHITE
r0.font.name = "Meiryo"
meiryo(r0)
ar = shp(374, Y0 + 50, 30, 13, NAVY, shape=MSO_SHAPE.RIGHT_ARROW)
txt(300, Y0 + 12, 220, 18, ("速度：徐行〜約30km/h（台本）",
                            {"size": 9.5, "color": SUB}), wrap=False)
# 歩行者（紺丸＋金リング）
shp(180, Y0 + 118, 22, 22, None, line=GOLD, lw=2.0, shape=MSO_SHAPE.OVAL)
shp(185, Y0 + 123, 12, 12, NAVY, shape=MSO_SHAPE.OVAL)
txt(120, Y0 + 138, 200, 16, ("歩行者（胸に4chマイク）",
                             {"size": 9.5, "color": SUB}), wrap=False)
# 横距離の破線＋コーン
shp(190, Y0 + 78, 1.2, 40, None, line=RED, lw=1.5, dash="dash")
txt(200, Y0 + 88, 220, 16, [[("横距離 2・3・5m", {"size": 10, "bold": True, "color": RED}),
                             ("（コーンで実測）", {"size": 9.5, "color": SUB})]], wrap=False)
for cx in (168, 214):
    shp(cx, Y0 + 72, 10, 10, GOLD, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
# 下注記
txt(84, Y0 + 166 - 4, 492, 16,
    ("速度×横距離を台本で変えて反復 ＝ 合成シナリオと同じ条件を実環境で再現",
     {"size": 9.5, "color": MUTED}), wrap=False)

# ================= 右: 確かめることカード =================
shp(610, Y0, 290, 158, WHITE, line=LINE)
shp(610, Y0, 6, 158, GOLD)
txt(628, Y0 + 8, 260, 20, ("この測定で確かめること",
                           {"size": 12, "bold": True, "color": NAVY}))
txt(628, Y0 + 34, 262, 120, [
    [("① ", {"size": 11, "bold": True, "color": GOLD}),
     ("合成で測った検出・方向・距離の性能が実環境でどれだけ保たれるか", {"size": 10.5, "color": SUB})],
    [("② ", {"size": 11, "bold": True, "color": GOLD}),
     ("誤警告が1時間あたり何回出るか（負例）", {"size": 10.5, "color": SUB})],
    [("③ ", {"size": 11, "bold": True, "color": GOLD}),
     ("切り分け実験（前頁）各版の共通試験になる", {"size": 10.5, "color": SUB})],
])

OUT = r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/図_実録幾何図_2026-08-13.pptx"
prs.save(OUT)
print("saved:", OUT)
