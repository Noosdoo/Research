# -*- coding: utf-8 -*-
"""14枚目「収録の流れ」を、伝わるデザインの2ルールで作り直した試作。

参考（2026-09-05に調査）:
  伝わるデザイン｜研究発表のユニバーサルデザイン（千葉大・高橋佑磨）
  https://tsutawarudesign.com/
  - 「使う色は、背景や文字の色を含めて４色ぐらいに」
  - 「メインの色は…出現頻度が高くなってもあまり不快にならない落ち着いた色」
  - 「強調の色は、最も重要な箇所のみに用いる」
  - 「囲い・枠が多くなると全体が煩雑になるので、濫用は避けましょう」
  - 「輪郭のはっきりしないイラストやグラフ」にだけ薄い灰色の四角で輪郭を作る
  - 「最低でも本文の文字の１文字分の余白」「揃えられるところはすべて揃える」
  - 「決してスライドをデコらないでください」

当てた2ルール:
  1. 装飾に使う色は「紫」と「灰」だけ。赤・金・緑は危険度を表すときだけ
     → この枚に危険度の話は無いので、6コマの赤・金・緑を全部やめた（4色→3色）
  2. 囲みは「輪郭のはっきりしない図」にだけ使う。文字の塊は見出しと余白で分ける
     → 囲み 9個 → 0個。流れは横罫と縦罫で示す

出力: md/seminar/p14_compare.pptx（1枚目=現行のコピー、2枚目=新案）
"""
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x22, 0x28, 0x38)
SUB = RGBColor(0x59, 0x5F, 0x6E)
MUTED = RGBColor(0x8A, 0x8F, 0x9A)
HAIR = RGBColor(0xDD, 0xDF, 0xE4)
PUR = RGBColor(0x7E, 0x6F, 0x98)
CHIP = RGBColor(0xE8, 0xE8, 0xE8)
NAVY = RGBColor(0x23, 0x2B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = Path(r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/p14_compare.pptx")
SECTIONS = ["近況報告", "研究の前提", "振り返り", "実録の計画", "相談", "まとめ"]
W, H = 960.0, 540.0


def meiryo(run, spc=None):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", "メイリオ")
    rPr.append(ea)
    if spc:
        rPr.set("spc", str(spc))


def txt(sl, x, y, w, h, runs, size=12, align=PP_ALIGN.LEFT, spc=None, line=None):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    if isinstance(runs, list) and runs and isinstance(runs[0], tuple):
        runs = [runs]
    for pi, prun in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        for t, o in prun:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.color.rgb = o.get("color", INK)
            f.name = "Meiryo"
            meiryo(r, spc=o.get("spc", spc))
    return tb


def rect(sl, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def label(sh, t, size=11, color=WHITE, bold=True):
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Meiryo"
    meiryo(r)


def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(
            copy.deepcopy(shp._element), "p:extLst")
    return dst


def build_new(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    # ---- テンプレの骨格（元のまま） ----
    rect(sl, 44, 14, 1.4, H - 28, PUR)
    rect(sl, 58, 46, 30, 6, PUR)
    txt(sl, 98, 30, W - 170, 40,
        [("収 録 の 流 れ", {"size": 22, "bold": True, "spc": 250})])
    rect(sl, 44, 90, W - 90, 1.2, INK)

    txt(sl, 70, 102, 820, 20,
        [[("1地点あたり6つの手順。", {"size": 12.5}),
          ("④だけをテイクごとに繰り返す。", {"size": 12.5, "bold": True})]])

    # ---- 流れ: 囲まず、横罫と縦罫だけで示す ----
    FY, FW = 146.0, 820.0 / 6
    rect(sl, 70, FY, 820, 1.2, INK)              # 流れの基線
    steps = [("①", "前日", "15分", "設定を家で全部固定する"),
             ("②", "到着", "2分", "定位置に立つ"),
             ("③", "儀式", "90秒", "手拍子・騒音計・ベル4方位"),
             ("④", "本番", "20〜30秒 ×n", "無言。1通過＝1テイク"),
             ("⑤", "直後", "30秒", "5秒数えてから声で記録"),
             ("⑥", "確認", "5分", "その場で再生して品質確認")]
    for i, (no, nm, tm, desc) in enumerate(steps):
        x = 70 + i * FW
        if i:
            rect(sl, x, FY + 8, 0.8, 92, HAIR)   # 区切りは細い縦罫だけ
        txt(sl, x + 12, FY + 12, 30, 26,
            [(no, {"size": 19, "bold": True, "color": PUR})])
        txt(sl, x + 40, FY + 17, FW - 52, 22, [(nm, {"size": 14, "bold": True})])
        txt(sl, x + 12, FY + 44, FW - 24, 18, [(tm, {"size": 10, "color": MUTED})])
        txt(sl, x + 12, FY + 64, FW - 24, 40,
            [(desc, {"size": 10.5, "color": SUB})], line=1.35)
    rect(sl, 70, FY + 104, 820, 0.8, HAIR)

    x4 = 70 + 3 * FW
    txt(sl, x4 + 12, FY + 108, FW * 1.4, 18,
        [("↻ 車が通るたびに④を繰り返す", {"size": 9.5, "color": PUR, "bold": True})])

    # ---- 下段: 囲まず、見出しと余白と縦罫で分ける ----
    BY = 300.0
    rect(sl, 70, BY, 820, 0.8, HAIR)
    rect(sl, 480, BY + 14, 0.8, 128, HAIR)

    txt(sl, 70, BY + 14, 380, 24,
        [[("③ 儀式の中身", {"size": 15, "bold": True}),
          ("　セッション冒頭に1回・90秒", {"size": 10, "color": MUTED})]])
    y = BY + 48
    for i, (a, b) in enumerate([
        ("手拍子を1回", "録音と動画の時刻を合わせる"),
        ("騒音計と並べて無言で60秒", "何デシベルだったかの基準を取る"),
        ("ベルを前後左右から1打ずつ", "マイクの正面と体の正面のズレを測る"),
    ]):
        txt(sl, 70, y, 20, 20,
            [("%d" % (i + 1), {"size": 12, "bold": True, "color": PUR})])
        txt(sl, 92, y, 340, 20, [(a, {"size": 11.5, "bold": True})])
        txt(sl, 92, y + 19, 340, 18, [(b, {"size": 10, "color": SUB})])
        y += 42

    txt(sl, 510, BY + 14, 380, 24,
        [[("帰宅後、その日のうちに", {"size": 15, "bold": True}),
          ("　半日", {"size": 10, "color": MUTED})]])
    y = BY + 50
    for t in ["録音を学習データと同じ形式に変換する",
              "1件ずつの10秒クリップに切り出す",
              "紙の記録を表に打ち込む",
              "検査スクリプトで形式ミスを機械的に見つける"]:
        rect(sl, 510, y + 6, 5, 5, PUR)
        txt(sl, 524, y, 366, 20, [(t, {"size": 11.5, "color": SUB})])
        y += 26

    # ---- フッター（元のまま） ----
    txt(sl, 58, 506, 110, 20, [("2026/09/15", {"size": 11, "color": MUTED})])
    ws = [16 + len(s) * 12.0 for s in SECTIONS]
    x = (W - (sum(ws) + 8 * (len(SECTIONS) - 1))) / 2
    for s, wd in zip(SECTIONS, ws):
        on = (s == "実録の計画")
        c = rect(sl, x, 502, wd, 22, NAVY if on else CHIP)
        label(c, s, size=10.5, color=WHITE if on else MUTED, bold=on)
        x += wd + 8
    txt(sl, W - 90, 506, 40, 20,
        [("14", {"size": 11, "color": MUTED})], align=PP_ALIGN.RIGHT)
    return sl


def main() -> int:
    s = Presentation(str(SRC))
    out = Presentation()
    out.slide_width = s.slide_width
    out.slide_height = s.slide_height
    clone(list(s.slides)[13], out)          # 現行の14枚目
    build_new(out)                          # 新案
    out.save(str(OUT))
    print("saved:", OUT)
    print("1枚目=現行 / 2枚目=新案（囲み9→0・装飾の色4→2）")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
