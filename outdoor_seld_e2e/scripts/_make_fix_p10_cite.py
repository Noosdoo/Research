# -*- coding: utf-8 -*-
"""p10「②③触覚の先行調査と前方の記録」の出典表記を原典に直す（2026-09-07）。

「腰の振動ベルトで測った先行研究の値(bioRxiv 2021)」の数値
  8振動子・間隔107mm→98% / 12振動子・間隔72mm→74%
は、bioRxiv 2021（Fadaei Jouybari et al.; 後に Exp Brain Res 2021, doi:10.1007/s00221-021-06181-x）
が **Cholewiak et al. (2004) を引用して述べている値** で、その論文自身の結果ではない。
原典: Cholewiak, Brill & Schwab (2004) "Vibrotactile localization on the abdomen: Effects of
place and space", Perception & Psychophysics 66(6), 970-987. doi:10.3758/BF03194989
→ スライドの括弧書きを (Cholewiak et al., 2004) に差し替える。

本人のデッキから p10 を複製し、該当 run の文字列だけ置換する。書式は元のまま。
出力: md/seminar/修正_触覚の出典_2026-09-15.pptx（1枚）
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(r"C:/Users/satos/research/outdoor_seld_e2e")
SRC = Path(r"C:/Users/satos/iCloudDrive/５松澤研究室/データ解析ゼミ/2026/0915"
           r"/20260915_B4_松本鋭.pptx")
OUT = ROOT / "md/seminar/修正_触覚の出典_2026-09-15.pptx"
TITLE = "触覚の先行調査と前方の記録"
OLD, NEW = "bioRxiv 2021", "Cholewiak et al., 2004"


def find(prs, title):
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and title in sh.text_frame.text and sh.top / 12700 < 70:
                return sl
    raise SystemExit("見出し「%s」のスライドが見つからない" % title)


def clone(src, prs):
    dst = prs.slides.add_slide(prs.slide_layouts[6])
    for shp in src.shapes:
        dst.shapes._spTree.insert_element_before(copy.deepcopy(shp._element), "p:extLst")
    return dst


def fix_footer(sl):
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == "2026/09/15":
            sh.left, sh.top, sh.width, sh.height = int(Pt(66)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))
        elif t.isdigit() and sh.top / 12700 > 480:
            sh.left, sh.top, sh.width, sh.height = int(Pt(678)), int(Pt(500.5)), int(Pt(216)), int(Pt(28.8))


def main() -> int:
    src = Presentation(str(SRC))
    out = Presentation()
    out.slide_width, out.slide_height = src.slide_width, src.slide_height
    sl = clone(find(src, TITLE), out)

    hits = 0
    for sh in sl.shapes:
        if not sh.has_text_frame or OLD not in sh.text_frame.text:
            continue
        for p in sh.text_frame.paragraphs:
            runs = p.runs
            joined = "".join(r.text for r in runs)
            if OLD not in joined:
                continue
            # 1つの run に収まっていればそこだけ、またがっていれば段落を1 runにまとめて置換
            done = False
            for r in runs:
                if OLD in r.text:
                    r.text = r.text.replace(OLD, NEW)
                    done = True
            if not done:
                runs[0].text = joined.replace(OLD, NEW)
                for r in list(runs[1:]):
                    r._r.getparent().remove(r._r)
            hits += 1
    assert hits == 1, "「%s」を含む段落が %d 個" % (OLD, hits)
    fix_footer(sl)

    dst = OUT
    try:
        out.save(str(dst))
    except PermissionError:
        dst = OUT.with_name(OUT.stem + "_新" + OUT.suffix)
        out.save(str(dst))
    print("保存:", dst)
    print("置換: (%s) → (%s)" % (OLD, NEW))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
