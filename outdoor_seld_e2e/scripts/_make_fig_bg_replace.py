# -*- coding: utf-8 -*-
"""研究背景スライドの右図 差し替え案（2案）を図だけのpptxとして生成。

案A: レーダー型 — 全方位リング＋前方視野の扇＋接近矢印＋音の波紋（推し）
案B: 波紋型 — 矢印なし・音の広がりだけで見せるミニマル版
コピー方法: スライド上で Ctrl+A → Ctrl+C → 自分のデッキに Ctrl+V → Ctrl+G でグループ化
"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x16, 0x23, 0x3A)
INK2 = RGBColor(0x3A, 0x46, 0x58)
MUTED = RGBColor(0x66, 0x70, 0x7F)
HAIR = RGBColor(0xD9, 0xDA, 0xD2)
AMBER = RGBColor(0xE8, 0xA2, 0x00)
RED = RGBColor(0xC4, 0x43, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
VISION = RGBColor(0xDE, 0xE3, 0xEC)   # 前方視野の扇（薄い紺系）
PAPERISH = RGBColor(0xF7, 0xF7, 0xF3)


def blend(c, t):
    """色cを紙色PAPERISHへ割合tだけ寄せる（波紋の外側を淡く）。"""
    return RGBColor(round(c[0] + (PAPERISH[0] - c[0]) * t),
                    round(c[1] + (PAPERISH[1] - c[1]) * t),
                    round(c[2] + (PAPERISH[2] - c[2]) * t))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def meiryo(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", "メイリオ")


def put_text(shape, text, size, color, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = "Meiryo"
    meiryo(r)


def oval(sl, cx, cy, rx, ry, fill=None, line=None, line_w=1.0, dash=None):
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - rx), Inches(cy - ry),
                             Inches(2 * rx), Inches(2 * ry))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
        if dash:
            ln = sh.line._get_or_add_ln()
            d = ln.makeelement(qn("a:prstDash"), {"val": dash})
            ln.append(d)
    sh.shadow.inherit = False
    return sh


def wedge(sl, cx, cy, r, a0, a1, fill):
    """中心(cx,cy)から半径r・角度a0→a1（度、上=90）の扇形をフリーフォームで描く。"""
    e = 914400
    pts = []
    steps = 24
    for i in range(steps + 1):
        a = math.radians(a0 + (a1 - a0) * i / steps)
        pts.append((int((cx + r * math.cos(a)) * e), int((cy - r * math.sin(a)) * e)))
    fb = sl.shapes.build_freeform(Emu(int(cx * e)), Emu(int(cy * e)), scale=1.0)
    fb.add_line_segments([(Emu(x), Emu(y)) for x, y in pts], close=True)
    sh = fb.convert_to_shape()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def chip(sl, cx, cy, w, h, text, fill, size=11):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - w / 2),
                             Inches(cy - h / 2), Inches(w), Inches(h))
    sh.adjustments[0] = 0.5
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    put_text(sh, text, size, WHITE, bold=True)
    return sh


def arrow_to_center(sl, cx, cy, theta_deg, r_out, r_in, color, w_pt=0.24):
    """角度theta（度、上=90）の方向から中心へ向かう矢印。"""
    th = math.radians(theta_deg)
    rm = (r_out + r_in) / 2
    length = r_out - r_in
    px = cx + rm * math.cos(th)
    py = cy - rm * math.sin(th)
    sh = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(px - length / 2),
                             Inches(py - w_pt / 2), Inches(length), Inches(w_pt))
    sh.adjustments[0] = 0.55
    sh.adjustments[1] = 0.55
    sh.rotation = math.degrees(math.atan2(math.sin(th), -math.cos(th))) % 360
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def ripples(sl, cx, cy, theta_deg, r_at, color, n=3, r0=0.10, dr=0.11):
    """音源位置に波紋（同心円アウトライン、外ほど淡い）。"""
    th = math.radians(theta_deg)
    px = cx + r_at * math.cos(th)
    py = cy - r_at * math.sin(th)
    for i in range(n):
        c = blend(color, 0.25 + 0.28 * i)
        oval(sl, px, py, r0 + dr * i, r0 + dr * i, fill=None, line=c, line_w=1.4)
    oval(sl, px, py, 0.035, 0.035, fill=color)


def person(sl, cx, cy):
    oval(sl, cx, cy - 0.10, 0.085, 0.085, fill=INK)
    body = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - 0.14),
                               Inches(cy + 0.005), Inches(0.28), Inches(0.19))
    body.adjustments[0] = 0.5
    body.fill.solid()
    body.fill.fore_color.rgb = INK
    body.line.fill.background()
    body.shadow.inherit = False
    tb = sl.shapes.add_textbox(Inches(cx - 0.5), Inches(cy + 0.20), Inches(1.0), Inches(0.24))
    put_text(tb, "ユーザー", 9.5, INK2, bold=True)


def note(sl, x, y, text, size=10.5, color=MUTED, bold=False, align=PP_ALIGN.LEFT, w=5.0):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.28))
    put_text(tb, text, size, color, bold=bold, align=align)
    return tb


# 音源の配置（角度: 上=90）
SOURCES = [
    ("サイレン", 90, RED),
    ("緊急車両", 210, RED),
    ("車", 262, INK2),
    ("自転車", 328, AMBER),
]

CX, CY, ROUT = 6.667, 3.6, 2.05

# ---------- 案A: レーダー型（リング＋扇＋矢印＋波紋） ----------
sl = prs.slides.add_slide(BLANK)
note(sl, 0.4, 0.25, "案A レーダー型（矢印＋波紋）— スライド上で Ctrl+A → コピーして貼り付け → Ctrl+G でグループ化", 11, MUTED)

for rr, dash in [(0.85, None), (1.45, None), (2.05, "dash")]:
    oval(sl, CX, CY, rr, rr, fill=None, line=HAIR, line_w=1.0, dash=dash)
wedge(sl, CX, CY, 2.05, 58, 122, VISION)
note(sl, CX - 2.45, CY - 1.80, "前方視野（見える範囲）", 9.5, MUTED, align=PP_ALIGN.RIGHT, w=1.7)

for name, th, col in SOURCES:
    arrow_to_center(sl, CX, CY, th, 1.85, 1.05, col)
    ripples(sl, CX, CY, th, 2.02, col, n=2, r0=0.09, dr=0.09)
    lx = CX + 2.55 * math.cos(math.radians(th))
    ly = CY - 2.55 * math.sin(math.radians(th))
    chip(sl, lx, ly, 1.15, 0.36, name, col)
person(sl, CX, CY)

note(sl, CX - 2.6, CY + 2.78, "視覚がカバーできるのは前方の扇だけ", 10.5, INK2, w=5.2, align=PP_ALIGN.CENTER)
note(sl, CX - 2.6, CY + 3.05, "音の危険は全方位から届く — その察知は聴覚頼み", 10.5, RED, bold=True, w=5.2, align=PP_ALIGN.CENTER)

# ---------- 案B: 波紋型（矢印なしミニマル） ----------
sl = prs.slides.add_slide(BLANK)
note(sl, 0.4, 0.25, "案B 波紋型（音の広がりだけで見せるミニマル版）", 11, MUTED)

oval(sl, CX, CY, ROUT, ROUT, fill=None, line=HAIR, line_w=1.25, dash="dash")
wedge(sl, CX, CY, ROUT, 58, 122, VISION)
note(sl, CX - 2.45, CY - 1.80, "前方視野（見える範囲）", 9.5, MUTED, align=PP_ALIGN.RIGHT, w=1.7)

for name, th, col in SOURCES:
    ripples(sl, CX, CY, th, 1.55, col, n=3, r0=0.11, dr=0.13)
    lx = CX + 2.45 * math.cos(math.radians(th))
    ly = CY - 2.45 * math.sin(math.radians(th))
    chip(sl, lx, ly, 1.15, 0.36, name, col)
person(sl, CX, CY)

note(sl, CX - 2.6, CY + 2.78, "見えるのは前方だけ／音は全方位に広がって届く", 10.5, INK2, w=5.2, align=PP_ALIGN.CENTER)

OUT = r"C:/Users/satos/research/outdoor_seld_e2e/md/seminar/図_研究背景_差し替え案_2026-08-13.pptx"
prs.save(OUT)
print("saved:", OUT)
