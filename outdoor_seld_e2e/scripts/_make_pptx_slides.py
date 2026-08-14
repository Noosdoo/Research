# -*- coding: utf-8 -*-
"""中間発表スライドの.pptx生成（本スクリプトと出力pptxが正本。
HTML版=中間発表スライド_2026-08-10.htmlは2026-08-10時点の旧版で監査修正
未反映のため参照禁止）。

デザイン: 道路標識アイデンティティ（標識紺×警戒黄）・メイリオ・16:9。
図はすべてPowerPointネイティブ図形（後から編集可能）。
使い方: python scripts/_make_pptx_slides.py
出力: md/seminar/中間発表スライド_2026-08-10.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "md" / "seminar" / "中間発表スライド_2026-08-10.pptx"

INK = RGBColor(0x16, 0x23, 0x3A)
INK2 = RGBColor(0x3A, 0x46, 0x58)
MUTED = RGBColor(0x66, 0x70, 0x7F)
PAPER = RGBColor(0xF7, 0xF7, 0xF3)
PAPER2 = RGBColor(0xEF, 0xEF, 0xEA)
HAIR = RGBColor(0xD9, 0xDA, 0xD2)
AMBER = RGBColor(0xE8, 0xA2, 0x00)
AMBERD = RGBColor(0xB3, 0x7E, 0x00)
RED = RGBColor(0xC4, 0x43, 0x2B)
GREEN = RGBColor(0x2F, 0x7D, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHIP = RGBColor(0xE4, 0xE5, 0xDE)
RED22 = RGBColor(0xEC, 0xCF, 0xC7)   # 意味色を紙に22%載せた近似
AMB18 = RGBColor(0xF4, 0xE8, 0xC7)
GRN10 = RGBColor(0xE3, 0xEB, 0xE6)
WAVE = RGBColor(0xD3, 0xD5, 0xD5)
HILITE = RGBColor(0xFF, 0xF8, 0xE8)

STEPS = ["背景・目的", "基礎知識", "関連研究", "提案手法", "検証", "今後・まとめ"]
W, H = 960, 540   # pt

prs = Presentation()
prs.slide_width = Pt(W)
prs.slide_height = Pt(H)
BLANK = prs.slide_layouts[6]


def _jp(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", "メイリオ")


def style(run, size, bold=False, color=INK):
    f = run.font
    f.name = "Meiryo"
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    _jp(run)


def box(sl, x, y, w, h):
    tb = sl.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tb.text_frame.word_wrap = True
    return tb


def para(tf, runs, size=15, bold=False, color=INK2, align=PP_ALIGN.LEFT,
         before=0, after=6, line=None, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if line:
        p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for t, kw in runs:
        r = p.add_run()
        r.text = t
        style(r, kw.get("size", size), kw.get("bold", bold), kw.get("color", color))
    return p


def rect(sl, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
         dash=None):
    sp = sl.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
        if dash:
            ln = sp.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    sp.text_frame.word_wrap = True
    return sp


def conn(sl, x1, y1, x2, y2, color, w=2.0, dash=None, arrow=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    if arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return c


def new_slide():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, -2, -2, W + 4, H + 4, fill=PAPER)
    return sl


PAGENO = [0]


def footer(sl, active):
    PAGENO[0] += 1
    conn(sl, 0, 492, W, 492, HAIR, 1.0)
    x = 54
    for i, s in enumerate(STEPS):
        wd = 14 + len(s) * 11.5
        c = rect(sl, x, 504, wd, 22, fill=(INK if i == active else CHIP),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        tf.margin_top = tf.margin_bottom = Pt(1)
        para(tf, s, size=10.5, bold=(i == active),
             color=(WHITE if i == active else MUTED), align=PP_ALIGN.CENTER, after=0,
             first=True)
        x += wd + 8
    tb = box(sl, W - 140, 504, 90, 22)
    para(tb.text_frame, f"{PAGENO[0]} / 18", size=10.5, color=MUTED,
         align=PP_ALIGN.RIGHT, after=0, first=True)


def header(sl, kicker, title):
    tb = box(sl, 54, 30, W - 108, 24)
    para(tb.text_frame, kicker, size=11, bold=True, color=AMBERD, after=0, first=True)
    tb2 = box(sl, 54, 52, W - 108, 44)
    para(tb2.text_frame, title, size=27, bold=True, color=INK, after=0, first=True)


def bullets(sl, x, y, w, items, size=15, gap=8):
    tb = box(sl, x, y, w, 400)
    tf = tb.text_frame
    first = True
    for it in items:
        mark = it.get("mark", "■ ")
        mc = it.get("mcolor", INK)
        runs = [(mark, {"color": mc, "bold": True, "size": it.get("size", size)})]
        for t, kw in it["runs"]:
            kw.setdefault("size", it.get("size", size))
            runs.append((t, kw))
        para(tf, runs, size=it.get("size", size), color=INK2, after=gap,
             line=1.28, first=first)
        first = False
    return tb


def card(sl, x, y, w, h, title=None, border=HAIR, lw=1.0):
    c = rect(sl, x, y, w, h, fill=WHITE, line=border, lw=lw)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if title:
        para(tf, title, size=14, bold=True, color=INK, after=4, first=True)
    return c


def stripe(sl, y, h=8):
    x = 0.0
    i = 0
    while x < W:
        rect(sl, x, y, 14, h, fill=(AMBER if i % 2 == 0 else INK),
             shape=MSO_SHAPE.PARALLELOGRAM)
        x += 12
        i += 1


def label(sl, x, y, w, text, size=10.5, color=MUTED, align=PP_ALIGN.CENTER,
          bold=False):
    tb = box(sl, x, y, w, 18)
    tb.text_frame.word_wrap = False
    para(tb.text_frame, text, size=size, bold=bold, color=color, align=align,
         after=0, first=True)


def tagchip(sl, x, y, text, color):
    wd = 16 + len(text) * 11
    c = rect(sl, x, y, wd, 20, fill=None, line=color, lw=1.2)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    para(tf, text, size=10, bold=True, color=color, align=PP_ALIGN.CENTER,
         after=0, first=True)
    return wd


# ============ 1 表紙 ============
sl = new_slide()
rect(sl, -2, -2, W + 4, 74, fill=INK)
tb = box(sl, 54, 24, 500, 26)
para(tb.text_frame, "2026 夏ゼミ ・ 中間発表", size=11.5, bold=True,
     color=RGBColor(0xC8, 0xCE, 0xDA), after=0, first=True)
for r_ in (26, 56, 90, 126):
    o = rect(sl, 750 - r_, 170 - r_, r_ * 2, r_ * 2, fill=None, line=WAVE, lw=1.5,
             shape=MSO_SHAPE.OVAL)
rect(sl, 742, 162, 16, 16, fill=AMBER, shape=MSO_SHAPE.OVAL)
tb = box(sl, 54, 130, 700, 200)
para(tb.text_frame, "SELDモデルを用いた", size=40, bold=True, color=INK,
     line=1.3, first=True)
para(tb.text_frame, "難聴歩行者向け屋外危険音", size=40, bold=True, color=INK, line=1.3)
para(tb.text_frame, "通知システムの構築", size=40, bold=True, color=INK, line=1.3)
tb = box(sl, 54, 360, 400, 40)
para(tb.text_frame, [("B4　", {"size": 16, "color": INK2}),
                     ("松本 鋭", {"size": 20, "bold": True, "color": INK})],
     after=0, first=True)
stripe(sl, 470)
footer(sl, -1)

# ============ 2 研究背景 ============
sl = new_slide()
header(sl, "BACKGROUND", "研究背景 — 歩行者の安全は聴覚に支えられている")
bullets(sl, 54, 120, 440, [
    {"runs": [("歩行者は接近車・サイレン・踏切などの", {}),
              ("危険察知を聴覚に依存", {"bold": True, "color": INK}), ("している", {})]},
    {"runs": [("難聴者はその情報が欠落し、屋外歩行の事故リスクが高い", {})]},
    {"mark": "・", "mcolor": MUTED,
     "runs": [("イヤホン歩行者・高齢者にも共通する課題", {"color": MUTED})], "size": 13},
    {"mcolor": AMBER,
     "runs": [("機械の耳は疲れず全方位を常時監視でき、可聴域の外の信号も扱える\n→「耳の代わり」に", {}),
              ("とどまらない支援", {"bold": True, "color": AMBERD}),
              ("まで狙える", {})]},
])
fx, fy = 560, 112   # 図の原点
rect(sl, fx, fy, 340, 360 - 10, fill=WHITE, line=HAIR)
tri = rect(sl, fx + 100, fy + 18, 140, 96, fill=RGBColor(0xE2, 0xE4, 0xE6),
           shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
tri.rotation = 180
label(sl, fx + 110, fy + 34, 120, "前方視野", size=10.5)
rect(sl, fx + 156, fy + 118, 28, 28, fill=INK, shape=MSO_SHAPE.OVAL)
label(sl, fx + 120, fy + 150, 100, "ユーザー", size=10.5, color=INK2)
conn(sl, fx + 40, fy + 300, fx + 150, fy + 152, RED, 2.5, arrow=True)
conn(sl, fx + 170, fy + 316, fx + 170, fy + 152, RED, 2.5, arrow=True)
conn(sl, fx + 300, fy + 300, fx + 190, fy + 152, RED, 2.5, arrow=True)
conn(sl, fx + 120, fy + 24, fx + 160, fy + 112, AMBER, 2.5, arrow=True)
label(sl, fx + 8, fy + 6, 180, "サイレン（前方も対象）", size=10, color=AMBERD,
      bold=True, align=PP_ALIGN.LEFT)
label(sl, fx + 4, fy + 306, 90, "緊急車両", size=10.5, color=INK, bold=True)
label(sl, fx + 126, fy + 320, 90, "車・バイク", size=10.5, color=INK, bold=True)
label(sl, fx + 246, fy + 306, 110, "キックボード", size=10.5, color=INK, bold=True)
footer(sl, 0)

# ============ 3 研究目的 ============
sl = new_slide()
header(sl, "OBJECTIVE", "研究目的 — 種類・方向・距離を、危険度に変えて届ける")
cw = (W - 108 - 40) / 3
for i, (t, body) in enumerate([
    ("検出対象＝8クラス",
     "サイレン（救急・パト・消防）／クラクション／バック音／自転車ベル／車（EV・大型含む）／踏切＋列車／キックボード／バイク"),
    ("伝え方",
     "音の種類・方向・距離を推定し、危険度3段階（至近警告／注意／抑制）を将来的には首元の振動で伝える構成を想定。行動の判断は本人に委ねる"),
    ("前提",
     "視覚は健常と仮定。検出・通知は全方位——絞るのは方向でなく「視覚で代替できる情報」。前方でも音でしか得られない情報（サイレン・踏切警報・EV）は対象。第1対象は難聴者"),
]):
    c = card(sl, 54 + i * (cw + 20), 118, cw, 190, title=t)
    para(c.text_frame, body, size=13, color=INK2, line=1.3)
rect(sl, 54, 330, 6, 120, fill=AMBER)
c = rect(sl, 60, 330, W - 114, 120, fill=WHITE, line=HAIR)
tf = c.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(10)
para(tf, "このシステムができること", size=14, bold=True, color=INK, after=4,
     first=True)
para(tf, [("危険音の", {}), ("種類・方向・距離", {"bold": True, "color": INK}),
          ("を同時に推定し、安全な車は鳴らさず、", {}),
          ("危険な接近だけを3段階で通知する（提示は首元振動を将来構想）", {"bold": True, "color": AMBERD}),
          ("。音源仕様は日本の法規・省令・実測研究に根拠を持たせている", {})],
     size=14, color=INK2, line=1.35)
footer(sl, 0)

# ============ 4 SELDとは ============
sl = new_slide()
header(sl, "PRELIMINARIES", "SELDとは — 画像でいう物体検出の、音・方向版")
tb = box(sl, 54, 104, W - 108, 30)
para(tb.text_frame,
     [("SELD", {"bold": True, "color": INK}),
      ("（Sound Event Localization and Detection）＝音の", {}),
      ("種類", {"bold": True, "color": INK}), ("と", {}),
      ("時間", {"bold": True, "color": INK}), ("と", {}),
      ("方向", {"bold": True, "color": INK}),
      ("を、複数の音源について同時に当てるタスク", {})],
     size=15, color=INK2, after=0, first=True)
gx, gy, gw, gh = 54, 148, 410, 240
rect(sl, gx, gy, gw, gh, fill=WHITE, line=HAIR)
conn(sl, gx + 24, gy + 180, gx + gw - 24, gy + 180, MUTED, 1.5)
label(sl, gx + gw - 90, gy + 186, 70, "時間 →", size=10, align=PP_ALIGN.RIGHT)
for (bx, by, bw, col, t) in [(60, 40, 140, RED, "車"),
                             (150, 82, 170, GREEN, "サイレン"),
                             (40, 124, 120, INK, "自転車ベル")]:
    b = rect(sl, gx + bx, gy + by, bw, 26, fill=col)
    tf = b.text_frame
    tf.margin_left = Pt(6)
    tf.margin_top = Pt(1)
    para(tf, t, size=11, color=WHITE, after=0, first=True)
label(sl, gx, gy + gh - 26, gw, "SED：「今、何が鳴っているか」（種類×時間）", size=11)
gx2 = 500
rect(sl, gx2, gy, 410, gh, fill=WHITE, line=HAIR)
cxc, cyc = gx2 + 205, gy + 108
rect(sl, cxc - 66, cyc - 66, 132, 132, fill=None, line=HAIR, lw=1.5,
     shape=MSO_SHAPE.OVAL)
rect(sl, cxc - 8, cyc - 8, 16, 16, fill=INK, shape=MSO_SHAPE.OVAL)
conn(sl, cxc, cyc, cxc + 58, cyc - 44, AMBER, 3.0, arrow=True)
label(sl, cxc + 40, cyc - 74, 150, "右斜め前 30°", size=11, color=AMBERD,
      align=PP_ALIGN.LEFT, bold=True)
label(sl, cxc - 20, cyc - 96, 40, "前", size=10)
label(sl, gx2, gy + gh - 26, 410, "SSL：「どっちから鳴っているか」（方位・仰角）", size=11)
tb = box(sl, 54, 402, W - 108, 30)
para(tb.text_frame,
     [("本研究はここに", {}), ("音源距離推定（SDE）を拡張", {"bold": True, "color": AMBERD}),
      ("し、種類＋方向＋", {}), ("距離", {"bold": True, "color": INK}),
      ("の同時推定にする", {})], size=15, color=INK2, after=0, first=True)
footer(sl, 1)

# ============ 5 PSELDNets ============
sl = new_slide()
header(sl, "RELATED WORK 1/2", "PSELDNets（2024）— SELDの基盤モデル")
bullets(sl, 54, 116, 400, [
    {"runs": [("大規模合成データで事前学習した", {}),
              ("SELD基盤モデル", {"bold": True, "color": INK}),
              ("（1,167時間・170クラス、主要3データセットでSOTA）", {})]},
    {"runs": [("省データ適応手法 AdapterBit を導入", {})]},
    {"runs": [("出力は multi-ACCDOA（同一クラス最大3つ同時検出）", {})]},
    {"mcolor": AMBER,
     "runs": [("本研究はこの基盤に", {}),
              ("距離推定ヘッドを拡張", {"bold": True, "color": AMBERD}),
              ("し、屋外・移動音源・8クラスへ応用", {})]},
], size=14)
fx, fy = 490, 112
rect(sl, fx, fy, 420, 300, fill=WHITE, line=HAIR)
rect(sl, fx + 12, fy + 26, 112, 78, fill=PAPER2, line=HAIR)
r1 = rect(sl, fx + 18, fy + 20, 112, 78, fill=WHITE, line=HAIR)
tf = r1.text_frame
tf.margin_top = Pt(6)
para(tf, "大規模合成データ", size=11, bold=True, color=INK,
     align=PP_ALIGN.CENTER, after=2, first=True)
para(tf, "1,167時間\n170クラス", size=10, color=INK2, align=PP_ALIGN.CENTER)
conn(sl, fx + 134, fy + 58, fx + 168, fy + 58, INK2, 2.0, arrow=True)
label(sl, fx + 122, fy + 30, 64, "事前学習", size=9)
r2 = rect(sl, fx + 172, fy + 28, 100, 62, fill=INK)
tf = r2.text_frame
tf.margin_top = Pt(8)
para(tf, "SELD\n基盤モデル", size=12, bold=True, color=WHITE,
     align=PP_ALIGN.CENTER, after=0, first=True)
conn(sl, fx + 276, fy + 58, fx + 310, fy + 58, INK2, 2.0, arrow=True)
label(sl, fx + 254, fy + 14, 100, "少データ適応", size=9)
label(sl, fx + 254, fy + 26, 100, "(AdapterBit)", size=9)
r3 = rect(sl, fx + 314, fy + 18, 94, 84, fill=HILITE, line=AMBER, lw=2.0)
tf = r3.text_frame
tf.margin_top = Pt(4)
para(tf, "本研究", size=11.5, bold=True, color=AMBERD,
     align=PP_ALIGN.CENTER, after=2, first=True)
para(tf, "屋外・移動音源\n8クラス", size=9.5, color=INK, align=PP_ALIGN.CENTER,
     after=1)
para(tf, "＋距離ヘッド", size=10, bold=True, color=AMBERD, align=PP_ALIGN.CENTER)
conn(sl, fx + 360, fy + 104, fx + 360, fy + 140, AMBERD, 1.5, dash="dash",
     arrow=True)
rect(sl, fx + 170, fy + 144, 240, 128, fill=WHITE, line=HAIR)
label(sl, fx + 180, fy + 152, 220, "出力ヘッド（multi-ACCDOA）", size=9.5,
      align=PP_ALIGN.LEFT)
for i, t in enumerate(["x", "y", "z"]):
    b = rect(sl, fx + 182 + i * 52, fy + 176, 44, 24, fill=PAPER2, line=HAIR)
    para(b.text_frame, t, size=10.5, color=INK, align=PP_ALIGN.CENTER,
         after=0, first=True)
b = rect(sl, fx + 182 + 3 * 52, fy + 176, 52, 24, fill=HILITE, line=AMBER, lw=1.5)
para(b.text_frame, "距離", size=10.5, bold=True, color=AMBERD,
     align=PP_ALIGN.CENTER, after=0, first=True)
tb = box(sl, fx + 180, fy + 208, 224, 56)
para(tb.text_frame, "方向(x,y,z)に距離の1軸を追加。\nSELD部分は等価性テストで不変を担保",
     size=9.5, color=MUTED, after=0, line=1.25, first=True)
label(sl, fx, fy + 278, 420, "基盤モデルの転移＋距離ヘッド拡張（黄＝本研究の追加）",
      size=10)
footer(sl, 2)

# ============ 6 DynamicSound ============
sl = new_slide()
header(sl, "RELATED WORK 2/2", "DynamicSound（2026）— 屋外音響の物理シミュレータ")
bullets(sl, 54, 112, 400, [
    {"runs": [("室内RIRツールの限界を屋外向けに解いたオープンソース", {})]},
    {"runs": [("自由音場・任意マイクアレイ・音源の連続3D移動", {})]},
    {"runs": [("屋外物理を明示的にモデル化：伝搬遅延・ドップラー・距離減衰・大気吸収・一次反射", {})]},
], size=14)
fx, fy = 490, 104
rect(sl, fx, fy, 420, 216, fill=WHITE, line=HAIR)
for r_ in (14, 26, 40):
    rect(sl, fx + 64 - r_, fy + 52 - r_, r_ * 2, r_ * 2, fill=None, line=WAVE,
         lw=1.5, shape=MSO_SHAPE.OVAL)
rect(sl, fx + 57, fy + 45, 14, 14, fill=AMBER, shape=MSO_SHAPE.OVAL)
conn(sl, fx + 80, fy + 30, fx + 118, fy + 20, AMBERD, 2.0, arrow=True)
label(sl, fx + 84, fy + 2, 200, "移動する音源（3D・連続）", size=9.5, color=AMBERD,
      align=PP_ALIGN.LEFT, bold=True)
conn(sl, fx + 78, fy + 60, fx + 330, fy + 128, INK2, 2.0, arrow=True)
for i, t in enumerate(["伝搬遅延", "ドップラー", "距離減衰", "大気吸収"]):
    b = rect(sl, fx + 66 + i * 64, fy + 66 + i * 16, 62, 20, fill=WHITE, line=HAIR)
    para(b.text_frame, t, size=9, color=INK, align=PP_ALIGN.CENTER, after=0,
         first=True)
conn(sl, fx + 16, fy + 172, fx + 404, fy + 172, MUTED, 2.0)
label(sl, fx + 16, fy + 178, 50, "地面", size=9, align=PP_ALIGN.LEFT)
conn(sl, fx + 74, fy + 62, fx + 200, fy + 172, MUTED, 1.5, dash="dash")
conn(sl, fx + 200, fy + 172, fx + 336, fy + 140, MUTED, 1.5, dash="dash",
     arrow=True)
label(sl, fx + 160, fy + 180, 90, "一次反射", size=9)
rect(sl, fx + 330, fy + 112, 72, 34, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
for dx, dy in ((14, 13), (30, 7), (46, 13), (30, 20)):
    rect(sl, fx + 330 + dx, fy + 112 + dy, 7, 7, fill=PAPER, shape=MSO_SHAPE.OVAL)
label(sl, fx + 310, fy + 150, 120, "4chマイクアレイ", size=9.5, color=INK, bold=True)
rect(sl, 54, 340, 6, 108, fill=AMBER)
c = rect(sl, 60, 340, W - 114, 108, fill=WHITE, line=HAIR)
tf = c.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(8)
para(tf, "それでも残る空白＝本研究の位置", size=14, bold=True, color=INK, after=3,
     first=True)
para(tf, [("歩行者装着 × 複数危険クラス同時 × 距離まで", {"bold": True, "color": INK}),
          ("を屋外で扱った研究は調査した範囲で見当たらない。分業型（分類CNN＋信号処理DOA）は複数音源の同時発生で対応付けが壊れやすい → 同時推定型SELDを採用",
           {})], size=13, color=INK2, line=1.3)
footer(sl, 2)

# ============ 7 提案手法1 ============
sl = new_slide()
header(sl, "PROPOSED METHOD 1/2", "提案手法 — 合成データで学習を成立させる")
cw = (W - 108 - 40) / 3
for i, (t, body) in enumerate([
    ("物理シミュレーション", "伝搬遅延／ドップラー／距離減衰／大気吸収／地面反射を自作実装（物理検証済み）"),
    ("4ch音＋方向・距離ラベル", "同一経路から同時に出力（FOA）。生成誤差ゼロの教師ラベル（シミュレータ内部基準・実世界とのモデル化誤差は残る）"),
    ("学習・評価", "PSELDNetsをファインチューニング。学習10,200＋評価1,800クリップ"),
]):
    c = card(sl, 54 + i * (cw + 20), 116, cw, 130, title=t)
    para(c.text_frame, body, size=12.5, color=INK2, line=1.3)
    if i < 2:
        conn(sl, 54 + (i + 1) * (cw + 20) - 16, 181, 54 + (i + 1) * (cw + 20) - 2,
             181, INK2, 2.0, arrow=True)
bullets(sl, 54, 268, W - 108, [
    {"runs": [("なぜ合成か：① 屋外・歩行者視点・距離つきの学習データが調査した範囲で存在しない　② 物理を自分で持てば1つずつ外せる＝", {}),
              ("ablationができる土俵", {"bold": True, "color": INK}), ("になる", {})]},
    {"runs": [("音源の音量・周波数は", {}),
              ("日本の法規・省令・実測研究で出典を層別管理", {"bold": True, "color": INK}),
              ("（規則値と設計値を混同しない）", {})]},
    {"mcolor": AMBER,
     "runs": [("空間手がかりが理想的過ぎる限界 → ", {}),
              ("9月に実録で検証", {"bold": True, "color": AMBERD}),
              ("（測定方法は後述）", {})]},
], size=14)
footer(sl, 3)

# ============ 8 提案手法2 ============
sl = new_slide()
header(sl, "PROPOSED METHOD 2/2", "2層の仕組み — 推定距離で「いつ・どの強さで」を決める")
x0, cw2 = 54, 400
c = card(sl, x0, 110, cw2, 64, title="知覚層（SELD＋SDE）")
para(c.text_frame, "8クラスの検出・方向・距離を0.1秒ごとに推定", size=12.5, color=INK2)
label(sl, x0, 176, cw2, "▼", size=11)
c = card(sl, x0, 194, cw2, 118, title="通知層（規則v3.4）— 推定距離で出し分け", border=AMBER,
         lw=1.5)
para(c.text_frame, [("至近警告", {"bold": True, "color": RED}),
                    ("　推定距離≤1.5mが2フレーム連続（同一物体は方位連結±60°の近似）", {})],
     size=12, color=INK2, line=1.3)
para(c.text_frame, [("注意", {"bold": True, "color": AMBERD}), ("　≤3.0m　　", {}),
                    ("抑制", {"bold": True, "color": GREEN}),
                    ("　>3.2m＝安全な車は鳴らさない", {})], size=12, color=INK2,
     line=1.3)
label(sl, x0, 314, cw2, "▼", size=11)
c = card(sl, x0, 332, cw2, 62, title="首元振動デバイス（今後）")
para(c.text_frame, "方向×振動でユーザに伝える", size=12.5, color=INK2)
tb = box(sl, x0, 400, cw2, 60)
para(tb.text_frame, "知覚層が「何が・どこから・どれだけ近くに」を推定し、通知層が「いつ・どの強さで」伝えるかを決める",
     size=11.5, color=MUTED, line=1.3, after=0, first=True)
cx, cy = 700, 260
rect(sl, cx - 150, cy - 150, 300, 300, fill=GRN10, line=GREEN, lw=1.5,
     shape=MSO_SHAPE.OVAL, dash="dash")
rect(sl, cx - 96, cy - 96, 192, 192, fill=AMB18, line=AMBERD, lw=1.5,
     shape=MSO_SHAPE.OVAL)
rect(sl, cx - 48, cy - 48, 96, 96, fill=RED22, line=RED, lw=1.5,
     shape=MSO_SHAPE.OVAL)
rect(sl, cx - 8, cy - 8, 16, 16, fill=INK, shape=MSO_SHAPE.OVAL)
label(sl, cx - 80, cy - 40, 160, "≤1.5m 至近警告", size=11, color=RED, bold=True)
label(sl, cx - 80, cy - 84, 160, "≤3.0m 注意", size=11, color=AMBERD, bold=True)
label(sl, cx - 80, cy - 136, 160, ">3.2m 抑制", size=11, color=GREEN, bold=True)
label(sl, cx - 170, cy + 158, 340, "全部鳴らすと使えない → 通知の頻度そのものを設計対象に",
      size=10.5)
footer(sl, 3)

# ============ 9 データセットと評価手順（2026-08-13挿入・Sol指摘対応） ============
sl = new_slide()
header(sl, "VALIDATION SETUP", "検証の前に — データセットと評価手順")
rows = [
    ("区分", "中身", "本研究での役割"),
    ("事前学習", "1,167時間・170クラス（PSELDNets付属）", "出発点の基盤モデル（既存を利用）"),
    ("学習（微調整）", "本研究の合成10,200クリップ（8クラス・約28時間）",
     "距離ヘッド込みで全体をファインチューニング"),
    ("検証", "学習と同一方式の別クリップ", "モデル選択・通知閾値の決定（発表数値には不使用）"),
    ("評価", "同一設計・新乱数の1,800クリップ（学習・検証に未使用）",
     "最終結果の測定 — 1回だけ採点"),
]
gt = sl.shapes.add_table(5, 3, Pt(54), Pt(128), Pt(W - 108), Pt(252)).table
gt.columns[0].width = Pt(150)
gt.columns[1].width = Pt(372)
gt.columns[2].width = Pt(330)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gt.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = (PAPER2 if ri == 0 else
                                    (HILITE if ri == 4 else WHITE))
        cell.margin_top = cell.margin_bottom = Pt(4)
        tf = cell.text_frame
        tf.word_wrap = True
        para(tf, val, size=12.5, bold=(ri == 0 or ci == 0),
             color=INK if (ri and ci == 0) else INK2,
             align=PP_ALIGN.LEFT, after=0, first=True)
b = rect(sl, 54, 410, W - 108, 54, fill=HILITE)
para(b.text_frame,
     [("基準を先に決める → 未使用データを新たに生成 → 1回だけ評価", {"bold": True, "color": INK}),
      ("　＝ 答案を見てから基準を変えられない手順", {"color": INK2})],
     size=13.5, after=0, first=True)
footer(sl, 4)

# ============ 10 検証1 ============
sl = new_slide()
header(sl, "RESULTS 1/2", "検証 — 8クラスの検出・方向・距離")
tb = box(sl, 54, 98, W - 108, 30)
para(tb.text_frame,
     [("評価データ：", {}), ("学習に未使用の1,800クリップ",
                        {"bold": True, "color": INK}),
      ("（合成・種類/方向/距離の正解つき）で", {}),
      ("1回だけ", {"bold": True, "color": INK}), ("評価", {})],
     size=13.5, color=INK2, line=1.3, after=0, first=True)
rows = [
    ("クラス", "検出率（可聴フレーム）", "方向誤差（中央値）"),
    ("サイレン（救急・パト・消防）", "99.5%", "3.2°"),
    ("クラクション", "99.2%", "2.2°"),
    ("バック音", "97.2%", "5.1°"),
    ("自転車ベル", "97.6%", "4.1°"),
    ("車（EV・大型含む）", "99.4%", "2.0°"),
    ("踏切・列車", "100.0%", "3.5°"),
    ("キックボード", "97.3%", "1.0°"),
    ("バイク", "99.8%", "1.0°"),
]
gt = sl.shapes.add_table(9, 3, Pt(54), Pt(132), Pt(620), Pt(292)).table
gt.columns[0].width = Pt(280)
gt.columns[1].width = Pt(180)
gt.columns[2].width = Pt(160)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gt.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER2 if ri == 0 else WHITE
        cell.margin_top = cell.margin_bottom = Pt(2)
        tf = cell.text_frame
        tf.word_wrap = True
        para(tf, val, size=12 if ri else 11,
             bold=(ri == 0 or ci == 0),
             color=INK if (ri and ci == 0) else INK2,
             align=(PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT),
             after=0, first=True)
c = card(sl, 700, 132, 206, 292, title="距離推定（車）")
para(c.text_frame, [("至近（≤5m）の\n中央絶対誤差", {})], size=12, color=INK2,
     line=1.3)
para(c.text_frame, [("0.21m", {"bold": True, "color": INK, "size": 26})],
     size=26, after=8)
para(c.text_frame, "静音EVだけの場面でも検出率 99.0%", size=11.5, color=INK2,
     line=1.3)
tb = box(sl, 54, 434, W - 108, 30)
para(tb.text_frame,
     "可聴フレーム＝背景騒音に対してSNR≥0dBの区間。検出率はその全フレームに対する割合、方向誤差は全検出の中央値",
     size=11.5, color=MUTED, after=0, first=True)
footer(sl, 4)

# ============ 10 検証2 ============
sl = new_slide()
header(sl, "RESULTS 2/2", "検証 — 遠方の危険への「気づき」を距離で測る")
fx, fy = 54, 108
rect(sl, fx, fy, W - 108, 150, fill=WHITE, line=HAIR)
ax0, ax1, ay = fx + 40, fx + W - 108 - 40, fy + 92
conn(sl, ax0, ay, ax1, ay, INK, 2.0)


def dpos(m):
    return ax0 + (ax1 - ax0) * m / 300.0


for m in (0, 100, 200, 300):
    conn(sl, dpos(m), ay - 5, dpos(m), ay + 5, MUTED, 1.0)
    label(sl, dpos(m) - 30, ay + 10, 60, f"{m}m", size=10)
rect(sl, ax0 - 7, ay - 7, 14, 14, fill=INK, shape=MSO_SHAPE.OVAL)
label(sl, ax0 - 34, ay - 30, 60, "歩行者", size=10, color=INK2)
conn(sl, dpos(208), ay, dpos(208), ay - 52, RED, 2.5)
label(sl, dpos(208) - 70, ay - 74, 140, "中央 208m", size=14, color=RED, bold=True)
conn(sl, dpos(279), ay, dpos(279), ay - 38, AMBERD, 2.0)
label(sl, dpos(279) - 60, ay - 58, 120, "p90 279m", size=11, color=AMBERD, bold=True)
conn(sl, dpos(303), ay, dpos(303), ay - 24, MUTED, 1.5)
label(sl, dpos(303) - 90, ay - 44, 110, "最大 303m", size=10, color=MUTED,
      align=PP_ALIGN.RIGHT)
label(sl, fx + 12, fy + 122, W - 132,
      "サイレン初検知距離（100〜300mの遠方配置・未検知ゼロ・開発評価セット）— 常時監視による早期の気づきを距離で定量化",
      size=10.5, align=PP_ALIGN.LEFT)
cw = (W - 108 - 40) / 3
for i, (t, lines) in enumerate([
    ("通知層（車）", [[("1.5m以内まで近づく車に至近警告が出た割合 ", {}),
                  ("69.9%", {"bold": True, "color": INK, "size": 15})],
                 [("3.2mより遠い車を鳴らさなかった割合 ", {}),
                  ("90.4%", {"bold": True, "color": INK, "size": 15})],
                 [("安全な車への誤った至近警告 1.3%", {})]]),
    ("静音の脅威にも", [[("1.5m以内まで近づく", {}), ("キックボード", {"bold": True, "color": INK})],
                 [("への至近警告 ", {}), ("88.7%", {"bold": True, "color": INK, "size": 15})]]),
    ("可聴域の外へ", [[("40kHz超音波", {"bold": True, "color": INK}), ("の検出系を実装", {})],
                [("（仮定受信機・設計音圧の合成試験で10m・実機は9月確認）", {})]]),
]):
    c = card(sl, 54 + i * (cw + 20), 286, cw, 150, title=t)
    for ln in lines:
        para(c.text_frame, ln, size=12.5, color=INK2, line=1.3)
footer(sl, 4)

# ============ 11 ablation ============
sl = new_slide()
header(sl, "NEXT 1/2", "今後の方針① — 物理ablationで問いに答える（9月）")
tb = box(sl, 54, 104, W - 108, 34)
para(tb.text_frame,
     [("動機：どの物理要因が、どの誤り（", {}),
      ("見逃し／誤通知／方向外れ", {"bold": True, "color": INK}),
      ("）を生むのかを分離する", {})], size=15, color=INK2, after=0, first=True)
cw = (W - 108 - 40) / 3
for i, (t, body) in enumerate([
    ("規模", "物理を1つだけ外した版 ×4（ドップラー・大気吸収・幾何減衰・地面反射）"),
    ("測り方", "各版を同一の評価セット＋実環境録音の共通土俵で採点 → 落ちた分＝その物理の大事さ"),
    ("位置づけ", "「合成データのどの物理が実環境性能に必要か」という問いに、切り分け実験で直接答える"),
]):
    c = card(sl, 54 + i * (cw + 20), 152, cw, 160, title=t)
    para(c.text_frame, body, size=13, color=INK2, line=1.35)
footer(sl, 5)

# ============ 12 実録 ============
sl = new_slide()
header(sl, "NEXT 2/2", "今後の方針② — 実録評価の測定方法（9月）")
cw = (W - 108 - 40) / 3
for i, (tag, col, body) in enumerate([
    ("① 統制録音", INK2, "走行車20・弱点20（EV/発進/自転車/見通し不良）・キックボード20。速度×横距離×左右を散らし、幾何を合成シナリオに揃える"),
    ("② 固定・機会", AMBERD, "踏切8＋バック音の統制4＋機会枠8。サイレン・クラクションは遭遇したら他を中断して最優先記録"),
    ("③ 負例", GREEN, "静穏・繁華街・雨上がり・風 計20本＋連続録音100分。誤警告率／時の実世界試験"),
]):
    c = rect(sl, 54 + i * (cw + 20), 108, cw, 120, fill=WHITE, line=HAIR)
    tf = c.text_frame
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(34)
    tagchip(sl, 54 + i * (cw + 20) + 14, 118, tag, col)
    para(tf, body, size=12, color=INK2, line=1.3, first=True)
bullets(sl, 54, 244, W - 108, [
    {"runs": [("機材：4ch FOAマイク（Zoom H3-VR・レンタル）・96kHz・風防。テイク前に方位校正音＋騒音計LAeqを記録", {})]},
    {"runs": [("注釈（車種・目視速度・横距離・風速・LAeq）を正解として", {}),
              ("同じモデル・同じ通知閾値", {"bold": True, "color": INK}),
              ("で比較（正解は実録注釈の粒度・計100本＋静止/歩行対比20本＝計画値）", {})]},
], size=13.5)
rect(sl, 54, 336, 6, 110, fill=AMBER)
c = rect(sl, 60, 336, W - 114, 110, fill=WHITE, line=HAIR)
tf = c.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(8)
para(tf, "サイレンはスピーカ再生で測らない（設計判断）", size=14, bold=True, color=INK,
     after=3, first=True)
para(tf, [("緊急車両との誤認リスクに加え、規定音圧90–120dB@20mに市販スピーカでは届かず「近くの偽物」を測ることになる → ",
           {}), ("機会捕捉が正しい測定法", {"bold": True, "color": INK})],
     size=13, color=INK2, line=1.3)
footer(sl, 5)

# ============ 13 限界 ============
sl = new_slide()
header(sl, "LIMITATIONS", "正直な限界")
items = [
    ("実録", "実環境は未検証。9月に120本（ablation共通の静止100＋静止/歩行対比20）で"
             "測るが、測れるのは「胸部装着・実環境での通知性能」まで"),
    ("至近", "車の至近帯（≤1.5m）は安全上そもそも再現できない — "
             "至近警告の実環境確認は自転車・キックボードでの部分検証にとどまる"),
    ("指標", "総合SELDスコアは事前基準に僅差未達（0.087 vs 0.085）。"
             "実録では0.1秒刻みの正解位置が無いため実測SELDスコアは出さない（通知評価と幾何近似の誤差のみ）"),
    ("通知", "至近警告の到達は約7割。未達の主因は距離推定が2フレーム連続で1.5mに届かないケース"
             "（失敗例は全数分類済み・安全な車への誤帰属はほぼ無し）"),
    ("学習", "キックボードの高騒音下検出は88.7% — 距離重視の学習とのトレードオフが残る"),
    ("設計", "通知の同一物体判定は方位連結（±60°）による近似（完全なトラッキングは将来課題）"),
]
y = 112
for tag, body in items:
    tagchip(sl, 54, y + 2, tag, INK2)
    tb = box(sl, 116, y - 4, W - 170, 56)
    para(tb.text_frame, body, size=13, color=INK2, line=1.25, after=0, first=True)
    y += 58
tb = box(sl, 54, y + 4, W - 108, 30)
para(tb.text_frame, "限界は「どのケースが・何台・なぜ」まで個票で分解済み。"
     "測れないものは測れないと先に宣言してから9月の収録に入る。",
     size=12, color=MUTED, after=0, first=True)
footer(sl, 5)

# ============ 14 まとめ ============
sl = new_slide()
header(sl, "SUMMARY", "まとめ")
bullets(sl, 54, 120, W - 108, [
    {"runs": [("8クラス・", {}), ("距離つき", {"bold": True, "color": INK}),
              ("屋外SELDを日本準拠の合成データで構築し、", {}),
              ("学習に未使用の評価データで性能を確認", {"bold": True, "color": INK}),
              ("（事前設定の目標5項目中4項目を達成）", {})]},
    {"runs": [("通知層が", {}), ("推定距離で3段階", {"bold": True, "color": INK}),
              ("に出し分け — 至近警告到達 ", {}),
              ("69.9%", {"bold": True, "color": INK, "size": 17}),
              ("・安全車の抑制 ", {}),
              ("90.4%", {"bold": True, "color": INK, "size": 17}),
              ("・キックボード ", {}),
              ("88.7%", {"bold": True, "color": INK, "size": 17}),
              ("（すべて合成データでの評価）", {})]},
    {"runs": [("遠方サイレンの初検知 中央 ", {}),
              ("208m", {"bold": True, "color": INK, "size": 30}),
              ("（開発評価）・40kHz検出系も実装（合成試験）", {})]},
    {"runs": [("次：9月の実環境録音——", {}),
              ("静止100本を全ablation版の共通試験", {"bold": True, "color": INK}),
              ("、", {}), ("歩行20本は別枠で評価", {"bold": True, "color": INK}),
              ("——で「どの物理が必要か」と「実環境でどれだけ落ちるか」を測る", {})]},
], size=16, gap=20)
stripe(sl, 456)
footer(sl, 5)

# ============ 15 参考文献 ============
sl = new_slide()
header(sl, "REFERENCES", "参考文献")
refs = [
    '[1] J. Hu et al., "PSELDNets: Pre-trained Neural Networks on Large-scale Synthetic Datasets for Sound Event Localization and Detection," arXiv:2411.06399, 2024.',
    '[2] L. Barbisan et al., "DynamicSound simulator for simulating moving sources and microphone arrays," arXiv:2601.15433, 2026.',
    '[3] E. Fonseca et al., "FSD50K: An Open Dataset of Human-Labeled Sound Events," IEEE/ACM TASLP, 2022.',
    '[4] K. Shimada et al., "Ensemble of ACCDOA- and EINV2-based Systems…," DCASE2021 Challenge, 2021.',
    '[5] Q. Wang et al., "A Four-Stage Data Augmentation Approach to ResNet-Conformer Based Acoustic Modeling…," IEEE/ACM TASLP, 2023.',
    "[6] DCASE2024 Challenge Task 3: Sound Event Localization and Detection with Source Distance Estimation.",
]
tb = box(sl, 54, 116, W - 108, 360)
first = True
for t in refs:
    para(tb.text_frame, t, size=12.5, color=INK2, after=10, line=1.3, first=first)
    first = False
footer(sl, -1)

# ============ 16 付録A ============
sl = new_slide()
header(sl, "APPENDIX A", "付録 — 評価指標")
rows = [
    ("指標", "意味"),
    ("ER ↓", "検出のミス率：正解と比べて何割ミスがあるか"),
    ("F値 ↑", "検出の正確さ：誤検出の少なさと見逃しの少なさのバランス"),
    ("LE_CD ↓", "定位の誤差：正解と比べて何度ズレているか"),
    ("LR_CD ↑", "定位の検出率：方向まで含めて正確に検出できた割合"),
    ("SELDスコア ↓", "上記4つの統合：{ ER＋(1−F)＋LE/180＋(1−LR) } / 4"),
]
gt = sl.shapes.add_table(6, 2, Pt(54), Pt(116), Pt(W - 108), Pt(280)).table
gt.columns[0].width = Pt(180)
gt.columns[1].width = Pt(W - 108 - 180)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gt.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER2 if ri == 0 else WHITE
        tf = cell.text_frame
        tf.word_wrap = True
        para(tf, val, size=13, bold=(ri == 0 or ci == 0), color=INK2, after=0,
             first=True)
tb = box(sl, 54, 410, W - 108, 40)
para(tb.text_frame,
     "距離・通知は別指標：至近≤5mの中央絶対誤差／危険度3段階の到達・抑制・誤警告率（車＝トラック単位）。",
     size=12, color=MUTED, after=0, first=True)
footer(sl, -1)

# ============ 17 付録B ============
sl = new_slide()
header(sl, "APPENDIX B", "付録 — 確定評価の仕組み")
bullets(sl, 54, 120, W - 108, [
    {"runs": [("確定評価セット", {"bold": True, "color": INK}),
              ("：検証セットと同一設計・新乱数・未使用フォールドの1,800本。", {}),
              ("生成前に計画をコミット", {"bold": True, "color": INK}),
              ("（事前登録）し、推論・採点は1回だけ", {})]},
    {"runs": [("開発値と確定値の乖離は全指標±2pt以内 ＝ 同一生成方式の未知シードへ汎化している直接証拠（実環境への汎化は9月に検証）", {})]},
    {"runs": [("採点系の検証", {"bold": True, "color": INK}),
              ("：採点スクリプトは単体テストと独立実装の突き合わせで検証済み。通知ルールの同一物体連続性（±60°連結）もこの過程で追加された", {})]},
    {"runs": [("失敗ケースは", {}), ("個票監査", {"bold": True, "color": INK}),
              ("で全数分類（至近警告に届かなかった重大車153台の内訳まで記録）", {})]},
], size=14.5, gap=16)
footer(sl, -1)

prs.save(OUT)
print(f"saved: {OUT} ({OUT.stat().st_size/1024:.0f} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
